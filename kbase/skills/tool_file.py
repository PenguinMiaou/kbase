"""File tools — PPTX/DOCX read/write, file info, output management."""
import json
import shutil
import time
from pathlib import Path
from .tools import SkillTool, register_tool

OUTPUTS_DIR = Path.home() / ".kbase" / "default" / "outputs"


class FileInfoTool(SkillTool):
    name = "file_info"
    description = "Get file metadata: path, type, size, location in KBase index."
    is_read_only = True
    input_schema = {
        "type": "object",
        "properties": {
            "file_id": {"type": "string", "description": "KBase file ID"},
            "file_path": {"type": "string", "description": "Or direct file path"},
        },
    }

    def call(self, params: dict) -> str:
        fp = params.get("file_path")
        if not fp and params.get("file_id"):
            from kbase.store import KBaseStore
            store = KBaseStore()
            conn = store.conn
            row = conn.execute("SELECT file_path, file_name, file_type, file_size FROM files WHERE file_id = ?",
                               (params["file_id"],)).fetchone()
            if row:
                fp = row["file_path"]
            else:
                return f"Error: file_id {params['file_id']} not found"

        if not fp or not Path(fp).exists():
            return f"Error: File not found: {fp}"

        p = Path(fp)
        return json.dumps({
            "path": str(p),
            "name": p.name,
            "type": p.suffix,
            "size_mb": round(p.stat().st_size / 1048576, 2),
        }, ensure_ascii=False)


class OutputFileTool(SkillTool):
    name = "output_file"
    description = "Register a generated file as skill output. Returns download URL."
    is_read_only = False
    input_schema = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Path to the generated file"},
            "description": {"type": "string", "description": "What was done to the file"},
        },
        "required": ["file_path"],
    }

    def call(self, params: dict) -> str:
        fp = Path(params["file_path"])
        if not fp.exists():
            return f"Error: File not found: {fp}"

        OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
        # Copy to outputs dir with timestamp
        ts = time.strftime("%Y%m%d_%H%M%S")
        out_name = f"{fp.stem}_{ts}{fp.suffix}"
        out_path = OUTPUTS_DIR / out_name
        shutil.copy2(str(fp), str(out_path))

        return json.dumps({
            "status": "ok",
            "output_path": str(out_path),
            "output_name": out_name,
            "download_url": f"/api/skill/download/{out_name}",
            "description": params.get("description", ""),
        }, ensure_ascii=False)


class PptxReadTool(SkillTool):
    name = "pptx_read"
    description = "Read PPTX structure: slide count, text per slide, layout info."
    is_read_only = True
    max_result_chars = 10000
    input_schema = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "slides": {"type": "string", "description": "Slide range: 'all', '1-5', or '3'"},
        },
        "required": ["file_path"],
    }

    def call(self, params: dict) -> str:
        from pptx import Presentation
        fp = params["file_path"]
        if not Path(fp).exists():
            return f"Error: File not found: {fp}"

        prs = Presentation(fp)
        slides_param = params.get("slides", "all")
        result = {"file": Path(fp).name, "slide_count": len(prs.slides), "slides": []}

        for i, slide in enumerate(prs.slides, 1):
            if slides_param != "all":
                if "-" in slides_param:
                    s, e = map(int, slides_param.split("-"))
                    if i < s or i > e:
                        continue
                elif i != int(slides_param):
                    continue

            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            result["slides"].append({"slide": i, "layout": slide.slide_layout.name, "text": texts})

        return json.dumps(result, ensure_ascii=False, indent=2)


class PptxWriteTool(SkillTool):
    name = "pptx_write"
    description = "Modify PPTX slides: update text in specific shapes, add notes."
    is_read_only = False
    input_schema = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "changes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "slide": {"type": "integer", "description": "Slide number (1-based)"},
                        "shape_index": {"type": "integer", "description": "Shape index in slide"},
                        "text": {"type": "string", "description": "New text content"},
                        "notes": {"type": "string", "description": "Speaker notes to add"},
                    },
                },
            },
        },
        "required": ["file_path", "changes"],
    }

    def call(self, params: dict) -> str:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        fp = params["file_path"]
        prs = Presentation(fp)
        changed = 0

        for change in params.get("changes", []):
            slide_idx = change.get("slide", 1) - 1
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]

            if "notes" in change:
                if not slide.has_notes_slide:
                    slide.notes_slide
                slide.notes_slide.notes_text_frame.text = change["notes"]
                changed += 1

            if "text" in change and "shape_index" in change:
                shapes = list(slide.shapes)
                si = change["shape_index"]
                if si < len(shapes) and shapes[si].has_text_frame:
                    tf = shapes[si].text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = RGBColor(0xFF, 0, 0)
                    tf.paragraphs[0].text = change["text"]
                    if tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0, 0)
                    changed += 1

        prs.save(fp)
        return f"OK: {changed} changes applied to {Path(fp).name}"


class DocxReadTool(SkillTool):
    name = "docx_read"
    description = "Read DOCX structure: paragraphs, headings, tables."
    is_read_only = True
    max_result_chars = 10000
    input_schema = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "section": {"type": "string", "description": "'all', 'headings', 'tables', or paragraph range '1-20'"},
        },
        "required": ["file_path"],
    }

    def call(self, params: dict) -> str:
        from docx import Document
        fp = params["file_path"]
        if not Path(fp).exists():
            return f"Error: File not found: {fp}"

        doc = Document(fp)
        section = params.get("section", "all")
        result = {"file": Path(fp).name, "paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}

        if section == "headings":
            result["headings"] = [
                {"index": i, "level": p.style.name, "text": p.text}
                for i, p in enumerate(doc.paragraphs) if p.style.name.startswith("Heading")
            ]
        elif section == "tables":
            tables = []
            for ti, table in enumerate(doc.tables):
                rows = []
                for row in table.rows[:20]:
                    rows.append([cell.text for cell in row.cells])
                tables.append({"table": ti, "rows": rows})
            result["table_data"] = tables
        else:
            paras = []
            for i, p in enumerate(doc.paragraphs):
                if p.text.strip():
                    paras.append({"index": i, "style": p.style.name, "text": p.text[:200]})
            result["content"] = paras[:100]

        return json.dumps(result, ensure_ascii=False, indent=2)


class DocxWriteTool(SkillTool):
    name = "docx_write"
    description = "Modify DOCX: update paragraphs, add content, modify tables. Changes in red font."
    is_read_only = False
    input_schema = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string"},
            "changes": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {"type": "string", "enum": ["replace_paragraph", "insert_after", "append"]},
                        "paragraph_index": {"type": "integer"},
                        "text": {"type": "string"},
                        "comment": {"type": "string"},
                    },
                    "required": ["type", "text"],
                },
            },
        },
        "required": ["file_path", "changes"],
    }

    def call(self, params: dict) -> str:
        from docx import Document
        from docx.shared import RGBColor

        fp = params["file_path"]
        doc = Document(fp)
        changed = 0

        for change in params.get("changes", []):
            ctype = change["type"]
            text = change["text"]

            if ctype == "replace_paragraph" and "paragraph_index" in change:
                idx = change["paragraph_index"]
                if idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]
                    run = p.add_run(f" {text}")
                    run.font.color.rgb = RGBColor(0xFF, 0, 0)
                    changed += 1

            elif ctype == "insert_after" and "paragraph_index" in change:
                idx = change["paragraph_index"]
                if idx < len(doc.paragraphs):
                    new_p = doc.paragraphs[idx]._element
                    from docx.oxml.ns import qn
                    from lxml import etree
                    new_para = etree.SubElement(new_p.getparent(), qn('w:p'))
                    new_p.addnext(new_para)
                    # Simplified: just add text
                    p = doc.paragraphs[idx + 1] if idx + 1 < len(doc.paragraphs) else doc.add_paragraph()
                    run = p.add_run(text)
                    run.font.color.rgb = RGBColor(0xFF, 0, 0)
                    changed += 1

            elif ctype == "append":
                p = doc.add_paragraph()
                run = p.add_run(text)
                run.font.color.rgb = RGBColor(0, 0, 0xFF)
                changed += 1

        doc.save(fp)
        return f"OK: {changed} changes applied to {Path(fp).name}"


register_tool(FileInfoTool())
register_tool(OutputFileTool())
register_tool(PptxReadTool())
register_tool(PptxWriteTool())
register_tool(DocxReadTool())
register_tool(DocxWriteTool())
