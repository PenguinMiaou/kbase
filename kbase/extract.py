"""File content extractors for all supported formats."""
import csv
import io
import os
import re
from pathlib import Path
from typing import Optional


MAX_FILE_SIZE_MB = 500  # Skip files larger than this to prevent OOM

def extract_file(file_path: str | Path) -> dict:
    """Extract text and tables from a file. Returns {text, tables, metadata}."""
    p = Path(file_path)

    # Guard: skip extremely large files to prevent OOM
    try:
        size_mb = p.stat().st_size / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            return {
                "text": f"[File too large: {size_mb:.0f}MB, max {MAX_FILE_SIZE_MB}MB]",
                "tables": [],
                "metadata": {"type": p.suffix, "error": f"File too large ({size_mb:.0f}MB)", "file_size": p.stat().st_size},
            }
    except OSError:
        pass
    ext = p.suffix.lower()
    extractors = {
        ".md": _extract_markdown,
        ".txt": _extract_text,
        ".pptx": _extract_pptx,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xlsx,
        ".csv": _extract_csv,
        ".pdf": _extract_pdf,
        ".html": _extract_html,
        ".eml": _extract_eml,
        ".msg": _extract_eml,
        ".mbox": _extract_mbox,
        ".ppt": _extract_ppt_legacy,
        ".doc": _extract_doc_legacy,
        # Audio formats (whisper ASR)
        ".mp3": _extract_audio,
        ".m4a": _extract_audio,
        ".wav": _extract_audio,
        ".mp4": _extract_audio,
        ".ogg": _extract_audio,
        ".flac": _extract_audio,
        ".webm": _extract_audio,
        # Archives
        ".zip": _extract_archive,
        ".rar": _extract_archive,
        ".7z": _extract_archive,
        ".tar": _extract_archive,
        ".gz": _extract_archive,
        ".tgz": _extract_archive,
    }
    extractor = extractors.get(ext)
    if not extractor:
        return {"text": "", "tables": [], "metadata": {"type": ext, "error": "unsupported"}}
    try:
        # Check readability first
        if not os.access(str(p), os.R_OK):
            return {
                "text": "", "tables": [],
                "metadata": {"type": ext, "file_name": p.name, "file_path": str(p),
                             "error": "Permission denied — grant Full Disk Access in System Settings"},
            }
        result = extractor(p)
        result.setdefault("tables", [])
        result.setdefault("metadata", {})
        result["metadata"]["type"] = ext
        result["metadata"]["file_name"] = p.name
        result["metadata"]["file_path"] = str(p)
        result["metadata"]["file_size"] = p.stat().st_size
        return result
    except PermissionError:
        return {
            "text": "", "tables": [],
            "metadata": {"type": ext, "file_name": p.name, "file_path": str(p),
                         "error": "Permission denied — grant Full Disk Access in System Settings"},
        }
    except Exception as e:
        return {
            "text": "", "tables": [],
            "metadata": {"type": ext, "file_name": p.name, "file_path": str(p), "error": str(e)},
        }


def _extract_markdown(p: Path) -> dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    # Extract title from first heading
    title_match = re.search(r"^#\s+(.+)", text, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else p.stem
    return {"text": text, "metadata": {"title": title}}


def _extract_text(p: Path) -> dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    return {"text": text, "metadata": {"title": p.stem}}


def _extract_pptx(p: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(p))
    slides_text = []
    tables = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {i}]"]
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_parts.append(f"## {slide.shapes.title.text}")
        except Exception:
            pass

        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    rows = []
                    for row in table.rows[1:]:
                        rows.append([cell.text.strip() for cell in row.cells])
                    tables.append({
                        "source": f"slide_{i}",
                        "headers": headers,
                        "rows": rows,
                    })
                    slide_parts.append(_table_to_markdown(headers, rows))
            except Exception:
                # Skip problematic shapes (grouped shapes, SmartArt, etc.)
                continue

        slides_text.append("\n".join(slide_parts))

    title = ""
    try:
        if prs.slides and prs.slides[0].shapes.title:
            title = prs.slides[0].shapes.title.text or ""
    except Exception:
        pass

    return {
        "text": "\n\n".join(slides_text),
        "tables": tables,
        "metadata": {"title": title or p.stem, "slide_count": len(prs.slides)},
    }


def _extract_docx(p: Path) -> dict:
    from docx import Document

    doc = Document(str(p))
    parts = []
    tables = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]
        if tag == "p":
            # Paragraph
            from docx.oxml.ns import qn
            style_elem = element.find(qn("w:pPr"))
            text = element.text or ""
            # Get full text including runs
            for run in element.iter():
                if run.tag.endswith("}t"):
                    if run.text:
                        text = ""  # reset, will use paragraph approach
                        break
            # Use python-docx paragraph approach
            pass

        elif tag == "tbl":
            pass

    # Simpler approach: iterate paragraphs and tables in order
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style and para.style.name and "Heading" in para.style.name:
                level = para.style.name.replace("Heading", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 2
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)

    for i, table in enumerate(doc.tables):
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        rows = []
        for row in table.rows[1:]:
            rows.append([cell.text.strip() for cell in row.cells])
        tables.append({"source": f"table_{i+1}", "headers": headers, "rows": rows})
        parts.append(_table_to_markdown(headers, rows))

    title = doc.paragraphs[0].text.strip() if doc.paragraphs else p.stem
    return {
        "text": "\n\n".join(parts),
        "tables": tables,
        "metadata": {"title": title},
    }


def _extract_xlsx(p: Path) -> dict:
    # Handle old .xls format with xlrd
    if p.suffix.lower() == '.xls':
        return _extract_xls_legacy(p)

    from openpyxl import load_workbook

    wb = load_workbook(str(p), read_only=True, data_only=True)
    all_text_parts = []
    tables = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            row_vals = [str(v).strip() if v is not None else "" for v in row]
            if any(row_vals):
                rows_data.append(row_vals)

        if not rows_data:
            continue

        headers = rows_data[0]
        data_rows = rows_data[1:]

        tables.append({
            "source": sheet_name,
            "headers": headers,
            "rows": data_rows,
            "file_name": p.stem,
        })

        # Also create text representation
        all_text_parts.append(f"## Sheet: {sheet_name}")
        all_text_parts.append(_table_to_markdown(headers, data_rows[:50]))  # cap for text
        if len(data_rows) > 50:
            all_text_parts.append(f"... ({len(data_rows)} rows total)")

    wb.close()
    return {
        "text": "\n\n".join(all_text_parts),
        "tables": tables,
        "metadata": {"title": p.stem, "sheet_count": len(wb.sheetnames)},
    }


def _extract_csv(p: Path) -> dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return {"text": "", "tables": []}

    headers = rows[0]
    data_rows = rows[1:]
    tables = [{"source": "csv", "headers": headers, "rows": data_rows, "file_name": p.stem}]

    return {
        "text": _table_to_markdown(headers, data_rows),
        "tables": tables,
        "metadata": {"title": p.stem},
    }


def _extract_pdf(p: Path) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(str(p))
    parts = []
    tables = []
    page_count = len(doc)

    for i in range(page_count):
        try:
            page = doc.load_page(i)
            text = page.get_text("text").strip()
            if text:
                parts.append(f"[Page {i+1}]\n{text}")
            page_tables = _extract_pdf_tables(page, i + 1)
            tables.extend(page_tables)
        except Exception:
            continue

    doc.close()

    title = parts[0].split("\n")[1][:100] if parts else p.stem
    return {
        "text": "\n\n".join(parts),
        "tables": tables,
        "metadata": {"title": title, "page_count": page_count},
    }


def _extract_pdf_tables(page, page_num: int) -> list:
    """Best-effort table extraction from PDF page using PyMuPDF."""
    try:
        tabs = page.find_tables()
        results = []
        for i, tab in enumerate(tabs):
            data = tab.extract()
            if data and len(data) > 1:
                headers = [str(c) if c else "" for c in data[0]]
                rows = [[str(c) if c else "" for c in row] for row in data[1:]]
                results.append({
                    "source": f"page_{page_num}_table_{i+1}",
                    "headers": headers,
                    "rows": rows,
                })
        return results
    except Exception:
        return []


def _extract_html(p: Path) -> dict:
    text = p.read_text(encoding="utf-8", errors="replace")
    # Simple HTML tag stripping
    clean = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
    clean = re.sub(r"<style[^>]*>.*?</style>", "", clean, flags=re.DOTALL)
    clean = re.sub(r"<[^>]+>", " ", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    return {"text": clean, "metadata": {"title": p.stem}}


def _extract_audio(p: Path, whisper_model: str = None) -> dict:
    """Extract text from audio using Whisper (configurable model)."""
    from kbase.config import load_settings, WHISPER_MODELS, DEFAULT_WHISPER_MODEL

    settings = load_settings()
    model_key = whisper_model or settings.get("whisper_model", DEFAULT_WHISPER_MODEL)
    model_info = WHISPER_MODELS.get(model_key, {"name": "base", "type": "local"})

    if model_info["type"] == "openai-api":
        return _extract_audio_openai_api(p)
    elif model_info["type"] == "faster-whisper":
        return _extract_audio_faster_whisper(p, model_info["name"])
    else:
        return _extract_audio_whisper(p, model_info["name"])


def _extract_audio_whisper(p: Path, model_size: str = "base") -> dict:
    """Local OpenAI Whisper."""
    try:
        import whisper
    except ImportError:
        return {
            "text": f"[Audio: {p.name}]",
            "metadata": {"title": p.stem, "error": "Install: pip install openai-whisper"},
        }

    model = whisper.load_model(model_size)
    result = model.transcribe(str(p), language=None)

    segments = result.get("segments", [])
    parts = []
    for seg in segments:
        mins, secs = divmod(int(seg["start"]), 60)
        parts.append(f"[{mins:02d}:{secs:02d}] {seg['text'].strip()}")

    return {
        "text": "\n".join(parts),
        "metadata": {
            "title": p.stem,
            "duration_seconds": segments[-1]["end"] if segments else 0,
            "language": result.get("language", "unknown"),
            "segment_count": len(segments),
            "whisper_model": model_size,
        },
    }


def _extract_audio_faster_whisper(p: Path, model_size: str = "large-v3") -> dict:
    """Faster Whisper (CTranslate2 accelerated)."""
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return {
            "text": f"[Audio: {p.name}]",
            "metadata": {"title": p.stem, "error": "Install: pip install faster-whisper"},
        }

    model = WhisperModel(model_size, device="auto", compute_type="auto")
    segments_gen, info = model.transcribe(str(p))

    parts = []
    last_end = 0
    for seg in segments_gen:
        mins, secs = divmod(int(seg.start), 60)
        parts.append(f"[{mins:02d}:{secs:02d}] {seg.text.strip()}")
        last_end = seg.end

    return {
        "text": "\n".join(parts),
        "metadata": {
            "title": p.stem,
            "duration_seconds": last_end,
            "language": info.language,
            "segment_count": len(parts),
            "whisper_model": f"faster-{model_size}",
        },
    }


def _extract_audio_openai_api(p: Path) -> dict:
    """OpenAI Whisper API (cloud)."""
    import os
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return {
            "text": f"[Audio: {p.name}]",
            "metadata": {"title": p.stem, "error": "OPENAI_API_KEY not set"},
        }

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        with open(str(p), "rb") as f:
            result = client.audio.transcriptions.create(
                model="whisper-1",
                file=f,
                response_format="verbose_json",
            )

        segments = result.segments or []
        parts = []
        for seg in segments:
            mins, secs = divmod(int(seg["start"]), 60)
            parts.append(f"[{mins:02d}:{secs:02d}] {seg['text'].strip()}")

        return {
            "text": "\n".join(parts) or result.text,
            "metadata": {
                "title": p.stem,
                "duration_seconds": result.duration,
                "language": result.language,
                "segment_count": len(segments),
                "whisper_model": "openai-api",
            },
        }
    except Exception as e:
        return {"text": "", "metadata": {"title": p.stem, "error": str(e)}}


def _extract_eml(p: Path) -> dict:
    """Extract email content from .eml files."""
    import email
    from email import policy

    with open(str(p), "rb") as f:
        msg = email.message_from_binary_file(f, policy=policy.default)

    subject = msg.get("subject", "") or p.stem
    from_addr = msg.get("from", "")
    to_addr = msg.get("to", "")
    date = msg.get("date", "")

    # Extract body text
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                body += part.get_content() or ""
            elif ct == "text/html" and not body:
                html = part.get_content() or ""
                body = re.sub(r"<[^>]+>", " ", html)
                body = re.sub(r"\s+", " ", body).strip()
    else:
        ct = msg.get_content_type()
        content = msg.get_content() or ""
        if ct == "text/html":
            body = re.sub(r"<[^>]+>", " ", content)
            body = re.sub(r"\s+", " ", body).strip()
        else:
            body = content

    text = f"# {subject}\n\nFrom: {from_addr}\nTo: {to_addr}\nDate: {date}\n\n{body}"

    return {
        "text": text,
        "metadata": {
            "title": subject,
            "from": from_addr,
            "date": date,
        },
    }


def _extract_mbox(p: Path) -> dict:
    """Extract all emails from an mbox file (Feishu export format)."""
    import mailbox
    import email
    from email import policy
    from email.header import decode_header

    def decode_mime(s):
        if not s:
            return ""
        parts = decode_header(s)
        decoded = []
        for part, charset in parts:
            if isinstance(part, bytes):
                decoded.append(part.decode(charset or "utf-8", errors="replace"))
            else:
                decoded.append(part)
        return " ".join(decoded)

    mbox = mailbox.mbox(str(p))
    all_parts = []
    count = 0

    for msg in mbox:
        try:
            subject = decode_mime(msg.get("subject", "")) or "No Subject"
            from_addr = decode_mime(msg.get("from", ""))
            to_addr = decode_mime(msg.get("to", ""))
            date = msg.get("date", "")

            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or "utf-8"
                            body += payload.decode(charset, errors="replace")
                    elif ct == "text/html" and not body:
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or "utf-8"
                            html = payload.decode(charset, errors="replace")
                            body = re.sub(r"<[^>]+>", " ", html)
                            body = re.sub(r"\s+", " ", body).strip()
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    charset = msg.get_content_charset() or "utf-8"
                    body = payload.decode(charset, errors="replace")
                    if msg.get_content_type() == "text/html":
                        body = re.sub(r"<[^>]+>", " ", body)
                        body = re.sub(r"\s+", " ", body).strip()

            all_parts.append(f"## {subject}\n\nFrom: {from_addr}\nTo: {to_addr}\nDate: {date}\n\n{body[:2000]}")
            count += 1
        except Exception:
            continue

    mbox.close()

    return {
        "text": "\n\n---\n\n".join(all_parts),
        "metadata": {
            "title": f"{p.stem} ({count} emails)",
            "email_count": count,
        },
    }


def _extract_archive(p: Path) -> dict:
    """Extract contents from archive files (zip/tar/gz).

    Extracts to temp dir, then recursively extracts each supported file inside.
    """
    import tempfile
    import shutil
    import zipfile
    import tarfile

    ext = p.suffix.lower()
    all_text = []
    all_tables = []
    file_count = 0
    temp_dir = None

    try:
        temp_dir = tempfile.mkdtemp(prefix="kbase_archive_")

        if ext == ".zip":
            with zipfile.ZipFile(str(p), "r") as zf:
                zf.extractall(temp_dir)
        elif ext in (".tar", ".gz", ".tgz"):
            with tarfile.open(str(p), "r:*") as tf:
                tf.extractall(temp_dir)
        else:
            # .rar, .7z
            import subprocess
            if ext == ".rar":
                try:
                    import rarfile
                    with rarfile.RarFile(str(p)) as rf:
                        rf.extractall(temp_dir)
                except ImportError:
                    result = subprocess.run(["unrar", "x", "-y", str(p), temp_dir],
                                             capture_output=True, timeout=120)
            elif ext == ".7z":
                result = subprocess.run(["7z", "x", f"-o{temp_dir}", "-y", str(p)],
                                         capture_output=True, timeout=120)
            else:
                return {"text": "", "metadata": {"title": p.stem, "error": f"Unsupported archive: {ext}"}}

        # Now recursively extract all supported files inside
        from kbase.config import SUPPORTED_EXTENSIONS
        archive_exts = {".zip", ".rar", ".7z", ".tar", ".gz", ".tgz"}
        for root, dirs, files in os.walk(temp_dir):
            # Skip hidden dirs
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for fname in files:
                fpath = Path(root) / fname
                if fpath.suffix.lower() in SUPPORTED_EXTENSIONS and fpath.suffix.lower() not in archive_exts:
                    try:
                        result = extract_file(str(fpath))
                        if result.get("text"):
                            all_text.append(f"[Archive: {p.name} / {fname}]\n{result['text']}")
                            file_count += 1
                        all_tables.extend(result.get("tables", []))
                    except Exception:
                        continue

    except zipfile.BadZipFile:
        return {"text": "", "metadata": {"title": p.stem, "error": "Bad zip file"}}
    except Exception as e:
        return {"text": "", "metadata": {"title": p.stem, "error": str(e)}}
    finally:
        if temp_dir:
            shutil.rmtree(temp_dir, ignore_errors=True)

    return {
        "text": "\n\n---\n\n".join(all_text),
        "tables": all_tables,
        "metadata": {
            "title": p.stem,
            "archive_files": file_count,
        },
    }


def _extract_xls_legacy(p: Path) -> dict:
    """Extract old .xls format using xlrd."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(p))
        all_text = []
        tables = []
        for sheet in wb.sheets():
            rows_data = []
            for row_idx in range(sheet.nrows):
                row_vals = [str(sheet.cell_value(row_idx, col)).strip() for col in range(sheet.ncols)]
                if any(row_vals):
                    rows_data.append(row_vals)
            if not rows_data:
                continue
            headers = rows_data[0]
            data_rows = rows_data[1:]
            tables.append({"source": sheet.name, "headers": headers, "rows": data_rows, "file_name": p.stem})
            all_text.append(f"## Sheet: {sheet.name}")
            all_text.append(_table_to_markdown(headers, data_rows[:50]))
        return {"text": "\n\n".join(all_text), "tables": tables, "metadata": {"title": p.stem}}
    except ImportError:
        return {"text": f"[.xls file: {p.name} — install xlrd: pip install xlrd]", "metadata": {"title": p.stem, "error": "xlrd not installed"}}


def _extract_ppt_legacy(p: Path) -> dict:
    """Legacy .ppt files - try to read, fall back to filename only."""
    return {"text": f"[Legacy .ppt file: {p.name}]", "metadata": {"title": p.stem, "legacy": True}}


def _extract_doc_legacy(p: Path) -> dict:
    """Legacy .doc files - try to read, fall back to filename only."""
    return {"text": f"[Legacy .doc file: {p.name}]", "metadata": {"title": p.stem, "legacy": True}}


def _table_to_markdown(headers: list, rows: list) -> str:
    """Convert table data to markdown format."""
    if not headers:
        return ""
    lines = ["| " + " | ".join(headers) + " |"]
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
    for row in rows:
        # Pad row to match headers length
        padded = row + [""] * (len(headers) - len(row))
        lines.append("| " + " | ".join(padded[:len(headers)]) + " |")
    return "\n".join(lines)
