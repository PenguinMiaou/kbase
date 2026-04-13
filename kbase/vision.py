"""Vision module: extract and describe images from documents using Vision LLMs."""
import base64
import io
import json
import os
from pathlib import Path
from typing import Optional

from kbase.config import VISION_MODELS, DEFAULT_VISION_MODEL, load_settings


def describe_image(image_bytes: bytes, context: str = "", settings: dict = None) -> str:
    """Describe an image using the configured Vision LLM.

    Args:
        image_bytes: Raw image bytes (PNG/JPEG)
        context: Optional context (e.g., "Slide 3 of quarterly report.pptx")
        settings: Settings dict with API keys and vision_model selection

    Returns:
        Text description of the image content
    """
    settings = settings or load_settings()
    model_key = settings.get("vision_model", DEFAULT_VISION_MODEL)

    if model_key == "none" or not model_key:
        return ""

    model_info = VISION_MODELS.get(model_key)
    if not model_info:
        return ""

    model_type = model_info["type"]
    b64 = base64.b64encode(image_bytes).decode("utf-8")

    prompt = "Describe this image in detail. Include all text, numbers, labels, chart data, diagram structure, and relationships shown."
    if context:
        prompt += f"\nContext: {context}"
    prompt += "\nOutput in the same language as any text visible in the image. If mixed languages, use both."

    try:
        if model_type == "openai":
            return _call_openai_vision(model_info, b64, prompt, settings)
        elif model_type == "anthropic":
            return _call_anthropic_vision(model_info, b64, prompt, settings)
        elif model_type == "gemini":
            return _call_gemini_vision(model_info, b64, prompt, settings)
        elif model_type in ("dashscope", "openai-compatible"):
            return _call_openai_compatible_vision(model_info, b64, prompt, settings)
        elif model_type == "ollama":
            return _call_ollama_vision(model_info, b64, prompt, settings)
    except Exception as e:
        print(f"[KBase Vision] Error: {e}")
        return ""

    return ""


def extract_images_from_pptx(file_path: str) -> list[dict]:
    """Extract images from PPTX file with slide context."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            # Get slide title for context
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text
            except Exception:
                pass

            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img = shape.image
                        images.append({
                            "bytes": img.blob,
                            "content_type": img.content_type,
                            "context": f"Slide {i}: {title}" if title else f"Slide {i}",
                            "slide": i,
                        })
                except Exception:
                    continue
    except Exception as e:
        print(f"[KBase Vision] PPTX image extraction error: {e}")

    return images


def extract_images_from_pdf(file_path: str) -> list[dict]:
    """Extract images from PDF file with page context."""
    images = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images()
            for img_idx, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n > 4:  # CMYK
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    # Skip tiny images (icons, bullets)
                    if pix.width < 100 or pix.height < 100:
                        continue
                    images.append({
                        "bytes": pix.tobytes("png"),
                        "content_type": "image/png",
                        "context": f"Page {page_num + 1}",
                        "page": page_num + 1,
                    })
                except Exception:
                    continue
        doc.close()
    except Exception as e:
        print(f"[KBase Vision] PDF image extraction error: {e}")

    return images


def describe_document_images(file_path: str, settings: dict = None, max_images: int = 20) -> list[dict]:
    """Extract and describe all images in a document.

    Returns list of {text, context, slide/page} for each described image.
    """
    settings = settings or load_settings()
    model_key = settings.get("vision_model", DEFAULT_VISION_MODEL)
    if model_key == "none" or not model_key:
        return []

    p = Path(file_path)
    ext = p.suffix.lower()

    if ext == ".pptx":
        images = extract_images_from_pptx(file_path)
    elif ext == ".pdf":
        images = extract_images_from_pdf(file_path)
    else:
        return []

    # Limit images to process
    images = images[:max_images]

    descriptions = []
    for img in images:
        context = f"{p.name} — {img.get('context', '')}"
        desc = describe_image(img["bytes"], context=context, settings=settings)
        if desc:
            descriptions.append({
                "text": f"[Image: {img.get('context', '')}]\n{desc}",
                "context": img.get("context", ""),
                "slide": img.get("slide"),
                "page": img.get("page"),
            })

    return descriptions


# ── Vision LLM Backends ──

def _call_openai_vision(model_info, b64, prompt, settings):
    import urllib.request
    api_key = settings.get("openai_api_key", "") or os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return ""
    data = json.dumps({
        "model": model_info["model"],
        "messages": [{"role": "user", "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"}},
        ]}],
        "max_tokens": 500,
    }).encode()
    req = urllib.request.Request(
        "https://api.openai.com/v1/chat/completions",
        data=data,
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        result = json.loads(resp.read())
    return result["choices"][0]["message"]["content"]


def _call_anthropic_vision(model_info, b64, prompt, settings):
    import urllib.request
    api_key = settings.get("anthropic_api_key", "") or os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        return ""
    data = json.dumps({
        "model": model_info["model"],
        "max_tokens": 500,
        "messages": [{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}},
            {"type": "text", "text": prompt},
        ]}],
    }).encode()
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=data,
        headers={
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
            "Content-Type": "application/json",
        },
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        result = json.loads(resp.read())
    return result["content"][0]["text"]


def _call_gemini_vision(model_info, b64, prompt, settings):
    import urllib.request
    api_key = settings.get("gemini_api_key", "") or os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        return ""
    data = json.dumps({
        "contents": [{"parts": [
            {"text": prompt},
            {"inline_data": {"mime_type": "image/png", "data": b64}},
        ]}],
    }).encode()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_info['model']}:generateContent?key={api_key}"
    req = urllib.request.Request(url, data=data, headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=30) as resp:
        result = json.loads(resp.read())
    return result["candidates"][0]["content"]["parts"][0]["text"]


def _call_openai_compatible_vision(model_info, b64, prompt, settings):
    import urllib.request
    key_env = model_info.get("key_env", "").lower()
    api_key = settings.get(key_env, "") or os.environ.get(model_info.get("key_env", ""), "")
    base_url = model_info.get("base_url", "https://api.openai.com/v1")
    if not api_key:
        return ""
    data = json.dumps({
        "model": model_info["model"],
        "messages": [{"role": "user", "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
        ]}],
        "max_tokens": 500,
    }).encode()
    req = urllib.request.Request(
        f"{base_url}/chat/completions",
        data=data,
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        result = json.loads(resp.read())
    return result["choices"][0]["message"]["content"]


def _call_ollama_vision(model_info, b64, prompt, settings):
    import urllib.request
    model = settings.get("ollama_vision_model") or model_info.get("model", "minicpm-v")
    data = json.dumps({
        "model": model,
        "prompt": prompt,
        "images": [b64],
        "stream": False,
    }).encode()
    req = urllib.request.Request(
        "http://localhost:11434/api/generate",
        data=data,
        headers={"Content-Type": "application/json"},
    )
    with urllib.request.urlopen(req, timeout=60) as resp:
        result = json.loads(resp.read())
    return result.get("response", "")
