"""Web API and frontend for kbase. Copyright@PenguinMiaou"""
import json
import os
import signal
import subprocess
import sys
import time
import platform
from pathlib import Path

from fastapi import FastAPI, Query, UploadFile, File, Form, HTTPException, Request
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware

from kbase.store import KBaseStore
from kbase.ingest import ingest_directory, ingest_file
from kbase.search import hybrid_search, semantic_only, keyword_only, sql_search, get_table_context
from kbase.chat import (
    chat, clear_conversation, LLM_PROVIDERS, BUDDY_PRESETS,
    _conversations, _conv_titles, _save_conversations, _load_conversations,
    generate_title, get_conv_title, set_conv_title,
    get_memories, add_memory, delete_memory, extract_memories_from_conversation,
)
from kbase.config import (
    EMBEDDING_MODELS, WHISPER_MODELS, LANGUAGE_PROFILES,
    load_settings, save_settings, CHUNK_MAX_CHARS, CHUNK_OVERLAP_CHARS,
)


def _auto_save_research(report: str, question: str, workspace: str, settings: dict):
    """Auto-save research report to file and index it into the knowledge base.

    Creates a Markdown file in ~/.kbase/{workspace}/research_reports/
    and ingests it so future searches can find previous research outputs.
    """
    from kbase.config import get_workspace_dir

    ws_dir = get_workspace_dir(workspace)
    reports_dir = ws_dir / "research_reports"
    reports_dir.mkdir(parents=True, exist_ok=True)

    timestamp = time.strftime("%Y%m%d_%H%M%S")
    # Sanitize question for filename (take first 40 chars, remove special chars)
    safe_q = "".join(c for c in question[:40] if c.isalnum() or c in " _-").strip()
    safe_q = safe_q.replace(" ", "_") or "research"
    report_file = reports_dir / f"{timestamp}_{safe_q}.md"

    # Write report with metadata header
    content = (
        f"---\n"
        f"question: {question}\n"
        f"generated: {time.strftime('%Y-%m-%d %H:%M:%S')}\n"
        f"type: research_report\n"
        f"---\n\n"
        f"# {question}\n\n"
        f"{report}"
    )
    report_file.write_text(content, encoding="utf-8")

    # Auto-ingest the report into the knowledge base
    try:
        from kbase.store import KBaseStore
        from kbase.config import get_db_path, get_chroma_path
        store = KBaseStore(str(get_db_path(workspace)), str(get_chroma_path(workspace)))
        try:
            ingest_file(store, str(report_file))
        finally:
            store.close()
        print(f"[KBase] Research report saved and indexed: {report_file.name}")
    except Exception as e:
        print(f"[KBase] Research report saved but indexing failed: {e}")

    # Add reports_dir to ingest_dirs if not already there
    ingest_dirs = settings.get("ingest_dirs", [])
    reports_dir_str = str(reports_dir)
    if reports_dir_str not in ingest_dirs:
        ingest_dirs.append(reports_dir_str)
        settings["ingest_dirs"] = ingest_dirs
        save_settings(workspace, settings)


def _validate_file_path(file_path: str, store=None) -> str:
    """Security: validate file path is within allowed directories.

    Prevents path traversal attacks. Only allows:
    - Files registered in the DB (indexed files)
    - Files under user's home directory
    - Temp files (for slides/conversion cache)
    """
    import tempfile
    resolved = str(Path(file_path).resolve())
    allowed_roots = [
        str(Path.home()),
        tempfile.gettempdir(),
        str(Path.home() / ".kbase"),
    ]
    if not any(resolved.startswith(root) for root in allowed_roots):
        raise HTTPException(403, "Access denied: path outside allowed directories")
    # Block access to sensitive files
    sensitive = [".ssh", ".gnupg", ".aws", "credentials", ".env", "id_rsa", "shadow"]
    path_lower = resolved.lower()
    if any(s in path_lower for s in sensitive):
        raise HTTPException(403, "Access denied: sensitive file")
    return resolved


def _sanitize_filename(filename: str) -> str:
    """Security: sanitize uploaded filename to prevent path injection."""
    # Remove path separators and null bytes
    name = filename.replace("/", "_").replace("\\", "_").replace("\x00", "")
    # Keep only safe characters
    name = "".join(c for c in name if c.isalnum() or c in ".-_ ()")
    return name[:200] or "upload"


def create_app(workspace: str = "default") -> FastAPI:
    app = FastAPI(title="KBase", description="Local Knowledge Base API")
    # Security: restrict CORS to localhost only (not open to all origins)
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["http://localhost:8765", "http://127.0.0.1:8765", "http://localhost:*", "http://127.0.0.1:*"],
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # Static files for logos
    static_dir = Path(__file__).parent / "static"
    if static_dir.exists():
        app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

    def get_store():
        return KBaseStore(workspace)

    _stats_cache = [{"file_count": 0, "chunk_count": 0, "table_count": 0, "error_count": 0, "type_counts": {}}]

    # ---- Search API ----

    @app.get("/api/search")
    def api_search(
        q: str = Query(...), type: str = Query("auto"),
        top_k: int = Query(10), file_type: str = Query(None),
    ):
        store = get_store()
        try:
            if type == "semantic":
                return semantic_only(store, q, top_k=top_k, file_type=file_type)
            elif type == "keyword":
                return keyword_only(store, q, top_k=top_k)
            return hybrid_search(store, q, top_k=top_k)
        finally:
            store.close()

    @app.get("/api/sql")
    def api_sql(q: str = Query(...)):
        store = get_store()
        try:
            return sql_search(store, q)
        finally:
            store.close()

    @app.get("/api/tables")
    def api_tables():
        store = get_store()
        try:
            return get_table_context(store)
        finally:
            store.close()

    @app.get("/api/status")
    def api_status():
        try:
            store = get_store()
            try:
                stats = store.get_stats()
                _stats_cache[0] = stats
                return stats
            finally:
                store.close()
        except Exception:
            return {**_stats_cache[0], "workspace": workspace, "db_locked": True}

    @app.get("/api/files")
    def api_files(source_dir: str = Query(None)):
        store = get_store()
        try:
            return {"files": store.list_files(source_dir), "count": len(store.list_files(source_dir))}
        finally:
            store.close()

    @app.get("/api/browse-dir")
    def api_browse_dir():
        """Open native directory picker dialog using the best available method."""
        import subprocess, platform

        # Method 1: macOS — osascript (works in DMG, no dependencies)
        if platform.system() == "Darwin":
            try:
                result = subprocess.run(
                    ["osascript", "-e",
                     'tell application "System Events" to activate\n'
                     'set chosenFolder to choose folder with prompt "Select folder to index"\n'
                     'return POSIX path of chosenFolder'],
                    capture_output=True, text=True, timeout=120,
                )
                path = result.stdout.strip()
                if path:
                    # Test if we can actually read the directory
                    try:
                        os.listdir(path)
                    except PermissionError:
                        # Can't access — prompt user to grant permission
                        subprocess.run(["osascript", "-e",
                            'display dialog "KBase needs permission to access this folder.\\n\\n'
                            'Click OK to open System Settings, then add KBase to Full Disk Access." '
                            'buttons {"Cancel", "Open Settings"} default button "Open Settings"'],
                            capture_output=True, text=True, timeout=30,
                        )
                        subprocess.run(["open", "x-apple.systempreferences:com.apple.preference.security?Privacy_AllFiles"])
                        return {"path": "", "error": "permission_needed"}
                    return {"path": path}
            except Exception:
                pass

        # Method 2: Linux — zenity
        if platform.system() == "Linux":
            try:
                import shutil
                if shutil.which("zenity"):
                    result = subprocess.run(
                        ["zenity", "--file-selection", "--directory", "--title=Select folder to index"],
                        capture_output=True, text=True, timeout=120,
                    )
                    if result.returncode == 0:
                        return {"path": result.stdout.strip()}
            except Exception:
                pass

        # Method 3: Tkinter fallback
        try:
            import threading, queue as _queue
            q = _queue.Queue()
            def pick():
                try:
                    import tkinter as tk
                    from tkinter import filedialog
                    root = tk.Tk()
                    root.withdraw()
                    root.attributes('-topmost', True)
                    path = filedialog.askdirectory(title="Select folder to index")
                    root.destroy()
                    q.put(path or "")
                except Exception:
                    q.put("")
            t = threading.Thread(target=pick)
            t.start()
            t.join(timeout=120)
            return {"path": q.get(timeout=1) if not q.empty() else ""}
        except Exception:
            return {"path": ""}

    @app.post("/api/ingest")
    def api_ingest(directory: str = Form(...), force: bool = Form(False)):
        store = get_store()
        try:
            result = ingest_directory(store, directory, force=force)
            # Track directory
            settings_data = load_settings(workspace)
            dirs = settings_data.setdefault("ingest_dirs", {})
            dirs[directory] = {
                "enabled": True,
                "last_sync": time.time(),
                "file_count": result.get("total", 0),
                "status": "ok" if result.get("failed", 0) == 0 else "partial",
            }
            save_settings(workspace, settings_data)
            return result
        finally:
            store.close()

    @app.post("/api/ingest/pause")
    def api_ingest_pause():
        from kbase.ingest import pause_ingest
        pause_ingest()
        return {"status": "toggled"}

    @app.post("/api/ingest/stop")
    def api_ingest_stop():
        from kbase.ingest import stop_ingest
        stop_ingest()
        return {"status": "stopped"}

    @app.post("/api/ingest/resume")
    def api_ingest_resume():
        from kbase.ingest import resume_ingest
        resume_ingest()
        return {"status": "resumed"}

    @app.get("/api/ingest/status")
    def api_ingest_status():
        """Check if an ingest is currently running (for page restore after refresh)."""
        from kbase.ingest import _ingest_active, _ingest_progress
        return {"active": _ingest_active, "progress": _ingest_progress}

    @app.get("/api/ingest-stream")
    def api_ingest_stream(directory: str = Query(...), force: bool = Query(False)):
        """SSE endpoint for real-time ingest progress."""
        # Security: validate ingest directory is under user's home
        resolved = str(Path(directory).resolve())
        home = str(Path.home())
        if not resolved.startswith(home):
            raise HTTPException(403, "Can only ingest directories under your home folder")

        import queue, threading
        q = queue.Queue()

        def do_ingest():
            store = get_store()
            try:
                # Track directory in settings IMMEDIATELY (not after completion)
                # so it persists even if ingest is paused/interrupted
                settings_data = load_settings(workspace)
                dirs = settings_data.setdefault("ingest_dirs", {})
                if isinstance(dirs, list):
                    dirs = {d: {"enabled": True} for d in dirs if isinstance(d, str)}
                    settings_data["ingest_dirs"] = dirs
                if directory not in dirs:
                    dirs[directory] = {"enabled": True, "last_sync": 0, "file_count": 0, "status": "ingesting"}
                save_settings(workspace, settings_data)

                def cb(current, total, name, status):
                    q.put(json.dumps({"current": current, "total": total, "name": name, "status": status}))
                stats = ingest_directory(store, directory, force=force, progress_callback=cb)
                # Update directory status on completion
                settings_data = load_settings(workspace)
                dirs = settings_data.setdefault("ingest_dirs", {})
                dirs[directory] = {
                    "enabled": True,
                    "last_sync": time.time(),
                    "file_count": stats.get("total", 0),
                    "status": "ok" if stats.get("failed", 0) == 0 else "partial",
                }
                save_settings(workspace, settings_data)
                q.put(json.dumps({"done": True, **stats}))
            except Exception as e:
                q.put(json.dumps({"done": True, "error": str(e)}))
            finally:
                store.close()
                q.put(None)

        threading.Thread(target=do_ingest, daemon=True).start()

        def event_stream():
            while True:
                msg = q.get()
                if msg is None:
                    break
                yield f"data: {msg}\n\n"

        return StreamingResponse(event_stream(), media_type="text/event-stream")

    @app.post("/api/add")
    async def api_add_file(file: UploadFile = File(...)):
        from kbase.config import SUPPORTED_EXTENSIONS
        # Security: validate file extension
        ext = Path(file.filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            raise HTTPException(400, f"File type '{ext}' not supported. Allowed: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        # Security: enforce size limit (500MB)
        MAX_UPLOAD_SIZE = 500 * 1024 * 1024
        content = await file.read()
        if len(content) > MAX_UPLOAD_SIZE:
            raise HTTPException(413, f"File too large ({len(content) // 1024 // 1024}MB). Max: 500MB")
        # Security: sanitize filename
        safe_name = _sanitize_filename(file.filename)
        if not Path(safe_name).suffix:
            safe_name += ext  # Ensure extension preserved
        upload_dir = Path.home() / ".kbase" / workspace / "uploads"
        upload_dir.mkdir(parents=True, exist_ok=True)
        # Security: atomic write via temp file (M10)
        import tempfile
        tmp = tempfile.NamedTemporaryFile(dir=str(upload_dir), suffix=ext, delete=False)
        try:
            tmp.write(content)
            tmp.close()
            dest = upload_dir / safe_name
            os.rename(tmp.name, str(dest))
        except Exception:
            os.unlink(tmp.name)
            raise
        store = get_store()
        try:
            return ingest_file(store, str(dest), force=True)
        finally:
            store.close()

    @app.delete("/api/files/{file_id}")
    def api_remove_file(file_id: str):
        store = get_store()
        try:
            files = store.list_files()
            target = next((f for f in files if f["file_id"] == file_id), None)
            if not target:
                raise HTTPException(404, "File not found")
            store.remove_file(target["file_path"])
            return {"status": "removed", "file": target["file_path"]}
        finally:
            store.close()

    @app.post("/api/files/remove")
    async def api_remove_file_by_path(request: Request):
        body = await request.json()
        path = body.get("path", "")
        if not path:
            raise HTTPException(400, "No path")
        store = get_store()
        try:
            store.remove_file(path)
            # Also clean graph edges
            fid = store.file_id(path)
            c = store.conn.cursor()
            c.execute("DELETE FROM document_edges WHERE source_file_id=? OR target_file_id=?", (fid, fid))
            c.execute("DELETE FROM graph_node_positions WHERE file_id=?", (fid,))
            store.conn.commit()
            return {"status": "removed"}
        except Exception as e:
            raise HTTPException(500, str(e))

    # ---- Ingested Directories API ----

    @app.get("/api/ingest-dirs")
    def api_ingest_dirs():
        """List tracked ingested directories with status (file count from DB, not settings)."""
        settings_data = load_settings(workspace)
        dirs = settings_data.get("ingest_dirs", {})
        import time as _time
        now = _time.time()
        result = []
        if isinstance(dirs, list):
            dirs = {d: {"enabled": True} for d in dirs if isinstance(d, str)}
        if not isinstance(dirs, dict):
            dirs = {}

        # Get real file counts from DB (not stale settings value)
        dir_counts = {}
        try:
            store = get_store()
            c = store.conn.cursor()
            c.execute("SELECT source_dir, COUNT(*) as cnt FROM files WHERE error IS NULL OR error = '' GROUP BY source_dir")
            for row in c.fetchall():
                sd = row["source_dir"] or ""
                for p in dirs:
                    # Normalize: strip trailing slashes for comparison
                    pn = p.rstrip("/").rstrip(os.sep)
                    if sd == pn or sd.startswith(pn + "/") or sd.startswith(pn + os.sep):
                        dir_counts[p] = dir_counts.get(p, 0) + row["cnt"]
            store.close()
        except Exception:
            pass

        for path, info in dirs.items():
            if not isinstance(info, dict):
                info = {"enabled": True}
            last_sync = info.get("last_sync", 0)
            if last_sync == 0:
                ago = "Not synced yet"
            elif (elapsed := now - last_sync) < 60:
                ago = f"{int(elapsed)}s ago"
            elif elapsed < 3600:
                ago = f"{int(elapsed/60)}m ago"
            elif elapsed < 86400:
                ago = f"{int(elapsed/3600)}h ago"
            else:
                ago = f"{int(elapsed/86400)}d ago"
            result.append({
                "path": path,
                "enabled": info.get("enabled", True),
                "last_sync": info.get("last_sync", 0),
                "ago": ago,
                "file_count": dir_counts.get(path, 0),
                "status": info.get("status", "unknown"),
            })
        return {"dirs": result}

    @app.post("/api/ingest-dirs/toggle")
    async def api_toggle_ingest_dir(request: Request):
        body = await request.json()
        path = body.get("path", "")
        enabled = body.get("enabled", True)
        settings_data = load_settings(workspace)
        dirs = settings_data.setdefault("ingest_dirs", {})
        if path in dirs:
            dirs[path]["enabled"] = enabled
            save_settings(workspace, settings_data)
        return {"status": "ok", "path": path, "enabled": enabled}

    @app.post("/api/ingest-dirs/remove")
    async def api_remove_ingest_dir(request: Request):
        """Remove a directory and all its files from the index.
        Responds immediately after removing from settings, cleans DB in background."""
        body = await request.json()
        path = body.get("path", "")
        if not path:
            raise HTTPException(400, "No path provided")

        # Step 1: Remove from settings IMMEDIATELY (fast, <1ms)
        settings_data = load_settings(workspace)
        dirs = settings_data.get("ingest_dirs", {})
        if path in dirs:
            del dirs[path]
            save_settings(workspace, settings_data)

        # Step 2: Clean DB in background thread (slow, but user doesn't wait)
        import threading
        def _cleanup():
            store = get_store()
            try:
                c = store.conn.cursor()
                c.execute("""SELECT file_id, file_path FROM files
                             WHERE source_dir = ? OR file_path LIKE ? OR source_dir LIKE ?""",
                          (path, path + "%", path + "%"))
                files = c.fetchall()
                for f in files:
                    try:
                        store.remove_file(f["file_path"])
                    except Exception:
                        pass
                # Clean graph edges
                file_ids = [f["file_id"] for f in files]
                if file_ids:
                    ph = ",".join("?" * len(file_ids))
                    c.execute(f"DELETE FROM document_edges WHERE source_file_id IN ({ph}) OR target_file_id IN ({ph})", file_ids + file_ids)
                    c.execute(f"DELETE FROM graph_node_positions WHERE file_id IN ({ph})", file_ids)
                    store.conn.commit()
                print(f"[KBase] Cleaned {len(files)} files from '{path}'")
            except Exception as e:
                print(f"[KBase] Cleanup error: {e}")
            finally:
                store.close()
        threading.Thread(target=_cleanup, daemon=True).start()

        return {"status": "ok", "path": path}

    # ---- Natural Language SQL API ----

    @app.post("/api/nl-sql")
    async def api_nl_sql(request: Request):
        """Convert natural language to SQL and execute."""
        body = await request.json()
        question = body.get("question", "")
        if not question:
            raise HTTPException(400, "question required")

        store = get_store()
        try:
            # Get table schemas for context
            table_ctx = get_table_context(store)
            schema_text = "\n".join(
                f"Table: {t['table_name']} ({t['row_count']} rows)\n  Columns: {', '.join(t['headers'])}"
                for t in table_ctx.get("tables", [])[:30]
            )

            # Use LLM to generate SQL
            settings_data = load_settings(workspace)
            provider_key = settings_data.get("llm_provider", "claude-sonnet")
            from kbase.chat import _call_llm, LLM_PROVIDERS
            provider = LLM_PROVIDERS.get(provider_key, {})

            prompt = f"""Given these SQLite tables:

{schema_text}

Convert this question to a SQL query. Return ONLY the SQL, no explanation.
Question: {question}"""

            try:
                sql = _call_llm(provider, [{"role": "user", "content": prompt}],
                                "You are a SQL expert. Return only valid SQLite SQL.", settings_data)
                # Clean up
                sql = sql.strip().strip('`').strip()
                if sql.lower().startswith('sql'):
                    sql = sql[3:].strip()

                # Execute
                result = store.sql_query(sql)
                return {"question": question, "sql": sql, "results": result}
            except Exception as e:
                return {"question": question, "error": f"LLM error: {str(e)[:200]}"}
        finally:
            store.close()

    # ---- Error details API ----

    @app.get("/api/errors")
    def api_errors():
        store = get_store()
        try:
            files = store.list_files()
            errors = [{"file_name": f["file_name"], "file_path": f["file_path"],
                        "file_type": f["file_type"], "error": f["error"]}
                       for f in files if f.get("error")]
            # Group by error type
            from collections import Counter
            summary = Counter()
            for e in errors:
                msg = e["error"][:80]
                summary[msg] = summary.get(msg, 0) + 1
            return {
                "total": len(errors),
                "errors": errors,
                "summary": dict(summary.most_common(20)),
            }
        finally:
            store.close()

    # ---- Chat API ----

    @app.post("/api/chat")
    async def api_chat(request: Request):
        body = await request.json()
        question = body.get("question", "")
        if not question:
            raise HTTPException(400, "question is required")
        conv_id = body.get("conversation_id", "default")
        store = get_store()
        try:
            settings_data = load_settings(workspace)
            settings_data.update(body.get("settings_override", {}))
            result = chat(
                store, question, settings=settings_data,
                top_k=body.get("top_k", 10),
                conversation_id=conv_id,
            )
            return result
        except (ValueError, RuntimeError) as e:
            return {"answer": f"LLM Error: {str(e)}", "sources": [], "error": True}
        except Exception as e:
            return {"answer": f"Unexpected error: {str(e)[:300]}", "sources": [], "error": True}
        finally:
            store.close()

    @app.get("/api/research-stream")
    def api_research_stream(question: str = Query(...), conv_id: str = Query("default")):
        """SSE stream for deep research with real-time progress."""
        import queue as _queue, threading
        q = _queue.Queue()
        settings_data = load_settings(workspace)

        def do_research():
            try:
                from kbase.agent_loop import AgentLoop
                from kbase.chat import _call_llm, LLM_PROVIDERS, _conversations, _load_conversations, _save_conversations
                from kbase.search import hybrid_search as hs

                if not _conversations:
                    _load_conversations(workspace)

                provider_key = settings_data.get("llm_provider", "claude-sonnet")
                provider = LLM_PROVIDERS.get(provider_key, LLM_PROVIDERS.get("claude-sonnet"))

                def llm_func(prompt):
                    return _call_llm(provider, [{"role": "user", "content": prompt}],
                                     "You are a research assistant.", settings_data)

                def kb_func(query):
                    store = get_store()
                    try:
                        r = hs(store, query, top_k=5, use_rerank=False)
                        return r.get("results", [])
                    finally:
                        store.close()

                agent = AgentLoop(llm_func, kb_func, max_rounds=10, urls_per_round=20, max_time_seconds=600)
                result = agent.run(question, progress_queue=q)

                # Save to conversation
                if conv_id not in _conversations:
                    _conversations[conv_id] = []
                _conversations[conv_id].append({"role": "user", "content": question})
                _conversations[conv_id].append({"role": "assistant", "content": result.get("report", "")})
                _save_conversations()

                # Auto-save research report to file and index it
                try:
                    report_text = result.get("report", "")
                    if report_text:
                        _auto_save_research(report_text, question, workspace, settings_data)
                except Exception as e:
                    print(f"[KBase] Research auto-save failed: {e}")

                q.put(json.dumps({
                    "type": "result",
                    "answer": result.get("report", ""),
                    "sources": [s for s in result.get("sources", []) if s.get("source") == "kb"][:8],
                    "web_sources": result.get("web_sources", [])[:10],
                    "stats": {
                        "rounds": result.get("rounds", 0),
                        "total_urls": result.get("total_urls", 0),
                        "elapsed": result.get("elapsed_seconds", 0),
                    },
                }))
            except Exception as e:
                q.put(json.dumps({"type": "error", "message": str(e)[:300]}))
            finally:
                q.put(None)

        threading.Thread(target=do_research, daemon=True).start()

        def stream():
            while True:
                msg = q.get()
                if msg is None:
                    break
                yield f"data: {msg}\n\n"

        return StreamingResponse(stream(), media_type="text/event-stream")

    @app.post("/api/chat/clear")
    async def api_chat_clear(request: Request):
        body = await request.json()
        conv_id = body.get("conversation_id", "default")
        clear_conversation(conv_id)
        return {"status": "cleared", "conversation_id": conv_id}

    @app.get("/api/llm-providers")
    def api_llm_providers():
        return {"providers": LLM_PROVIDERS, "buddy_presets": BUDDY_PRESETS}

    @app.get("/api/conversations")
    def api_conversations():
        """List all conversation sessions."""
        if not _conversations:
            _load_conversations(workspace)
        convs = []
        for cid, msgs in _conversations.items():
            if msgs:
                preview = msgs[0]["content"][:80] if msgs else ""
                title = _conv_titles.get(cid, "")
                convs.append({
                    "id": cid,
                    "turns": len(msgs) // 2,
                    "preview": preview,
                    "title": title,
                    "last_message": msgs[-1]["content"][:80] if msgs else "",
                })
        return {"conversations": convs}

    @app.post("/api/conversations/{conv_id}/generate-title")
    async def api_generate_title(conv_id: str):
        """Auto-generate a title for this conversation."""
        if not _conversations:
            _load_conversations(workspace)
        settings_data = load_settings(workspace)
        title = generate_title(conv_id, settings_data)
        return {"id": conv_id, "title": title}

    @app.put("/api/conversations/{conv_id}/title")
    async def api_set_title(conv_id: str, request: Request):
        """Manually set conversation title."""
        body = await request.json()
        title = body.get("title", "").strip()[:50]
        set_conv_title(conv_id, title)
        return {"id": conv_id, "title": title}

    @app.get("/api/conversations/{conv_id}")
    def api_get_conversation(conv_id: str):
        if not _conversations:
            _load_conversations(workspace)
        msgs = _conversations.get(conv_id, [])
        return {"id": conv_id, "messages": msgs, "turns": len(msgs) // 2}

    @app.delete("/api/conversations/{conv_id}")
    def api_delete_conversation(conv_id: str):
        clear_conversation(conv_id)
        return {"status": "deleted", "id": conv_id}

    @app.post("/api/chat/rewind")
    async def api_chat_rewind(request: Request):
        """Rewind conversation by removing the last N turns."""
        body = await request.json()
        conv_id = body.get("conversation_id", "default")
        turns = body.get("turns", 1)  # How many Q&A pairs to remove
        msgs = _conversations.get(conv_id, [])
        remove_count = turns * 2  # Each turn = 1 user + 1 assistant
        if remove_count > 0 and len(msgs) >= remove_count:
            _conversations[conv_id] = msgs[:-remove_count]
        elif msgs:
            _conversations[conv_id] = []
        _save_conversations()
        remaining = len(_conversations.get(conv_id, [])) // 2
        return {"status": "rewound", "turns_removed": turns, "remaining_turns": remaining}

    # ---- Model Status API ----

    @app.get("/api/model-status")
    def api_model_status():
        """Check which local models are already downloaded."""
        import shutil
        from pathlib import Path

        status = {}
        cache_dirs = [
            Path.home() / ".cache" / "huggingface" / "hub",
            Path.home() / ".cache" / "torch" / "sentence_transformers",
            Path.home() / ".cache" / "chroma" / "onnx_models",
        ]

        # Check embedding models
        for key, m in EMBEDDING_MODELS.items():
            if m.get("type") == "local":
                model_name = m["name"].replace("/", "--")
                found = any(
                    any(d.name.endswith(model_name) or model_name in d.name
                        for d in cd.iterdir() if d.is_dir())
                    for cd in cache_dirs if cd.exists()
                )
                status[key] = {"downloaded": found, "type": "embedding"}

        # Check whisper models
        whisper_cache = Path.home() / ".cache" / "whisper"
        for key, m in WHISPER_MODELS.items():
            if m.get("type") in ("local", "faster-whisper"):
                model_name = m["name"]
                found = (whisper_cache.exists() and
                         any(model_name in f.name for f in whisper_cache.iterdir())) if whisper_cache.exists() else False
                # Also check huggingface cache for faster-whisper
                if not found:
                    for cd in cache_dirs:
                        if cd.exists():
                            found = any(model_name in d.name for d in cd.iterdir() if d.is_dir())
                            if found:
                                break
                status[key] = {"downloaded": found, "type": "whisper"}

        # Check Ollama
        ollama_available = shutil.which("ollama") is not None
        status["ollama"] = {"downloaded": ollama_available, "type": "llm"}

        # Check CLI tools (claude, qwen, llm)
        # Expand PATH to find nvm/homebrew/cargo installed tools
        import os as _os
        extra_paths = ["/usr/local/bin", "/opt/homebrew/bin",
                       str(Path.home() / ".local" / "bin"),
                       str(Path.home() / ".npm-global" / "bin"),
                       str(Path.home() / ".cargo" / "bin")]
        nvm_dir = Path.home() / ".nvm" / "versions" / "node"
        if nvm_dir.exists():
            for node_ver in sorted(nvm_dir.iterdir(), reverse=True):
                extra_paths.append(str(node_ver / "bin"))
        orig_path = _os.environ.get("PATH", "")
        _os.environ["PATH"] = _os.pathsep.join(extra_paths) + _os.pathsep + orig_path

        cli_checks = {
            "claude-cli": "claude",
            "qwen-cli": "qwen",
            "llm-cli": "llm",
        }
        for key, cmd in cli_checks.items():
            status[key] = {"downloaded": shutil.which(cmd) is not None, "type": "llm"}

        _os.environ["PATH"] = orig_path  # Restore

        return {"status": status}

    # ---- Auto-Update API ----

    @app.get("/api/version")
    def api_version():
        """Return current version and check for updates."""
        from kbase import __version__
        import shutil
        is_git = Path(__file__).parent.parent.joinpath(".git").exists()
        is_frozen = getattr(sys, 'frozen', False)
        return {
            "version": __version__,
            "install_type": "dmg" if is_frozen else ("git" if is_git else "pip"),
            "is_frozen": is_frozen,
        }

    @app.get("/api/update/check")
    def api_update_check():
        """Check for new version from remote source."""
        from kbase import __version__
        import urllib.request
        settings = load_settings()
        DEFAULT_UPDATE_URL = "https://raw.githubusercontent.com/PenguinMiaou/kbase/main/version.json"
        update_url = settings.get("update_url", DEFAULT_UPDATE_URL).strip()
        if not update_url:
            update_url = DEFAULT_UPDATE_URL
        try:
            import ssl
            ctx = ssl.create_default_context()
            try:
                import certifi
                ctx.load_verify_locations(certifi.where())
            except ImportError:
                ctx.check_hostname = False
                ctx.verify_mode = ssl.CERT_NONE
            req = urllib.request.Request(update_url, headers={"User-Agent": "KBase"})
            with urllib.request.urlopen(req, timeout=5, context=ctx) as resp:
                import json as _json
                data = _json.loads(resp.read().decode())
            remote_ver = data.get("version", "0.0.0")
            download_url = data.get("download_url", "")
            changelog = data.get("changelog", "")
            # Simple tuple version compare (no packaging dependency)
            def ver_tuple(v):
                return tuple(int(x) for x in v.split(".")[:3] if x.isdigit())
            has_update = ver_tuple(remote_ver) > ver_tuple(__version__)
            return {
                "update_available": has_update,
                "current": __version__,
                "latest": remote_ver,
                "download_url": download_url,
                "changelog": changelog,
            }
        except Exception as e:
            return {"update_available": False, "current": __version__, "error": str(e)}

    @app.post("/api/update/apply")
    async def api_update_apply():
        """Apply update — git pull for source installs, download+install for binary."""
        import subprocess
        repo_dir = Path(__file__).parent.parent
        git_dir = repo_dir / ".git"
        is_frozen = getattr(sys, 'frozen', False)

        # Git install: git pull + pip install
        if git_dir.exists() and not is_frozen:
            try:
                result = subprocess.run(
                    ["git", "pull", "--ff-only"],
                    cwd=str(repo_dir), capture_output=True, text=True, timeout=30,
                )
                if result.returncode != 0:
                    return {"success": False, "message": result.stderr.strip()}
                pip_result = subprocess.run(
                    [sys.executable, "-m", "pip", "install", "-e", str(repo_dir), "-q"],
                    capture_output=True, text=True, timeout=120,
                )
                from kbase import __version__
                import importlib, kbase
                importlib.reload(kbase)
                return {
                    "success": True,
                    "message": f"Updated to {kbase.__version__}. Restart server to apply.",
                    "git_output": result.stdout.strip(),
                    "need_restart": True,
                }
            except subprocess.TimeoutExpired:
                return {"success": False, "message": "Update timed out"}
            except Exception as e:
                return {"success": False, "message": str(e)}

        # Binary install (DMG/EXE): download + prepare updater
        raise HTTPException(400, "Use /api/update/download for binary installs")

    @app.get("/api/update/download")
    async def api_update_download(request: Request):
        """Download and install update for binary installs (DMG/EXE). SSE stream."""
        import urllib.request, ssl, tempfile, threading

        settings = load_settings()
        DEFAULT_UPDATE_URL = "https://raw.githubusercontent.com/PenguinMiaou/kbase/main/version.json"
        update_url = settings.get("update_url", DEFAULT_UPDATE_URL).strip() or DEFAULT_UPDATE_URL

        def stream():
            try:
                # 1. Fetch version info
                yield f"data: {json.dumps({'stage': 'checking', 'message': 'Checking for updates...'})}\n\n"
                ctx = ssl.create_default_context()
                try:
                    import certifi
                    ctx.load_verify_locations(certifi.where())
                except ImportError:
                    ctx.check_hostname = False
                    ctx.verify_mode = ssl.CERT_NONE
                req = urllib.request.Request(update_url, headers={"User-Agent": "KBase"})
                with urllib.request.urlopen(req, timeout=10, context=ctx) as resp:
                    data = json.loads(resp.read().decode())

                is_mac = platform.system() == "Darwin"
                download_url = data.get("download_url_mac" if is_mac else "download_url_win", "")
                if not download_url:
                    download_url = data.get("download_url", "")
                if not download_url:
                    yield f"data: {json.dumps({'stage': 'error', 'message': 'No download URL in version manifest'})}\n\n"
                    return

                latest = data.get("version", "unknown")
                yield f"data: {json.dumps({'stage': 'downloading', 'message': f'Downloading v{latest}...', 'version': latest})}\n\n"

                # 2. Download the asset
                req2 = urllib.request.Request(download_url, headers={"User-Agent": "KBase"})
                with urllib.request.urlopen(req2, timeout=300, context=ctx) as resp2:
                    total = int(resp2.headers.get('Content-Length', 0))
                    ext = ".dmg" if is_mac else ".zip"
                    tmp_path = os.path.join(tempfile.gettempdir(), f"kbase-update{ext}")
                    downloaded = 0
                    with open(tmp_path, 'wb') as f:
                        while True:
                            chunk = resp2.read(1024 * 256)  # 256KB chunks
                            if not chunk:
                                break
                            f.write(chunk)
                            downloaded += len(chunk)
                            if total > 0:
                                pct = int(downloaded * 100 / total)
                                mb = downloaded / (1024 * 1024)
                                total_mb = total / (1024 * 1024)
                                yield f"data: {json.dumps({'stage': 'downloading', 'progress': pct, 'downloaded_mb': round(mb, 1), 'total_mb': round(total_mb, 1)})}\n\n"

                yield f"data: {json.dumps({'stage': 'downloaded', 'message': 'Download complete. Ready to install.', 'path': tmp_path})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'stage': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(stream(), media_type="text/event-stream")

    @app.post("/api/update/install")
    async def api_update_install():
        """Install downloaded update and restart."""
        import tempfile, threading
        is_mac = platform.system() == "Darwin"
        is_win = platform.system() == "Windows"
        ext = ".dmg" if is_mac else ".zip"
        tmp_path = os.path.join(tempfile.gettempdir(), f"kbase-update{ext}")

        if not os.path.exists(tmp_path):
            raise HTTPException(400, "No downloaded update found. Run download first.")

        if is_mac:
            # macOS: mount DMG, copy .app, unmount, restart
            updater_script = os.path.join(tempfile.gettempdir(), "kbase-updater.sh")
            with open(updater_script, 'w') as f:
                f.write(f"""#!/bin/bash
sleep 2
hdiutil attach "{tmp_path}" -nobrowse -quiet
VOL=$(ls /Volumes/ | grep KBase | head -1)
if [ -n "$VOL" ]; then
    rm -rf /Applications/KBase.app
    cp -R "/Volumes/$VOL/KBase.app" /Applications/
    hdiutil detach "/Volumes/$VOL" -quiet
fi
rm -f "{tmp_path}"
open /Applications/KBase.app
rm -f "{updater_script}"
""")
            os.chmod(updater_script, 0o755)
            subprocess.Popen(["/bin/bash", updater_script])

        elif is_win:
            # Windows: extract ZIP, replace files, restart
            import zipfile
            app_dir = os.path.dirname(sys.executable) if getattr(sys, 'frozen', False) else os.path.dirname(__file__)
            extract_dir = os.path.join(tempfile.gettempdir(), "kbase-update-extract")
            updater_script = os.path.join(tempfile.gettempdir(), "kbase-updater.bat")
            with open(updater_script, 'w') as f:
                f.write(f"""@echo off
timeout /t 3 /nobreak >nul
xcopy /s /e /y "{extract_dir}\\*" "{app_dir}\\" >nul 2>&1
rmdir /s /q "{extract_dir}" >nul 2>&1
del "{tmp_path}" >nul 2>&1
start "" "{os.path.join(app_dir, 'KBase.exe')}"
del "%~f0" >nul 2>&1
""")
            # Pre-extract ZIP
            with zipfile.ZipFile(tmp_path, 'r') as zf:
                zf.extractall(extract_dir)
            subprocess.Popen(["cmd", "/c", updater_script],
                             creationflags=0x00000008)  # DETACHED_PROCESS

        else:
            raise HTTPException(400, "Auto-install not supported on this platform")

        # Shutdown after launching updater
        def _do_shutdown():
            time.sleep(1)
            ppid = os.getppid()
            try:
                os.kill(ppid, signal.SIGTERM)
            except (ProcessLookupError, PermissionError):
                pass
            os.kill(os.getpid(), signal.SIGKILL)
        threading.Thread(target=_do_shutdown, daemon=True).start()
        return {"success": True, "message": "Installing update... KBase will restart."}

    # ---- File Preview API ----

    @app.get("/api/file-preview/{file_id}")
    def api_file_preview(file_id: str, max_chunks: int = Query(8)):
        """Get file content preview by file_id — returns original chunks from ChromaDB."""
        store = get_store()
        c = store.conn.cursor()
        c.execute("SELECT file_path, file_name, file_type, chunk_count, source_dir, summary FROM files WHERE file_id = ?", (file_id,))
        row = c.fetchone()
        if not row:
            raise HTTPException(404, "File not found")
        info = dict(row)
        import shutil
        info["can_convert"] = True  # always true — Python libs as fallback

        # Get original text from ChromaDB (not jieba-segmented FTS)
        try:
            results = store.collection.get(
                where={"file_id": file_id},
                include=["documents", "metadatas"],
            )
            # Sort by chunk_index, filter out parent chunks
            chunk_pairs = []
            for i, meta in enumerate(results.get("metadatas", [])):
                is_parent = meta.get("is_parent")
                if is_parent and str(is_parent).lower() in ("true", "1"):
                    continue
                idx = meta.get("chunk_index", i)
                chunk_pairs.append((idx, results["documents"][i], meta))
            chunk_pairs.sort(key=lambda x: x[0])
            chunks = [{"text": cp[1], "metadata": cp[2]} for cp in chunk_pairs[:max_chunks]]
        except Exception:
            chunks = []
        # Also get edges for this file
        c.execute("""
            SELECT e.edge_id, e.edge_type, e.label, e.direction, e.score, e.method,
                   CASE WHEN e.source_file_id = ? THEN f2.file_name ELSE f1.file_name END as neighbor_name,
                   CASE WHEN e.source_file_id = ? THEN e.target_file_id ELSE e.source_file_id END as neighbor_id
            FROM document_edges e
            LEFT JOIN files f1 ON e.source_file_id = f1.file_id
            LEFT JOIN files f2 ON e.target_file_id = f2.file_id
            WHERE e.source_file_id = ? OR e.target_file_id = ?
            ORDER BY e.score DESC LIMIT 10
        """, (file_id, file_id, file_id, file_id))
        edges = [dict(r) for r in c.fetchall()]
        return {**info, "chunks": chunks, "edges": edges}

    @app.get("/api/file-serve/{file_id}")
    def api_file_serve(file_id: str):
        """Serve the original file for preview (PDF/HTML/image)."""
        from fastapi.responses import FileResponse
        store = get_store()
        c = store.conn.cursor()
        c.execute("SELECT file_path, file_type FROM files WHERE file_id = ?", (file_id,))
        row = c.fetchone()
        if not row:
            raise HTTPException(404, "File not found")
        file_path = _validate_file_path(row["file_path"])
        if not os.path.isfile(file_path):
            raise HTTPException(404, "File no longer exists on disk")
        mime_map = {
            '.pdf': 'application/pdf',
            '.html': 'text/html', '.htm': 'text/html',
            '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
            '.gif': 'image/gif', '.svg': 'image/svg+xml', '.webp': 'image/webp',
            '.md': 'text/markdown', '.txt': 'text/plain',
            '.mp3': 'audio/mpeg', '.m4a': 'audio/mp4', '.wav': 'audio/wav',
        }
        ext = (row["file_type"] or "").lower()
        media_type = mime_map.get(ext, 'application/octet-stream')
        from starlette.responses import Response
        with open(file_path, 'rb') as f:
            content = f.read()
        headers = {"Content-Disposition": "inline"}  # inline, not download
        return Response(content=content, media_type=media_type, headers=headers)

    @app.get("/api/file-slides/{file_id}")
    def api_file_slides(file_id: str):
        """Convert PPTX to per-slide PNG images using LibreOffice."""
        import shutil, tempfile, glob as globmod
        store = get_store()
        c = store.conn.cursor()
        c.execute("SELECT file_path, file_type FROM files WHERE file_id = ?", (file_id,))
        row = c.fetchone()
        if not row:
            raise HTTPException(404, "File not found")
        file_path = _validate_file_path(row["file_path"])
        if not os.path.isfile(file_path):
            raise HTTPException(404, "File no longer exists")

        # Output directory per file
        slides_dir = os.path.join(tempfile.gettempdir(), "kbase-slides", file_id)
        os.makedirs(slides_dir, exist_ok=True)

        # Check if already converted
        existing = sorted(globmod.glob(os.path.join(slides_dir, "*.png")))
        if not existing:
            if not shutil.which("soffice"):
                raise HTTPException(400, "LibreOffice not installed")
            # Convert PPTX -> PDF first, then PDF -> PNG pages
            pdf_path = os.path.join(slides_dir, "slides.pdf")
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", slides_dir, file_path],
                capture_output=True, timeout=60,
            )
            # Find the generated PDF
            pdfs = globmod.glob(os.path.join(slides_dir, "*.pdf"))
            if pdfs:
                pdf_path = pdfs[0]
                # PDF -> PNG per page using PyMuPDF
                try:
                    import fitz
                    doc = fitz.open(pdf_path)
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=150)
                        pix.save(os.path.join(slides_dir, f"slide_{i:03d}.png"))
                    doc.close()
                except Exception as e:
                    raise HTTPException(500, f"PNG conversion failed: {e}")
            existing = sorted(globmod.glob(os.path.join(slides_dir, "*.png")))

        if not existing:
            raise HTTPException(500, "No slides generated")

        return {"total": len(existing), "slides": [
            f"/api/file-slide-img/{file_id}/{i}" for i in range(len(existing))
        ]}

    @app.get("/api/file-slide-img/{file_id}/{index}")
    def api_file_slide_img(file_id: str, index: int):
        """Serve a single slide PNG image."""
        import tempfile, glob as globmod
        slides_dir = os.path.join(tempfile.gettempdir(), "kbase-slides", file_id)
        pngs = sorted(globmod.glob(os.path.join(slides_dir, "*.png")))
        if index < 0 or index >= len(pngs):
            raise HTTPException(404, "Slide not found")
        with open(pngs[index], 'rb') as f:
            content = f.read()
        from starlette.responses import Response
        return Response(content=content, media_type="image/png")

    @app.get("/api/file-xlsx/{file_id}")
    def api_file_xlsx(file_id: str):
        """Parse XLSX into JSON for Luckysheet rendering."""
        store = get_store()
        c = store.conn.cursor()
        c.execute("SELECT file_path FROM files WHERE file_id = ?", (file_id,))
        row = c.fetchone()
        if not row:
            raise HTTPException(404, "File not found")
        file_path = row["file_path"]
        if not os.path.isfile(file_path):
            raise HTTPException(404, "File no longer exists")

        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for idx, name in enumerate(wb.sheetnames[:10]):
            ws = wb[name]
            celldata = []
            for r, row_data in enumerate(ws.iter_rows(max_row=500, values_only=False)):
                for c_idx, cell in enumerate(row_data):
                    if cell.value is not None:
                        entry = {"r": r, "c": c_idx, "v": {"v": cell.value, "m": str(cell.value)}}
                        if cell.font and cell.font.bold:
                            entry["v"]["bl"] = 1
                        celldata.append(entry)
            sheets.append({
                "name": name,
                "index": idx,
                "order": idx,
                "status": 1 if idx == 0 else 0,
                "celldata": celldata,
                "config": {},
            })
        wb.close()
        return sheets

    @app.get("/api/file-convert/{file_id}")
    def api_file_convert(file_id: str):
        """Convert PPTX/DOCX/XLSX to PDF (LibreOffice) or HTML (fallback)."""
        import shutil, tempfile
        store = get_store()
        c = store.conn.cursor()
        c.execute("SELECT file_path, file_type FROM files WHERE file_id = ?", (file_id,))
        row = c.fetchone()
        if not row:
            raise HTTPException(404, "File not found")
        file_path = row["file_path"]
        ext = (row["file_type"] or "").lower()
        if not os.path.isfile(file_path):
            raise HTTPException(404, "File no longer exists")

        # Try LibreOffice first — renders PPTX/DOCX faithfully as PDF
        if shutil.which("soffice"):
            tmp_dir = os.path.join(tempfile.gettempdir(), "kbase-convert")
            os.makedirs(tmp_dir, exist_ok=True)
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            pdf_out = os.path.join(tmp_dir, base_name + ".pdf")
            if not os.path.exists(pdf_out):
                try:
                    subprocess.run(
                        ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, file_path],
                        capture_output=True, timeout=30,
                    )
                except Exception:
                    pass
            if os.path.exists(pdf_out):
                with open(pdf_out, 'rb') as f:
                    content = f.read()
                from starlette.responses import Response
                return Response(content=content, media_type="application/pdf",
                                headers={"Content-Disposition": "inline"})

        html_parts = ['<html><head><meta charset="utf-8"><style>',
            'body{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif;padding:20px;color:#1f2937;line-height:1.6;max-width:800px;margin:0 auto;}',
            'h1,h2,h3{color:#4f46e5;margin:16px 0 8px;}',
            'table{border-collapse:collapse;width:100%;margin:12px 0;}',
            'th,td{border:1px solid #e5e7eb;padding:6px 10px;text-align:left;font-size:13px;}',
            'th{background:#f3f4f6;font-weight:600;}',
            '.slide{border:1px solid #e5e7eb;border-radius:8px;padding:16px 20px;margin:16px 0;background:#fafafa;}',
            '.slide-num{color:#6366f1;font-weight:700;font-size:12px;margin-bottom:8px;}',
            'img{max-width:100%;border-radius:4px;}',
            'p{margin:4px 0;}',
            '</style></head><body>']

        try:
            if ext in ('.pptx', '.ppt'):
                from pptx import Presentation
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    html_parts.append(f'<div class="slide"><div class="slide-num">Slide {i+1}</div>')
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if not text:
                                    continue
                                if para.level == 0 and len(text) < 80:
                                    html_parts.append(f'<h3>{text}</h3>')
                                else:
                                    html_parts.append(f'<p>{text}</p>')
                        if shape.has_table:
                            tbl = shape.table
                            html_parts.append('<table>')
                            for r, row in enumerate(tbl.rows):
                                tag = 'th' if r == 0 else 'td'
                                html_parts.append('<tr>' + ''.join(f'<{tag}>{cell.text}</{tag}>' for cell in row.cells) + '</tr>')
                            html_parts.append('</table>')
                    html_parts.append('</div>')

            elif ext in ('.docx', '.doc'):
                from docx import Document
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    style = para.style.name.lower() if para.style else ''
                    if 'heading 1' in style:
                        html_parts.append(f'<h1>{text}</h1>')
                    elif 'heading 2' in style:
                        html_parts.append(f'<h2>{text}</h2>')
                    elif 'heading 3' in style:
                        html_parts.append(f'<h3>{text}</h3>')
                    else:
                        # Check for bold runs
                        runs_html = ''
                        for run in para.runs:
                            t = run.text
                            if run.bold:
                                t = f'<strong>{t}</strong>'
                            if run.italic:
                                t = f'<em>{t}</em>'
                            runs_html += t
                        html_parts.append(f'<p>{runs_html or text}</p>')
                for table in doc.tables:
                    html_parts.append('<table>')
                    for r, row in enumerate(table.rows):
                        tag = 'th' if r == 0 else 'td'
                        html_parts.append('<tr>' + ''.join(f'<{tag}>{cell.text}</{tag}>' for cell in row.cells) + '</tr>')
                    html_parts.append('</table>')

            elif ext in ('.xlsx', '.xls', '.csv'):
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True, data_only=True)
                for sheet in wb.sheetnames[:5]:  # max 5 sheets
                    ws = wb[sheet]
                    html_parts.append(f'<h2>{sheet}</h2><table>')
                    for r, row in enumerate(ws.iter_rows(max_row=100, values_only=True)):
                        tag = 'th' if r == 0 else 'td'
                        html_parts.append('<tr>' + ''.join(f'<{tag}>{cell if cell is not None else ""}</{tag}>' for cell in row) + '</tr>')
                    html_parts.append('</table>')
                wb.close()
            else:
                raise HTTPException(400, f"Unsupported format: {ext}")

        except Exception as e:
            html_parts.append(f'<p style="color:red;">Preview error: {str(e)[:200]}</p>')

        html_parts.append('</body></html>')
        return HTMLResponse(''.join(html_parts))

    # ---- Knowledge Graph API ----

    @app.get("/api/graph")
    def api_graph(
        min_score: float = Query(0.0),
        edge_type: str = Query(None),
        file_type: str = Query(None),
        source_dir: str = Query(None),
    ):
        """Get full graph data for visualization."""
        from kbase.graph import get_graph_data
        store = get_store()
        edge_types = edge_type.split(",") if edge_type else None
        return get_graph_data(store, edge_types=edge_types, min_score=min_score,
                             file_type=file_type, source_dir=source_dir)

    @app.get("/api/graph/local/{file_id}")
    def api_graph_local(file_id: str, depth: int = Query(2), min_score: float = Query(0.0)):
        """Get local subgraph centered on a file."""
        from kbase.graph import get_local_graph
        store = get_store()
        return get_local_graph(store, file_id, depth=depth, min_score=min_score)

    @app.post("/api/graph/compute")
    def api_graph_compute():
        """Trigger graph relationship computation."""
        from kbase.graph import compute_graph
        store = get_store()
        result = compute_graph(store, threshold=0.65)
        return result

    @app.post("/api/summaries/generate")
    def api_summaries_generate():
        """Backfill: generate LLM summaries for files that don't have one yet.
        Returns SSE stream with progress."""
        import queue as _queue, threading
        q = _queue.Queue()
        settings_data = load_settings(workspace)

        def do_generate():
            try:
                store = get_store()
                files = store.get_files_without_summary(limit=50)
                if not files:
                    q.put(json.dumps({"type": "done", "generated": 0, "message": "All files already have summaries"}))
                    store.close()
                    return

                from kbase.chat import generate_document_summary
                generated = 0
                for i, f in enumerate(files):
                    try:
                        # Get file text from ChromaDB
                        results = store.collection.get(
                            where={"file_id": f["file_id"]},
                            include=["documents"],
                        )
                        text = "\n".join(results.get("documents", [])[:5])
                        if not text:
                            continue

                        summary = generate_document_summary(text, f["file_name"], settings_data)
                        if summary:
                            store.update_file_summary(f["file_id"], summary)
                            generated += 1

                        q.put(json.dumps({
                            "type": "progress",
                            "current": i + 1,
                            "total": len(files),
                            "file": f["file_name"],
                            "has_summary": bool(summary),
                        }))
                    except Exception as e:
                        q.put(json.dumps({"type": "error", "file": f["file_name"], "message": str(e)[:200]}))

                q.put(json.dumps({"type": "done", "generated": generated, "total": len(files)}))
                store.close()
            except Exception as e:
                q.put(json.dumps({"type": "error", "message": str(e)[:300]}))
            finally:
                q.put(None)

        threading.Thread(target=do_generate, daemon=True).start()

        def stream():
            while True:
                msg = q.get()
                if msg is None:
                    break
                yield f"data: {msg}\n\n"

        return StreamingResponse(stream(), media_type="text/event-stream")

    @app.get("/api/graph/stats")
    def api_graph_stats():
        """Get graph statistics."""
        from kbase.graph import get_graph_stats
        store = get_store()
        return get_graph_stats(store)

    @app.post("/api/graph/edge")
    async def api_graph_edge_create(request: Request):
        """Create or confirm an edge."""
        from kbase.graph import add_edge
        store = get_store()
        body = await request.json()
        return add_edge(
            store,
            source_id=body["source"],
            target_id=body["target"],
            edge_type=body.get("edge_type", "confirmed"),
            label=body.get("label", ""),
            direction=body.get("direction", "forward"),
        )

    @app.put("/api/graph/edge/{edge_id}")
    async def api_graph_edge_update(edge_id: str, request: Request):
        """Update an edge."""
        from kbase.graph import update_edge
        store = get_store()
        body = await request.json()
        return update_edge(store, edge_id, **body)

    @app.delete("/api/graph/edge/{edge_id}")
    def api_graph_edge_delete(edge_id: str):
        """Delete an edge."""
        from kbase.graph import delete_edge
        store = get_store()
        return delete_edge(store, edge_id)

    @app.put("/api/graph/positions")
    async def api_graph_positions(request: Request):
        """Save node positions for canvas mode."""
        from kbase.graph import save_positions
        store = get_store()
        body = await request.json()
        return save_positions(store, body.get("positions", []))

    # ---- Glossary API ----

    @app.get("/api/glossary")
    def api_get_glossary():
        from kbase.enhance import get_glossary
        return get_glossary()

    @app.post("/api/glossary")
    async def api_add_glossary(request: Request):
        body = await request.json()
        term = body.get("term", "").strip()
        synonyms = body.get("synonyms", [])
        if not term:
            raise HTTPException(400, "term is required")
        from kbase.enhance import add_glossary_term
        add_glossary_term(term, synonyms)
        return {"status": "added", "term": term}

    @app.delete("/api/glossary/{term}")
    def api_delete_glossary(term: str):
        from kbase.enhance import remove_glossary_term
        remove_glossary_term(term)
        return {"status": "removed", "term": term}

    @app.post("/api/glossary/extract")
    async def api_extract_glossary(request: Request):
        """Auto-extract glossary from indexed documents using LLM."""
        from kbase.enhance import auto_build_glossary
        from kbase.chat import _call_llm, LLM_PROVIDERS
        settings_data = load_settings(workspace)
        provider_key = settings_data.get("llm_provider", "claude-sonnet")
        provider = LLM_PROVIDERS.get(provider_key, LLM_PROVIDERS.get("claude-sonnet"))

        def llm_func(prompt):
            return _call_llm(provider, [{"role": "user", "content": prompt}], "", settings_data)

        # Get sample texts from indexed files
        store = get_store()
        try:
            files = store.list_files()
            texts = []
            for f in files[:20]:  # Sample up to 20 files
                chunks = store.get_file_chunks(f["file_path"]) if hasattr(store, "get_file_chunks") else []
                if chunks:
                    texts.append(" ".join(c.get("text", "")[:500] for c in chunks[:3]))
            new_count = auto_build_glossary(texts, llm_func)
            return {"status": "ok", "new_terms": new_count}
        except Exception as e:
            return {"status": "error", "message": str(e)[:300]}
        finally:
            store.close()

    # ---- Global Memory API ----

    @app.get("/api/memories")
    def api_get_memories():
        if not _conversations:
            _load_conversations(workspace)
        return {"memories": get_memories()}

    @app.post("/api/memories")
    async def api_add_memory(request: Request):
        body = await request.json()
        content = body.get("content", "").strip()
        if not content:
            raise HTTPException(400, "content is required")
        entry = add_memory(content, source=body.get("source", "manual"))
        return entry

    @app.delete("/api/memories/{mem_id}")
    def api_delete_memory(mem_id: str):
        delete_memory(mem_id)
        return {"status": "deleted", "id": mem_id}

    @app.post("/api/memories/extract/{conv_id}")
    def api_extract_memories(conv_id: str):
        """Auto-extract memories from a conversation."""
        if not _conversations:
            _load_conversations(workspace)
        settings_data = load_settings(workspace)
        new_mems = extract_memories_from_conversation(conv_id, settings_data)
        return {"extracted": new_mems, "total": len(get_memories())}

    # ---- Search Feedback / Click Tracking API (Harness Sensor) ----

    @app.post("/api/feedback/click")
    async def api_feedback_click(request: Request):
        """Record when user clicks a search result."""
        body = await request.json()
        store = get_store()
        try:
            store.record_click(
                query=body.get("query", ""),
                file_id=body.get("file_id", ""),
                file_name=body.get("file_name", ""),
                position=body.get("position", 0),
            )
            return {"status": "ok"}
        finally:
            store.close()

    @app.post("/api/feedback/rate")
    async def api_feedback_rate(request: Request):
        """Record thumbs up/down on an answer."""
        body = await request.json()
        store = get_store()
        try:
            store.record_feedback(
                query=body.get("query", ""),
                file_id=body.get("file_id", ""),
                action=body.get("action", "thumbs_up"),
            )
            return {"status": "ok"}
        finally:
            store.close()

    @app.get("/api/user-interests")
    def api_user_interests():
        """Get top user query interests (lightweight memory)."""
        store = get_store()
        try:
            return {"interests": store.get_top_interests(20)}
        finally:
            store.close()

    # ---- Directory Browser API ----

    @app.get("/api/browse")
    def api_browse(path: str = Query("~")):
        """List directories for the file browser."""
        import os
        target = os.path.expanduser(path)
        if not os.path.isdir(target):
            target = os.path.dirname(target)
        if not os.path.isdir(target):
            target = os.path.expanduser("~")
        try:
            entries = []
            # Parent directory
            parent = os.path.dirname(target)
            if parent != target:
                entries.append({"name": "..", "path": parent, "type": "parent"})
            for name in sorted(os.listdir(target)):
                full = os.path.join(target, name)
                if os.path.isdir(full) and not name.startswith("."):
                    entries.append({"name": name, "path": full, "type": "dir"})
            return {"current": target, "entries": entries}
        except PermissionError:
            return {"current": target, "entries": [], "error": "Permission denied"}

    # ---- Connectors API ----

    @app.get("/api/connectors")
    def api_connectors():
        from kbase.connectors.feishu import CONNECTORS
        settings_data = load_settings(workspace)
        # Add connection status
        result = {}
        for k, v in CONNECTORS.items():
            c = dict(v)
            connector_settings = settings_data.get(f"connector_{k}", {})
            c["connected"] = bool(connector_settings.get("connected"))
            result[k] = c
        return {"connectors": result}

    @app.post("/api/connectors/{name}/connect")
    async def api_connector_connect(name: str, request: Request):
        body = await request.json()
        settings_data = load_settings(workspace)
        existing = settings_data.get(f"connector_{name}", {})
        # Only mark as "configured", not "connected" — connected requires OAuth
        existing.update(body)
        existing["configured"] = True
        existing.setdefault("connected", False)  # Don't override if already connected via OAuth
        settings_data[f"connector_{name}"] = existing
        save_settings(workspace, settings_data)
        return {"status": "saved", "connector": name}

    @app.get("/api/connectors/feishu/guide", response_class=HTMLResponse)
    def api_feishu_guide():
        from kbase.connectors.feishu_guide import FEISHU_GUIDE_HTML
        # Wrap in a styled page
        return f"""<!DOCTYPE html><html><head><meta charset="UTF-8"><title>Feishu Setup Guide — KBase</title>
        <link rel="icon" href="/static/logos/kbase-logo.svg">
        <style>body{{font-family:-apple-system,sans-serif;background:#0f172a;color:#e2e8f0;padding:40px;}}
        a{{color:#818cf8;}} code{{background:rgba(99,102,241,0.15);padding:2px 6px;border-radius:4px;font-size:13px;}}
        h2,h3{{color:#e2e8f0;}} table{{width:100%;}} td,th{{text-align:left;}}
        @media(prefers-color-scheme:light){{body{{background:#f8fafc;color:#1e293b;}}h2,h3{{color:#1e293b;}}code{{background:#e2e8f0;}}}}</style>
        </head><body>{FEISHU_GUIDE_HTML}</body></html>"""

    @app.get("/api/connectors/feishu/oauth-url")
    def api_feishu_oauth_url(
        redirect_uri: str = Query("http://localhost:8765/api/connectors/feishu/callback"),
        scopes: str = Query(""),
    ):
        settings_data = load_settings(workspace)
        cfg = settings_data.get("connector_feishu", {})
        if not cfg.get("app_id"):
            raise HTTPException(400, "Feishu App ID not configured")
        from kbase.connectors.feishu import FeishuConnector
        fc = FeishuConnector(
            cfg["app_id"], cfg["app_secret"], workspace,
            use_lark=cfg.get("use_lark", False),
        )
        scope_list = scopes.split(",") if scopes else None
        url = fc.get_oauth_url(redirect_uri, scopes=scope_list)
        return {"oauth_url": url, "redirect_uri": redirect_uri}

    @app.get("/api/connectors/feishu/callback")
    def api_feishu_callback(code: str = Query(...), state: str = Query(None)):
        """OAuth callback from Feishu."""
        settings_data = load_settings(workspace)
        cfg = settings_data.get("connector_feishu", {})
        from kbase.connectors.feishu import FeishuConnector
        fc = FeishuConnector(cfg["app_id"], cfg["app_secret"], workspace,
                             use_lark=cfg.get("use_lark", False), custom_domain=cfg.get("custom_domain", ""))
        token_data = fc.exchange_code(code)
        settings_data["connector_feishu"]["connected"] = True
        settings_data["connector_feishu"]["user_name"] = token_data.get("name", "")
        save_settings(workspace, settings_data)
        return HTMLResponse("<h2>Feishu Connected! You can close this tab.</h2><script>window.close()</script>")

    @app.get("/api/connectors/feishu/debug")
    def api_feishu_debug():
        """Debug: test each Feishu API endpoint."""
        settings_data = load_settings(workspace)
        cfg = settings_data.get("connector_feishu", {})
        from kbase.connectors.feishu import FeishuConnector
        fc = FeishuConnector(cfg.get("app_id",""), cfg.get("app_secret",""), workspace,
                             use_lark=cfg.get("use_lark", False))
        results = {"has_user_token": bool(fc._user_token)}
        # Test each API
        import urllib.error
        for name, path in [
            ("list_chats", "/im/v1/chats?page_size=1"),
            ("list_msgs_test", "/im/v1/messages?container_id_type=chat&container_id=test&page_size=1"),
            ("list_drive", "/drive/v1/files?page_size=1"),
            ("root_folder", "/drive/explorer/v2/root_folder/meta"),
            ("user_info", "/authen/v1/user_info"),
        ]:
            try:
                r = fc._api("GET", path)
                results[name] = {"code": r.get("code"), "msg": r.get("msg", "")[:100]}
            except Exception as e:
                results[name] = {"error": str(e)[:200]}
        return results

    @app.post("/api/connectors/feishu/sync")
    def api_feishu_sync():
        """Sync all Feishu data and ingest."""
        settings_data = load_settings(workspace)
        cfg = settings_data.get("connector_feishu", {})
        if not cfg.get("app_id"):
            raise HTTPException(400, "Not configured")
        from kbase.connectors.feishu import FeishuConnector
        fc = FeishuConnector(cfg["app_id"], cfg["app_secret"], workspace,
                             use_lark=cfg.get("use_lark", False))
        # Pass selected scopes to determine which modules to sync
        selected_scopes = cfg.get("selected_scopes", [])
        sync_docs = any("doc" in s or "drive" in s for s in selected_scopes)
        sync_chats = any("im:" in s or "chat" in s for s in selected_scopes)
        sync_emails = any("mail" in s for s in selected_scopes)
        sync_stats = fc.sync_all(sync_docs=sync_docs, sync_chats=sync_chats, sync_emails=sync_emails)
        # Auto-ingest the synced files
        if sync_stats.get("output_dir"):
            store = get_store()
            try:
                ingest_stats = ingest_directory(store, sync_stats["output_dir"])
                sync_stats["ingest"] = ingest_stats
            finally:
                store.close()
        return sync_stats

    # ---- File Preview API ----

    @app.get("/api/preview")
    def api_preview(path: str = Query(...), page: int = Query(0), slide: int = Query(0), max_chars: int = Query(2000)):
        """Get preview — image for PDF/PPTX pages, text for others."""
        p = Path(path)
        if not p.exists():
            return {"preview": "File not found", "path": path, "type": "text"}
        ext = p.suffix.lower()
        try:
            # PDF: render page as image
            if ext == ".pdf" and page > 0:
                import fitz, base64
                doc = fitz.open(str(p))
                pg = doc.load_page(min(page - 1, len(doc) - 1))
                pix = pg.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                img_bytes = pix.tobytes("png")
                doc.close()
                return {
                    "image": base64.b64encode(img_bytes).decode(),
                    "title": f"{p.name} — Page {page}",
                    "type": "image", "path": str(p),
                }

            # PPTX: try to render slide via PDF conversion
            if ext == ".pptx" and slide > 0:
                import fitz, base64, subprocess, tempfile
                # Try LibreOffice conversion
                try:
                    with tempfile.TemporaryDirectory() as tmpdir:
                        subprocess.run(
                            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, str(p)],
                            capture_output=True, timeout=30,
                        )
                        pdf_path = Path(tmpdir) / f"{p.stem}.pdf"
                        if pdf_path.exists():
                            doc = fitz.open(str(pdf_path))
                            pg = doc.load_page(min(slide - 1, len(doc) - 1))
                            pix = pg.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                            img_bytes = pix.tobytes("png")
                            doc.close()
                            return {
                                "image": base64.b64encode(img_bytes).decode(),
                                "title": f"{p.name} — Slide {slide}",
                                "type": "image", "path": str(p),
                            }
                except Exception:
                    pass  # Fall through to text preview

            # Fallback: text preview
            from kbase.extract import extract_file
            result = extract_file(str(p))
            text = result.get("text", "")[:max_chars]
            return {
                "preview": text,
                "title": result.get("metadata", {}).get("title", p.name),
                "type": "text", "path": str(p),
            }
        except Exception as e:
            return {"preview": f"Error: {e}", "path": path, "type": "text"}

    @app.post("/api/open-file")
    async def api_open_file(request: Request):
        body = await request.json()
        file_path = _validate_file_path(body.get("path", ""))
        page = body.get("page", 0)
        slide = body.get("slide", 0)
        if not file_path or not Path(file_path).exists():
            raise HTTPException(404, "File not found")
        system = platform.system()
        ext = Path(file_path).suffix.lower()
        try:
            if system == "Darwin":
                if ext == ".pdf" and page > 0:
                    # macOS Preview can open to specific page
                    subprocess.Popen(["open", "-a", "Preview", file_path])
                else:
                    subprocess.Popen(["open", "-R", file_path])
            elif system == "Windows":
                subprocess.Popen(["explorer", "/select,", file_path])
            else:
                subprocess.Popen(["xdg-open", str(Path(file_path).parent)])
            return {"status": "opened", "path": file_path}
        except Exception as e:
            raise HTTPException(500, str(e))

    # ---- Model Download API ----

    @app.get("/api/model-status/check")
    def api_model_status_single(model_name: str = Query(...)):
        """Check if a single local model is already downloaded."""
        try:
            from sentence_transformers import SentenceTransformer
            import huggingface_hub
            # Check if model exists in cache
            try:
                path = huggingface_hub.snapshot_download(model_name, local_files_only=True)
                return {"status": "downloaded", "model": model_name, "path": str(path)}
            except Exception:
                return {"status": "not_downloaded", "model": model_name}
        except ImportError:
            return {"status": "unknown", "model": model_name}

    @app.get("/api/model-download")
    def api_model_download(model_name: str = Query(...)):
        """Download a model with SSE progress."""
        import threading, queue
        q = queue.Queue()

        def do_download():
            try:
                import subprocess
                q.put(json.dumps({"status": "downloading", "message": f"Downloading {model_name}...", "progress": 10}))

                # Use system python3 (not frozen binary) to download via huggingface_hub
                # This works even in DMG mode because macOS ships with python3
                sys_python = "/usr/bin/python3"
                if not Path(sys_python).exists():
                    import shutil
                    sys_python = shutil.which("python3") or "python3"

                # Step 1: Ensure huggingface_hub is installed
                q.put(json.dumps({"status": "downloading", "message": "Setting up download tools...", "progress": 20}))
                subprocess.run(
                    [sys_python, "-m", "pip", "install", "--user", "huggingface_hub", "-q"],
                    capture_output=True, text=True, timeout=120,
                )

                # Step 2: Download model using huggingface_hub snapshot_download
                q.put(json.dumps({"status": "downloading", "message": f"Downloading {model_name} (this may take a few minutes)...", "progress": 40}))
                dl_script = f"""
import huggingface_hub
path = huggingface_hub.snapshot_download('{model_name}')
print(path)
"""
                result = subprocess.run(
                    [sys_python, "-c", dl_script],
                    capture_output=True, text=True, timeout=600,
                )
                if result.returncode != 0:
                    q.put(json.dumps({"status": "error", "message": f"Download failed: {result.stderr[:300]}"}))
                    return

                model_path = result.stdout.strip()
                q.put(json.dumps({"status": "done", "message": f"Model downloaded to {model_path}. Restart KBase to use it.", "progress": 100}))
            except subprocess.TimeoutExpired:
                q.put(json.dumps({"status": "error", "message": "Download timed out (10 min limit)"}))
            except Exception as e:
                q.put(json.dumps({"status": "error", "message": str(e)[:300]}))
            finally:
                q.put(None)

        threading.Thread(target=do_download, daemon=True).start()

        def stream():
            while True:
                msg = q.get()
                if msg is None:
                    break
                yield f"data: {msg}\n\n"

        return StreamingResponse(stream(), media_type="text/event-stream")

    # ---- Settings API ----

    @app.get("/api/settings")
    def api_get_settings():
        settings = load_settings(workspace)
        # Security: mask API keys in response (show last 4 chars only)
        masked = dict(settings)
        for key in list(masked.keys()):
            if "api_key" in key.lower() or "secret" in key.lower():
                val = masked[key]
                if isinstance(val, str) and len(val) > 8:
                    masked[key] = "***" + val[-4:]
        from kbase.config import VISION_MODELS
        return {
            "settings": masked,
            "embedding_models": EMBEDDING_MODELS,
            "whisper_models": WHISPER_MODELS,
            "vision_models": VISION_MODELS,
            "llm_providers": LLM_PROVIDERS,
            "buddy_presets": BUDDY_PRESETS,
            "language_profiles": LANGUAGE_PROFILES,
            "defaults": {"chunk_max_chars": CHUNK_MAX_CHARS, "chunk_overlap_chars": CHUNK_OVERLAP_CHARS},
        }

    @app.post("/api/settings")
    async def api_save_settings(request: Request):
        body = await request.json()
        current = load_settings(workspace)
        # Security: don't overwrite real API keys with masked values (***xxxx)
        for key in list(body.keys()):
            if ("api_key" in key.lower() or "secret" in key.lower()):
                val = body[key]
                if isinstance(val, str) and val.startswith("***"):
                    del body[key]  # Keep the existing real key
        current.update(body)
        save_settings(workspace, current)
        return {"status": "saved", "settings": current}

    # ---- Shutdown ----

    @app.post("/api/shutdown")
    def api_shutdown():
        """Gracefully shutdown the entire app (server + launcher)."""
        import threading
        def _do_shutdown():
            time.sleep(0.5)
            # Kill parent process (launcher) which will also kill this server
            ppid = os.getppid()
            try:
                os.kill(ppid, signal.SIGTERM)
            except ProcessLookupError:
                pass
            # Kill self
            os.kill(os.getpid(), signal.SIGKILL)
        threading.Thread(target=_do_shutdown, daemon=True).start()
        return {"status": "shutting down"}

    # ---- Frontend ----

    @app.get("/", response_class=HTMLResponse)
    def index():
        # Serve new Claude-style UI
        new_ui = Path(__file__).parent / "static" / "index.html"
        if new_ui.exists():
            return HTMLResponse(new_ui.read_text(encoding="utf-8"))
        return FRONTEND_HTML

    @app.get("/v1", response_class=HTMLResponse)
    def index_v1():
        """Legacy UI"""
        return FRONTEND_HTML

    return app


# =========================================================================
# Frontend HTML
# =========================================================================
FRONTEND_HTML = r"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<link rel="icon" href="/static/logos/kbase-logo.svg">
<title>KBase</title>
<link href="https://cdn.jsdelivr.net/npm/tailwindcss@2.2.19/dist/tailwind.min.css" rel="stylesheet">
<style>
  :root { --accent: #818cf8; --accent2: #a78bfa; --bg: #0a0e1a; --card: rgba(15, 23, 42, 0.85); }
  body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background: var(--bg); color: #e2e8f0; overflow-x: hidden; }
  body::before { content:''; position:fixed; top:-50%; left:-50%; width:200%; height:200%; background: radial-gradient(ellipse at 20% 50%, rgba(99,102,241,0.08) 0%, transparent 50%), radial-gradient(ellipse at 80% 20%, rgba(168,85,247,0.06) 0%, transparent 50%), radial-gradient(ellipse at 50% 80%, rgba(59,130,246,0.05) 0%, transparent 50%); z-index:-1; animation: bgShift 20s ease infinite; }
  @keyframes bgShift { 0%,100%{transform:rotate(0deg)} 50%{transform:rotate(3deg)} }
  .glass { background: var(--card); backdrop-filter: blur(16px) saturate(180%); border: 1px solid rgba(148,163,184,0.08); box-shadow: 0 4px 30px rgba(0,0,0,0.3); }
  .glass:hover { border-color: rgba(148,163,184,0.15); }
  .result-card { transition: all 0.25s cubic-bezier(0.4,0,0.2,1); }
  .result-card:hover { background: rgba(51,65,85,0.5); transform: translateX(4px); border-left: 2px solid var(--accent); }
  .sinput { background: rgba(15,23,42,0.9); border: 1px solid rgba(99,102,241,0.2); transition: all 0.3s; }
  .sinput:focus { border-color: var(--accent); box-shadow: 0 0 0 3px rgba(99,102,241,0.1), 0 0 20px rgba(99,102,241,0.05); }
  .tab-active { border-bottom: 2px solid var(--accent); color: var(--accent); text-shadow: 0 0 10px rgba(129,140,248,0.3); }
  .hl { background: rgba(250,204,21,0.15); padding: 1px 4px; border-radius: 3px; border-bottom: 1px solid rgba(250,204,21,0.3); }
  .fade-in { animation: fadeIn 0.4s cubic-bezier(0.4,0,0.2,1); }
  @keyframes fadeIn { from { opacity:0; transform:translateY(12px) } to { opacity:1; transform:translateY(0) } }
  .stat-card { background: linear-gradient(135deg, rgba(99,102,241,0.08), rgba(168,85,247,0.08)); transition: all 0.3s; }
  .stat-card:hover { transform: translateY(-2px); box-shadow: 0 8px 25px rgba(99,102,241,0.15); }
  .model-card { cursor:pointer; transition:all 0.25s; }
  .model-card:hover { border-color:rgba(99,102,241,0.5); transform:translateY(-1px); }
  .model-card.sel { border-color:var(--accent); background:rgba(99,102,241,0.12); box-shadow:0 0 15px rgba(99,102,241,0.1); }
  .badge-local { background:rgba(34,197,94,0.15); color:#86efac; font-size:0.65rem; }
  .badge-cloud { background:rgba(251,191,36,0.15); color:#fde68a; font-size:0.65rem; }
  .chat-msg { max-width:85%; border-radius:12px; }
  .chat-user { background:linear-gradient(135deg, rgba(99,102,241,0.2), rgba(139,92,246,0.15)); margin-left:auto; border-bottom-right-radius:4px; }
  .chat-ai { background:rgba(15,23,42,0.9); margin-right:auto; border-bottom-left-radius:4px; border-left:2px solid rgba(99,102,241,0.3); }
  .src-link { cursor:pointer; transition:all 0.2s; border-radius:6px; }
  .src-link:hover { background:rgba(99,102,241,0.25); transform:scale(1.02); }
  .answer-content p { margin-bottom:0.5em; }
  .answer-content ul,.answer-content ol { margin-left:1.5em; margin-bottom:0.5em; }
  pre { white-space:pre-wrap; word-break:break-all; }
  .logo-glow { text-shadow: 0 0 40px rgba(129,140,248,0.4), 0 0 80px rgba(168,85,247,0.2); }
  /* File preview tooltip */
  .preview-tooltip { position:fixed; z-index:100; max-width:500px; max-height:400px; overflow-y:auto;
    background:var(--card); backdrop-filter:blur(16px); border:1px solid rgba(99,102,241,0.3);
    border-radius:12px; padding:16px; box-shadow:0 20px 60px rgba(0,0,0,0.5);
    font-size:12px; line-height:1.6; pointer-events:none; opacity:0; transition:opacity 0.2s; }
  .preview-tooltip.visible { opacity:1; }
  .preview-tooltip .preview-title { font-weight:600; font-size:13px; margin-bottom:8px; color:#818cf8; }
  .preview-tooltip .preview-text { color:#cbd5e1; white-space:pre-wrap; word-break:break-word; }
  body.light .preview-tooltip { background:#fff; border-color:#e2e8f0; box-shadow:0 10px 40px rgba(0,0,0,0.15); }
  body.light .preview-tooltip .preview-title { color:#4338ca; }
  body.light .preview-tooltip .preview-text { color:#334155; }
  /* Deep Thinking toggle (Apple style) */
  .deep-toggle { display:flex; align-items:center; cursor:pointer; gap:2px; }
  .deep-toggle input { display:none; }
  .deep-slider { width:36px; height:20px; background:#475569; border-radius:10px; position:relative; transition:all 0.3s; flex-shrink:0; }
  .deep-slider::after { content:''; width:16px; height:16px; background:#fff; border-radius:50%; position:absolute; top:2px; left:2px; transition:all 0.3s; }
  .deep-toggle input:checked + .deep-slider { background:linear-gradient(135deg, #6366f1, #a855f7); }
  .deep-toggle input:checked + .deep-slider::after { left:18px; }
  body.light .deep-slider { background:#cbd5e1; }
  /* Search mode buttons */
  .smode-btn { display:flex; align-items:center; gap:3px; padding:3px 8px; border-radius:6px; font-size:11px;
    background:transparent; border:1px solid rgba(148,163,184,0.2); cursor:pointer; transition:all 0.2s; position:relative; }
  .smode-btn:hover { border-color:rgba(99,102,241,0.5); background:rgba(99,102,241,0.08); }
  .smode-btn.smode-active { border-color:#6366f1; background:rgba(99,102,241,0.15); color:#818cf8; }
  .smode-btn::after { content:attr(data-tip); position:absolute; bottom:calc(100% + 6px); left:50%; transform:translateX(-50%);
    background:#1e293b; color:#e2e8f0; padding:6px 10px; border-radius:6px; font-size:11px; white-space:nowrap;
    opacity:0; pointer-events:none; transition:opacity 0.15s; z-index:50; box-shadow:0 4px 12px rgba(0,0,0,0.3); }
  .smode-btn:hover::after { opacity:1; }
  body.light .smode-btn::after { background:#fff; color:#334155; box-shadow:0 4px 12px rgba(0,0,0,0.1); border:1px solid #e2e8f0; }
  body.light .smode-btn.smode-active { background:rgba(99,102,241,0.08); color:#4338ca; border-color:#6366f1; }
  .pulse-dot { width:8px; height:8px; border-radius:50%; background:#34d399; display:inline-block; animation:pulse 2s infinite; }
  @keyframes pulse { 0%,100%{opacity:1;transform:scale(1)} 50%{opacity:0.5;transform:scale(0.8)} }
  ::-webkit-scrollbar { width:6px; }
  ::-webkit-scrollbar-track { background:transparent; }
  ::-webkit-scrollbar-thumb { background:rgba(99,102,241,0.3); border-radius:3px; }
  ::-webkit-scrollbar-thumb:hover { background:rgba(99,102,241,0.5); }
  /* Light theme - NUCLEAR override: force all elements */
  body.light { background:#f8fafc !important; color:#1e293b !important; }
  body.light::before { background:radial-gradient(ellipse at 20% 50%,rgba(99,102,241,0.03) 0%,transparent 50%) !important; }
  body.light *:not([id^="stat-"]) { color: inherit; }
  body.light .glass { background:#ffffff !important; border-color:#e2e8f0 !important; box-shadow:0 1px 8px rgba(0,0,0,0.06) !important; color:#334155 !important; }
  body.light .sinput { background:#f1f5f9 !important; border-color:#cbd5e1 !important; color:#1e293b !important; }
  body.light .sinput:focus { border-color:#6366f1 !important; }
  body.light .stat-card { background:linear-gradient(135deg,rgba(99,102,241,0.06),rgba(168,85,247,0.04)) !important; }
  /* ALL text — force dark on light */
  body.light div, body.light span, body.light p, body.light label,
  body.light h1, body.light h2, body.light h3, body.light h4,
  body.light li, body.light td, body.light th, body.light a,
  body.light input, body.light textarea, body.light select,
  body.light button, body.light code, body.light pre { color:#334155; }
  /* Specific color overrides for semantic colors */
  body.light [class*="text-indigo"] { color:#4338ca !important; }
  body.light [class*="text-purple"] { color:#6d28d9 !important; }
  body.light [class*="text-pink"] { color:#be185d !important; }
  body.light [class*="text-green"] { color:#047857 !important; }
  body.light [class*="text-red"] { color:#b91c1c !important; }
  body.light [class*="text-yellow"] { color:#b45309 !important; }
  body.light [class*="text-orange"] { color:#c2410c !important; }
  body.light [class*="text-blue"] { color:#1d4ed8 !important; }
  body.light [class*="text-gray-4"], body.light [class*="text-gray-5"] { color:#64748b !important; }
  body.light [class*="text-gray-6"] { color:#94a3b8 !important; }
  body.light [class*="text-white"] { color:#1e293b !important; }
  /* Stat cards */
  body.light #stat-files { color:#3730a3 !important; }
  body.light #stat-chunks { color:#5b21b6 !important; }
  body.light #stat-tables { color:#9d174d !important; }
  body.light #stat-errors { color:#b91c1c !important; }
  /* Chat bubbles */
  body.light .chat-user { background:#eef2ff !important; color:#1e293b !important; }
  body.light .chat-ai { background:#ffffff !important; border-left:2px solid #6366f1 !important; color:#334155 !important; }
  body.light .chat-ai * { color:#334155 !important; }
  body.light .chat-ai .text-indigo-300, body.light .chat-ai [class*="text-indigo"] { color:#4338ca !important; }
  body.light .chat-ai .text-gray-600, body.light .chat-ai [class*="text-gray-6"] { color:#94a3b8 !important; }
  body.light .answer-content, body.light .answer-content * { color:#334155 !important; }
  /* Cards/links */
  body.light .result-card:hover { background:#f1f5f9 !important; }
  body.light .model-card.sel { background:rgba(99,102,241,0.08) !important; border-color:#6366f1 !important; }
  body.light .src-link { color:#4338ca !important; }
  body.light .src-link:hover { background:rgba(99,102,241,0.1) !important; }
  /* Buttons - keep colored ones, fix neutral ones */
  body.light .bg-gray-700, body.light .bg-gray-800 { background:#e2e8f0 !important; color:#334155 !important; }
  body.light .bg-gray-700:hover, body.light .bg-gray-800:hover { background:#cbd5e1 !important; }
  body.light .bg-red-900 { background:#fee2e2 !important; color:#b91c1c !important; }
  body.light .bg-indigo-600 { color:#fff !important; }
  body.light .bg-green-600, body.light .bg-green-500 { color:#fff !important; }
  /* Badges */
  body.light .badge-local { background:rgba(34,197,94,0.1) !important; color:#047857 !important; }
  body.light .badge-cloud { background:rgba(251,191,36,0.1) !important; color:#b45309 !important; }
  /* Tables */
  body.light th { color:#4338ca !important; border-color:#e2e8f0 !important; }
  body.light td { color:#334155 !important; border-color:#f1f5f9 !important; }
  /* Borders */
  body.light [class*="border-gray"] { border-color:#e2e8f0 !important; }
  /* Logo - keep gradient visible */
  body.light .logo-glow { text-shadow:none !important; }
  body.light h1.logo-glow { color:transparent !important; }
  /* Scrollbar */
  body.light ::-webkit-scrollbar-thumb { background:rgba(99,102,241,0.2) !important; }
  /* Highlight */
  body.light .hl { background:rgba(250,204,21,0.3) !important; }
  /* Select/option */
  body.light select, body.light option { color:#1e293b !important; background:#fff !important; -webkit-text-fill-color:#1e293b !important; }
  select.sinput { min-width:100px; }
  /* Code blocks */
  body.light pre { background:#f1f5f9 !important; color:#334155 !important; }
  body.light code { background:#e2e8f0 !important; color:#334155 !important; }
  /* Pulse dot keep green */
  body.light .pulse-dot { background:#059669 !important; }
  /* Artifact panel light mode */
  body.light #artifact-panel { background:#fff !important; border-color:#e2e8f0 !important; box-shadow:-10px 0 40px rgba(0,0,0,0.08) !important; }
  body.light #artifact-panel * { color:#334155; }
  body.light #artifact-panel h1,body.light #artifact-panel h2,body.light #artifact-panel h3 { color:#1e293b; }
  body.light #artifact-panel code { background:#f1f5f9; color:#334155; }
  body.light #artifact-panel pre { background:#f8fafc; color:#334155; }
  body.light #artifact-panel table th { color:#4338ca; border-color:#e2e8f0; }
  body.light #artifact-panel table td { border-color:#f1f5f9; }
  /* Hover states in light mode */
  body.light .hover\:bg-gray-700:hover, body.light [class*="hover:bg-gray-7"]:hover { background:#e2e8f0 !important; color:#1e293b !important; }
  body.light .hover\:text-white:hover { color:#1e293b !important; }
  body.light .hover\:bg-gray-600:hover { background:#cbd5e1 !important; }
  /* Tab hover */
  body.light .hover\:text-gray-200:hover { color:#1e293b !important; }
</style>
</head>
<body class="min-h-screen">
<div class="max-w-6xl mx-auto px-4 py-6">

  <!-- Header -->
  <div class="flex justify-between items-center mb-6">
    <div>
      <div class="flex items-center gap-3">
        <img src="/static/logos/kbase-logo.svg" alt="KBase" style="width:36px;height:36px;">
        <h1 class="text-3xl font-bold bg-clip-text text-transparent bg-gradient-to-r from-indigo-400 via-purple-400 to-pink-400 logo-glow">KBase</h1>
      </div>
      <p class="text-gray-500 text-xs mt-1"><span class="pulse-dot mr-1"></span> RAG + Text2SQL + FTS | Enhanced Pipeline</p>
    </div>
    <div class="flex items-center gap-3">
      <select id="ui-lang" onchange="switchLang(this.value)" class="sinput px-2 py-1 rounded text-xs text-gray-300 outline-none">
        <option value="zh">中文</option>
        <option value="en">English</option>
      </select>
      <button onclick="toggleTheme()" id="theme-btn" class="px-3 py-2 bg-gray-800 hover:bg-gray-700 rounded-lg text-sm transition" title="Toggle theme">
        &#9790;
      </button>
      <button onclick="shutdown()" class="px-4 py-2 bg-red-900 hover:bg-red-700 rounded-lg text-sm text-red-300 transition" title="Stop server">
        Exit
      </button>
    </div>
  </div>

  <!-- Stats -->
  <div class="grid grid-cols-4 gap-3 mb-5">
    <div class="stat-card glass rounded-lg p-3 text-center">
      <div class="text-xl font-bold text-indigo-300" id="stat-files">-</div><div class="text-xs text-gray-500">Files</div>
    </div>
    <div class="stat-card glass rounded-lg p-3 text-center">
      <div class="text-xl font-bold text-purple-300" id="stat-chunks">-</div><div class="text-xs text-gray-500">Chunks</div>
    </div>
    <div class="stat-card glass rounded-lg p-3 text-center">
      <div class="text-xl font-bold text-pink-300" id="stat-tables">-</div><div class="text-xs text-gray-500">Tables</div>
    </div>
    <div class="stat-card glass rounded-lg p-3 text-center cursor-pointer" onclick="showErrors()" title="Click to view errors">
      <div class="text-xl font-bold text-red-400" id="stat-errors">-</div><div class="text-xs text-gray-500">Errors</div>
    </div>
  </div>

  <!-- Error Modal -->
  <div id="error-modal" class="hidden fixed inset-0 z-50 flex items-center justify-center" style="background:rgba(0,0,0,0.7)">
    <div class="glass rounded-xl p-6 max-w-3xl w-full mx-4 max-h-[80vh] overflow-y-auto">
      <div class="flex justify-between items-center mb-4">
        <h3 class="text-lg font-medium">Error Details</h3>
        <button onclick="document.getElementById('error-modal').classList.add('hidden')" class="text-gray-400 hover:text-white text-xl">&times;</button>
      </div>
      <div id="error-summary" class="mb-4"></div>
      <div id="error-list"></div>
    </div>
  </div>

  <!-- Tabs -->
  <div class="flex gap-5 mb-5 border-b border-gray-700 pb-1">
    <button onclick="switchTab('chat')" id="tab-chat" class="pb-2 text-sm font-medium tab-active">Chat</button>
    <button onclick="switchTab('search')" id="tab-search" class="pb-2 text-sm font-medium text-gray-400">Search</button>
    <button onclick="switchTab('sql')" id="tab-sql" class="pb-2 text-sm font-medium text-gray-400">SQL</button>
    <button onclick="switchTab('files')" id="tab-files" class="pb-2 text-sm font-medium text-gray-400">Files</button>
    <button onclick="switchTab('ingest')" id="tab-ingest" class="pb-2 text-sm font-medium text-gray-400">Ingest</button>
    <button onclick="switchTab('connectors')" id="tab-connectors" class="pb-2 text-sm font-medium text-gray-400">Connectors</button>
    <button onclick="switchTab('settings')" id="tab-settings" class="pb-2 text-sm font-medium text-gray-400">Settings</button>
  </div>

  <!-- ============ CHAT ============ -->
  <div id="panel-chat">
    <!-- Buddy greeting -->
    <div id="chat-welcome" class="text-center py-8">
      <div class="text-5xl mb-3" id="buddy-avatar"></div>
      <div class="text-lg text-gray-300" id="buddy-greeting">Hey! I'm your knowledge base buddy. Ask me anything about your files.</div>
      <div class="text-xs text-gray-500 mt-2" id="buddy-status"></div>
    </div>
    <!-- History panel -->
    <div id="chat-history-panel" class="hidden glass rounded-xl p-4 mb-3 max-h-48 overflow-y-auto">
      <div class="flex justify-between items-center mb-2">
        <span class="text-sm font-medium">Conversation History</span>
        <button onclick="document.getElementById('chat-history-panel').classList.add('hidden')" class="text-xs text-gray-500 hover:text-red-400">&times;</button>
      </div>
      <div id="chat-history-list" class="space-y-1"></div>
    </div>
    <div id="chat-messages" class="space-y-3 mb-4" style="max-height:55vh; overflow-y:auto;"></div>
    <div class="glass rounded-xl p-4">
      <div class="flex gap-2 mb-2">
        <select id="chat-provider" class="sinput px-2 py-1 rounded text-gray-300 outline-none text-xs"></select>
        <select id="chat-buddy" class="sinput px-2 py-1 rounded text-gray-300 outline-none text-xs"></select>
        <div class="flex gap-1" id="search-mode-btns">
          <button onclick="setSearchMode('kb')" class="smode-btn smode-active" data-mode="kb" data-tip="Search local files only / 仅搜索本地文件">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M4 19.5A2.5 2.5 0 016.5 17H20"/><path d="M6.5 2H20v20H6.5A2.5 2.5 0 014 19.5v-15A2.5 2.5 0 016.5 2z"/></svg>
            KB
          </button>
          <button onclick="setSearchMode('web')" class="smode-btn" data-mode="web" data-tip="Search the internet / 搜索互联网">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><circle cx="12" cy="12" r="10"/><path d="M2 12h20"/><path d="M12 2a15 15 0 014 10 15 15 0 01-4 10 15 15 0 01-4-10 15 15 0 014-10z"/></svg>
            Web
          </button>
          <button onclick="setSearchMode('hybrid')" class="smode-btn" data-mode="hybrid" data-tip="Search both KB and web / 知识库+网络混合搜索">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M16 3h5v5"/><path d="M8 3H3v5"/><path d="M21 3l-7 7"/><path d="M3 3l7 7"/><path d="M3 21l7-7"/><path d="M21 21l-7-7"/></svg>
            Hybrid
          </button>
          <button onclick="setSearchMode('research')" class="smode-btn" data-mode="research" data-tip="Deep multi-step research / 拆解问题+多轮搜索+综合分析">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><circle cx="11" cy="11" r="8"/><path d="M21 21l-4.35-4.35"/><path d="M11 8v6"/><path d="M8 11h6"/></svg>
            Research
          </button>
        </div>
        <input type="hidden" id="chat-search-mode" value="kb">
        <!-- Deep Thinking toggle (Apple style) -->
        <label class="deep-toggle" title="Deep Thinking: multi-round iterative research / 深度思考：多轮迭代搜索研究">
          <input type="checkbox" id="deep-thinking-toggle">
          <span class="deep-slider"></span>
          <span class="text-xs ml-1" style="user-select:none;">Deep</span>
        </label>
        <span class="text-xs text-gray-500 flex items-center" id="chat-memory-badge">Memory: 0 turns</span>
        <button onclick="showHistory()" class="ml-auto text-xs text-gray-500 hover:text-indigo-300 transition">History</button>
        <button onclick="clearChat()" class="text-xs text-gray-500 hover:text-red-400 transition">Clear</button>
      </div>
      <div class="flex gap-3">
        <input id="chat-input" type="text" placeholder="Ask your knowledge base..."
          class="sinput flex-1 px-4 py-3 rounded-lg text-white outline-none"
          onkeydown="if(event.key==='Enter'&&!event.shiftKey&&!event.isComposing){event.preventDefault();doChat()}">
        <button onclick="doChat()" id="chat-btn" class="px-6 py-3 bg-indigo-600 hover:bg-indigo-500 rounded-lg font-medium transition">Send</button>
      </div>
    </div>
  </div>

  <!-- ============ SEARCH ============ -->
  <div id="panel-search" class="hidden">
    <div class="glass rounded-xl p-5 mb-4">
      <div class="flex gap-3">
        <input id="search-input" type="text" placeholder="Search..." class="sinput flex-1 px-4 py-3 rounded-lg text-white outline-none text-lg" onkeydown="if(event.key==='Enter'&&!event.isComposing)doSearch()">
        <select id="search-type" class="sinput px-3 py-3 rounded-lg text-gray-300 outline-none">
          <option value="auto">Auto</option><option value="semantic">Semantic</option><option value="keyword">Keyword</option>
        </select>
        <button onclick="doSearch()" class="px-6 py-3 bg-indigo-600 hover:bg-indigo-500 rounded-lg font-medium transition">Search</button>
      </div>
    </div>
    <div id="results" class="space-y-3"></div>
  </div>

  <!-- ============ SQL ============ -->
  <div id="panel-sql" class="hidden">
    <div class="glass rounded-xl p-5 mb-4">
      <textarea id="sql-input" rows="3" placeholder="SELECT * FROM table_name LIMIT 10" class="sinput w-full px-4 py-3 rounded-lg text-white outline-none font-mono text-sm"></textarea>
      <div class="flex gap-3 mt-3">
        <button onclick="doSQL()" class="px-4 py-2 bg-indigo-600 hover:bg-indigo-500 rounded-lg text-sm transition">Execute</button>
        <button onclick="loadTables()" class="px-4 py-2 bg-gray-700 hover:bg-gray-600 rounded-lg text-sm transition">Tables</button>
      </div>
    </div>
    <div id="sql-results"></div><div id="table-list"></div>
  </div>

  <!-- ============ FILES ============ -->
  <div id="panel-files" class="hidden"><div id="file-list"></div></div>

  <!-- ============ INGEST ============ -->
  <div id="panel-ingest" class="hidden">
    <div class="glass rounded-xl p-5">
      <h3 class="text-lg font-medium mb-3">Ingest Directory</h3>
      <div class="flex gap-3 mb-2">
        <input id="ingest-path" type="text" placeholder="/path/to/files" class="sinput flex-1 px-4 py-3 rounded-lg text-white outline-none">
        <button onclick="browseDir()" class="px-4 py-3 bg-gray-700 hover:bg-gray-600 rounded-lg text-sm transition" title="Browse directories">Browse</button>
        <label class="flex items-center gap-2 text-sm text-gray-400"><input type="checkbox" id="ingest-force"> Force</label>
        <button onclick="doIngest()" class="px-5 py-3 bg-green-600 hover:bg-green-500 rounded-lg font-medium transition">Ingest</button>
      </div>
      <!-- Directory browser -->
      <div id="dir-browser" class="hidden glass rounded-lg p-3 mb-3 max-h-64 overflow-y-auto">
        <div class="flex justify-between items-center mb-2">
          <span class="text-xs text-gray-400" id="dir-current">/</span>
          <button onclick="document.getElementById('dir-browser').classList.add('hidden')" class="text-xs text-gray-500 hover:text-gray-300">Close</button>
        </div>
        <div id="dir-list" class="space-y-1"></div>
      </div>
      <div id="ingest-result"></div>

      <h3 class="text-lg font-medium mt-5 mb-3">Upload Files</h3>
      <div class="space-y-3">
        <div class="flex gap-3 items-center">
          <label class="px-4 py-2 bg-gray-700 hover:bg-gray-600 rounded-lg text-sm transition cursor-pointer">
            Choose Files
            <input type="file" id="upload-file" class="hidden" multiple>
          </label>
          <span class="text-xs text-gray-500" id="upload-file-label">No files selected</span>
          <button onclick="doUpload()" class="px-4 py-2 bg-green-600 hover:bg-green-500 rounded-lg text-sm transition">Upload</button>
        </div>
        <div class="flex gap-3 items-center">
          <label class="px-4 py-2 bg-gray-700 hover:bg-gray-600 rounded-lg text-sm transition cursor-pointer">
            Choose Folder
            <input type="file" id="upload-folder" class="hidden" webkitdirectory mozdirectory directory multiple>
          </label>
          <span class="text-xs text-gray-500" id="upload-folder-label">No folder selected</span>
          <button onclick="doUploadFolder()" class="px-4 py-2 bg-green-600 hover:bg-green-500 rounded-lg text-sm transition">Upload Folder</button>
        </div>
      </div>
      <div id="upload-result" class="mt-3"></div>
    </div>
  </div>

  <!-- ============ CONNECTORS ============ -->
  <div id="panel-connectors" class="hidden">
    <div class="glass rounded-xl p-5">
      <h3 class="text-lg font-medium mb-1">Data Connectors (数据源连接)</h3>
      <p class="text-xs text-gray-500 mb-4">连接外部数据源，自动同步内容到知识库。</p>
      <div id="connector-list" class="grid grid-cols-1 md:grid-cols-2 gap-4"></div>
    </div>
    <!-- Connector config panel -->
    <div id="connector-config" class="hidden glass rounded-xl p-5 mt-4">
      <div class="flex justify-between items-center mb-3">
        <h3 class="text-lg font-medium" id="connector-config-title"></h3>
        <button onclick="document.getElementById('connector-config').classList.add('hidden')" class="text-gray-400 hover:text-white">&times;</button>
      </div>
      <div id="connector-config-fields"></div>
      <div id="connector-config-actions" class="mt-4 flex gap-3"></div>
      <div id="connector-sync-result" class="mt-3"></div>
    </div>
  </div>

  <!-- ============ SETTINGS ============ -->
  <div id="panel-settings" class="hidden">
    <div class="space-y-5">
      <!-- LLM Provider (API Key inline) -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">LLM Provider (大模型)</h3>
        <p class="text-xs text-gray-500 mb-3">选择模型后，下方会出现 API Key 输入框。本地模型无需 Key。</p>
        <div id="llm-models" class="grid grid-cols-1 md:grid-cols-2 lg:grid-cols-3 gap-3"></div>
        <!-- Inline API config (shown when model selected) -->
        <div id="llm-config" class="mt-4 hidden glass rounded-lg p-4" style="border-left:3px solid var(--accent)">
          <div id="llm-config-title" class="text-sm font-medium mb-2"></div>
          <div id="llm-config-fields"></div>
        </div>
      </div>
      <!-- Buddy & Memory -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">Chat Personality & Memory</h3>
        <div class="grid grid-cols-2 gap-4 mt-3">
          <div>
            <label class="text-sm text-gray-300 block mb-2">Buddy Preset (助手人格)</label>
            <div id="buddy-presets" class="space-y-2"></div>
          </div>
          <div>
            <label class="text-sm text-gray-300 block mb-2">Memory Turns (记忆轮数)</label>
            <input id="memory-turns" type="number" min="0" max="50" step="1" class="sinput w-full px-3 py-2 rounded-lg text-white outline-none">
            <div class="text-xs text-gray-500 mt-1">保留多少轮对话上下文 (0=无记忆, 推荐5~15)</div>
          </div>
        </div>
      </div>
      <!-- Embedding -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">Embedding Model (向量化)</h3>
        <p class="text-xs text-gray-500 mb-3">切换后需重新 Ingest 生效。</p>
        <div id="embedding-models" class="grid grid-cols-1 md:grid-cols-2 gap-3"></div>
        <div id="embedding-config" class="mt-3 hidden glass rounded-lg p-4" style="border-left:3px solid var(--accent)">
          <div id="embedding-config-title" class="text-sm font-medium mb-2"></div>
          <div id="embedding-config-fields"></div>
        </div>
      </div>
      <!-- Whisper -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">Whisper Model (语音识别)</h3>
        <p class="text-xs text-gray-500 mb-3">音频文件转文字。</p>
        <div id="whisper-models" class="grid grid-cols-1 md:grid-cols-2 gap-3"></div>
        <div id="whisper-config" class="mt-3 hidden glass rounded-lg p-4" style="border-left:3px solid var(--accent)">
          <div id="whisper-config-title" class="text-sm font-medium mb-2"></div>
          <div id="whisper-config-fields"></div>
        </div>
      </div>
      <!-- Chunk -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">Chunk Settings (切片)</h3>
        <div class="grid grid-cols-2 gap-4 mt-3">
          <div>
            <label class="text-sm text-gray-300 block mb-1">Max Size (字符)</label>
            <input id="chunk-max" type="number" min="500" max="5000" step="100" class="sinput w-full px-3 py-2 rounded-lg text-white outline-none">
            <div class="text-xs text-gray-500 mt-1">推荐 1000~2000</div>
          </div>
          <div>
            <label class="text-sm text-gray-300 block mb-1">Overlap (字符)</label>
            <input id="chunk-overlap" type="number" min="0" max="1000" step="50" class="sinput w-full px-3 py-2 rounded-lg text-white outline-none">
            <div class="text-xs text-gray-500 mt-1">推荐 max 的 10~20%</div>
          </div>
        </div>
      </div>
      <!-- Language Profile -->
      <div class="glass rounded-xl p-5">
        <h3 class="text-lg font-medium mb-1">Language (语言优化)</h3>
        <p class="text-xs text-gray-500 mb-3">选择文档主要语言，影响分词和检索策略。切换后需重新 Ingest。</p>
        <div id="language-profiles" class="grid grid-cols-1 md:grid-cols-3 gap-3"></div>
      </div>
      <button onclick="saveSettings()" class="w-full py-3 bg-indigo-600 hover:bg-indigo-500 rounded-lg font-medium transition">Save Settings</button>
      <div id="settings-msg" class="text-center"></div>
    </div>
  </div>
</div>

<!-- Artifact/Report panel (slide from right) -->
<div id="artifact-panel" class="hidden" style="position:fixed;top:0;right:0;width:50%;height:100vh;z-index:60;background:var(--card);border-left:1px solid rgba(148,163,184,0.15);overflow-y:auto;transition:transform 0.3s;box-shadow:-10px 0 40px rgba(0,0,0,0.3);">
  <div style="padding:20px;">
    <div class="flex justify-between items-center mb-4">
      <div class="flex items-center gap-3">
        <svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" style="color:#818cf8"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><path d="M14 2v6h6"/></svg>
        <span class="text-lg font-medium" id="artifact-title">Research Report</span>
      </div>
      <div class="flex items-center gap-2">
        <button onclick="downloadArtifact()" style="padding:6px 12px;border-radius:6px;font-size:12px;background:#4f46e5;color:#fff;" title="Download as Markdown">
          <svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" class="inline mr-1"><path d="M21 15v4a2 2 0 01-2 2H5a2 2 0 01-2-2v-4"/><path d="M7 10l5 5 5-5"/><path d="M12 15V3"/></svg>Download .md
        </button>
        <button onclick="closeArtifact()" style="font-size:20px;color:#94a3b8;padding:4px 8px;">&times;</button>
      </div>
    </div>
    <div id="artifact-content" style="line-height:1.8;font-size:14px;"></div>
  </div>
</div>

<!-- File preview tooltip -->
<div id="preview-tooltip" class="preview-tooltip"></div>

<script>
const API='';
let settings={};

async function fetchJSON(u,o){const r=await fetch(API+u,o);return r.json();}
function esc(s){return s.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');}
function escRe(s){return s.replace(/[.*+?^${}()|[\]\\]/g,'\\$&');}

// ---- Stats ----
async function loadStats(){
  const s=await fetchJSON('/api/status');
  document.getElementById('stat-files').textContent=s.file_count||0;
  document.getElementById('stat-chunks').textContent=s.chunk_count||0;
  document.getElementById('stat-tables').textContent=s.table_count||0;
  document.getElementById('stat-errors').textContent=s.error_count||0;
}

// ---- Chat ----
// Persist conversation ID in localStorage so refresh keeps the same conversation
let convId=localStorage.getItem('kbase-conv-id');
if(!convId){convId='conv-'+Date.now();localStorage.setItem('kbase-conv-id',convId);}
let chatTurns=0;
let chatAbort=null;

async function doChat(){
  // If already running, abort
  if(chatAbort){chatAbort.abort();chatAbort=null;document.getElementById('chat-btn').textContent=curLang==='zh'?'发送':'Send';return;}
  const isDeep=document.getElementById('deep-thinking-toggle').checked;
  const q=document.getElementById('chat-input').value.trim();
  if(!q)return;
  const el=document.getElementById('chat-messages');
  const btn=document.getElementById('chat-btn');
  const provider=document.getElementById('chat-provider').value;
  const buddy=document.getElementById('chat-buddy').value;

  // Hide welcome on first message
  const welcome=document.getElementById('chat-welcome');
  if(welcome)welcome.style.display='none';

  // User message
  el.innerHTML+=`<div class="chat-msg chat-user rounded-lg px-4 py-3 text-sm fade-in">${esc(q)}</div>`;
  document.getElementById('chat-input').value='';
  el.scrollTop=el.scrollHeight;

  // Loading
  const loadId='load-'+Date.now();
  const deepLoadMsg=isDeep?(curLang==='zh'?'Deep Thinking 启动...':'Deep Thinking starting...'):(curLang==='zh'?'正在检索并生成回答...':'Searching...');
  el.innerHTML+=`<div id="${loadId}" class="chat-msg chat-ai rounded-lg px-4 py-3 text-sm fade-in">
    <div class="text-gray-400">${deepLoadMsg}</div>
    ${isDeep?'<div id="deep-progress" class="mt-2 text-xs space-y-1" style="color:#94a3b8;"></div>':''}
  </div>`;
  el.scrollTop=el.scrollHeight;

  // === Deep Thinking: SSE stream ===
  if(isDeep){
    btn.textContent='Stop';
    const url='/api/research-stream?question='+encodeURIComponent(q)+'&conv_id='+convId;
    const evtSrc=new EventSource(url);
    const progEl=document.getElementById('deep-progress');
    evtSrc.onmessage=function(e){
      const d=JSON.parse(e.data);
      if(d.type==='round'&&progEl)progEl.innerHTML+=`<div>Round ${d.num}: ${(d.queries||[]).join(', ')} <span style="color:#64748b">(${d.total_urls} URLs)</span></div>`;
      else if(d.type==='round_done'&&progEl)progEl.innerHTML+=`<div style="color:#059669">+${d.new_findings} findings (${d.total_urls} total URLs)</div>`;
      else if(d.type==='phase'&&progEl)progEl.innerHTML+=`<div style="color:#818cf8;font-weight:600">${d.name_zh||d.name}...</div>`;
      else if((d.type==='sufficient'||d.type==='timeout')&&progEl)progEl.innerHTML+=`<div style="color:#d97706">${d.reason||'Done'}</div>`;
      else if(d.type==='result'){
        evtSrc.close();
        const answer=formatAnswer(d.answer||'',(d.sources||[]));
        // Save report for artifact preview
        lastReport={text:d.answer||'',sources:d.sources||[],web:d.web_sources||[],stats:d.stats||{}};
        let srcHtml=buildSourcesHtml(d.sources||[],d.web_sources||[]);
        const st=d.stats||{};
        document.getElementById(loadId).innerHTML=`
          <div class="answer-content leading-relaxed">${answer}</div>
          ${srcHtml}
          <div class="flex justify-between items-center mt-2 pt-1 border-t border-gray-800">
            <div class="text-xs" style="color:#64748b">Deep: ${st.rounds||0} rounds | ${st.total_urls||0} URLs | ${st.elapsed||0}s</div>
            <div class="flex gap-2">
              <button onclick="showArtifact()" class="text-xs flex items-center gap-1" style="color:#818cf8" title="View as document"><svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><path d="M14 2v6h6"/></svg>Report</button>
              <button onclick="rewindChat()" class="text-xs flex items-center gap-1" style="color:#64748b"><svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M1 4v6h6"/><path d="M3.51 15a9 9 0 102.13-9.36L1 10"/></svg>Rewind</button>
            </div>
          </div>`;
        chatTurns++;
        document.getElementById('chat-memory-badge').textContent='Memory: '+chatTurns+' turns';
      }else if(d.type==='error'){evtSrc.close();document.getElementById(loadId).innerHTML=`<div style="color:#dc2626">Error: ${esc(d.message||'')}</div>`;}
      el.scrollTop=el.scrollHeight;
    };
    evtSrc.onerror=function(){evtSrc.close();};
    btn.textContent=curLang==='zh'?'发送':'Send';
    return;
  }

  // === Normal mode ===
  chatAbort=new AbortController();
  btn.textContent='Stop';

  try{
    const resp=await fetch('/api/chat',{
      method:'POST',
      headers:{'Content-Type':'application/json'},
      signal:chatAbort.signal,
      body:JSON.stringify({
        question:q,
        conversation_id:convId,
        settings_override:{llm_provider:provider, buddy_preset:buddy, search_mode:document.getElementById('chat-search-mode').value},
        top_k:10
      })
    });
    const data=await resp.json();

    chatTurns=data.history_turns||0;
    document.getElementById('chat-memory-badge').textContent=`Memory: ${chatTurns} turns`;

    // Sources
    let srcHtml='';
    if(data.sources&&data.sources.length>0){
      srcHtml='<div class="mt-3 pt-2 border-t border-gray-700"><div class="text-xs text-gray-500 mb-1">Sources:</div><div class="flex flex-wrap gap-1">';
      data.sources.forEach(s=>{
        const fname=s.name||'?', fpath=s.path||'';
        srcHtml+=`<span class="src-link text-xs px-2 py-1 rounded bg-gray-800 text-indigo-300 flex items-center gap-1" onclick="openFile('${fpath.replace(/'/g,"\\'")}')">
          <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M3 7v10a2 2 0 002 2h14a2 2 0 002-2V9a2 2 0 00-2-2h-6l-2-2H5a2 2 0 00-2 2z"/></svg>
          ${esc(fname)}</span>`;
      });
      srcHtml+='</div>';
      // Web sources
      if(data.web_sources&&data.web_sources.length){
        srcHtml+='<div class="flex flex-wrap gap-1 mt-1">';
        data.web_sources.forEach(w=>{
          srcHtml+=`<a href="${w.url}" target="_blank" rel="noopener" class="src-link text-xs px-2 py-1 rounded bg-gray-800 flex items-center gap-1" style="color:#60a5fa;">
            <svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><circle cx="12" cy="12" r="10"/><path d="M2 12h20"/><path d="M12 2a15.3 15.3 0 014 10 15.3 15.3 0 01-4 10 15.3 15.3 0 01-4-10 15.3 15.3 0 014-10z"/></svg>
            ${esc(w.name||'').substring(0,40)}</a>`;
        });
        srcHtml+='</div>';
      }
      srcHtml+='</div>';
    }

    const answer=formatAnswer(data.answer||'No response', data.sources||[]);
    const turnNum=data.history_turns||chatTurns;
    document.getElementById(loadId).outerHTML=`<div class="chat-msg chat-ai rounded-lg px-4 py-3 text-sm fade-in">
      <div class="answer-content leading-relaxed">${answer}</div>
      ${srcHtml}
      <div class="flex justify-between items-center mt-2 pt-1 border-t border-gray-800">
        <div class="text-xs text-gray-600">${data.provider} | ${data.context_chunks} chunks</div>
        <button onclick="rewindChat()" class="text-xs text-gray-500 hover:text-yellow-400 transition flex items-center gap-1" title="Rewind: undo this exchange">
          <svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M1 4v6h6"/><path d="M3.51 15a9 9 0 1 0 2.13-9.36L1 10"/></svg>
          Rewind
        </button>
      </div>
    </div>`;
  }catch(e){
    if(e.name==='AbortError'){
      document.getElementById(loadId).outerHTML=`<div class="chat-msg chat-ai rounded-lg px-4 py-3 text-sm text-gray-500 fade-in">[Stopped]</div>`;
    }else{
      document.getElementById(loadId).outerHTML=`<div class="chat-msg chat-ai rounded-lg px-4 py-3 text-sm text-red-400 fade-in">Error: ${esc(e.message||String(e))}</div>`;
    }
  }
  chatAbort=null;
  btn.textContent=curLang==='zh'?'发送':'Send';
  el.scrollTop=el.scrollHeight;
}

// --- Artifact / Report ---
let lastReport=null;

function buildSourcesHtml(sources,webSources){
  let h='';
  if((sources&&sources.length)||(webSources&&webSources.length)){
    h='<div class="mt-3 pt-2 border-t border-gray-700"><div class="text-xs text-gray-500 mb-1">Sources:</div><div class="flex flex-wrap gap-1">';
    (sources||[]).forEach(s=>{
      const p=(s.path||'').replace(/'/g,"\\'");
      h+=`<span class="src-link text-xs px-2 py-1 rounded bg-gray-800 text-indigo-300" onclick="openFile('${p}')"><svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" class="inline mr-1"><path d="M3 7v10a2 2 0 002 2h14a2 2 0 002-2V9a2 2 0 00-2-2h-6l-2-2H5a2 2 0 00-2 2z"/></svg>${esc(s.name||'')}</span>`;
    });
    (webSources||[]).forEach(w=>{
      h+=`<a href="${w.url}" target="_blank" class="src-link text-xs px-2 py-1 rounded bg-gray-800" style="color:#60a5fa"><svg width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" class="inline mr-1"><circle cx="12" cy="12" r="10"/><path d="M2 12h20"/></svg>${esc((w.name||'').substring(0,40))}</a>`;
    });
    h+='</div></div>';
  }
  return h;
}

function showArtifact(){
  if(!lastReport)return;
  const panel=document.getElementById('artifact-panel');
  panel.classList.remove('hidden');
  document.getElementById('artifact-title').textContent='Research Report';
  // Render as formatted HTML
  const rendered=formatAnswer(lastReport.text,[]);
  document.getElementById('artifact-content').innerHTML=rendered;
}

function closeArtifact(){
  document.getElementById('artifact-panel').classList.add('hidden');
}

function downloadArtifact(){
  if(!lastReport)return;
  const blob=new Blob([lastReport.text],{type:'text/markdown'});
  const url=URL.createObjectURL(blob);
  const a=document.createElement('a');
  a.href=url;
  a.download=`kbase-research-${new Date().toISOString().slice(0,10)}.md`;
  a.click();
  URL.revokeObjectURL(url);
}

function setSearchMode(mode){
  document.getElementById('chat-search-mode').value=mode;
  document.querySelectorAll('.smode-btn').forEach(b=>{
    b.classList.toggle('smode-active',b.dataset.mode===mode);
  });
}

async function rewindChat(){
  // Get the last user message before removing
  const el=document.getElementById('chat-messages');
  const userMsgs=el.querySelectorAll('.chat-user');
  const lastUserMsg=userMsgs.length?userMsgs[userMsgs.length-1].textContent:'';

  await fetchJSON('/api/chat/rewind',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({conversation_id:convId,turns:1})});
  const msgs=el.querySelectorAll('.chat-msg');
  if(msgs.length>=2){msgs[msgs.length-1].remove();msgs[msgs.length-2].remove();}
  else if(msgs.length>=1){msgs[msgs.length-1].remove();}
  chatTurns=Math.max(0,chatTurns-1);
  document.getElementById('chat-memory-badge').textContent=`Memory: ${chatTurns} turns`;

  // Put user's question back in input
  if(lastUserMsg){
    document.getElementById('chat-input').value=lastUserMsg.trim();
    document.getElementById('chat-input').focus();
  }
}

async function clearChat(){
  if(!confirm('Clear conversation?'))return;
  await fetchJSON('/api/chat/clear',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({conversation_id:convId})});
  // New conversation ID
  convId='conv-'+Date.now();
  localStorage.setItem('kbase-conv-id',convId);
  document.getElementById('chat-messages').innerHTML='';
  const welcome=document.getElementById('chat-welcome');
  if(welcome)welcome.style.display='block';
  chatTurns=0;
  document.getElementById('chat-memory-badge').textContent='Memory: 0 turns';
}

function formatAnswer(text, sources){
  // Full markdown → HTML renderer
  let lines=(text||'').split('\n');
  let html='';
  let inCode=false, inTable=false, inList=false, listType='';

  for(let i=0;i<lines.length;i++){
    let line=lines[i];

    // Code blocks
    if(line.trim().startsWith('```')){
      if(inCode){html+='</pre>';inCode=false;}
      else{html+='<pre class="bg-gray-900 rounded p-3 my-2 text-xs overflow-x-auto">';inCode=true;}
      continue;
    }
    if(inCode){html+=esc(line)+'\n';continue;}

    // Close table if line doesn't start with |
    if(inTable && !line.trim().startsWith('|')){html+='</tbody></table></div>';inTable=false;}
    // Close list if empty line or non-list line
    if(inList && !line.trim().match(/^[-*]\s|^\d+\.\s/) && line.trim()!==''){
      html+=listType==='ul'?'</ul>':'</ol>';inList=false;
    }

    let escaped=esc(line);

    // Table separator — skip
    if(escaped.trim().match(/^\|[\s\-:|]+\|?$/) || escaped.trim().match(/^[\s\-:|]+$/)){continue;}

    // Table row (flexible: starts with | or has multiple |)
    const pipeCount=(escaped.match(/\|/g)||[]).length;
    const looksLikeTable=escaped.trim().startsWith('|') && pipeCount>=2;
    if(looksLikeTable){
      let raw=escaped.trim();
      if(raw.startsWith('|'))raw=raw.slice(1);
      if(raw.endsWith('|'))raw=raw.slice(0,-1);
      const cells=raw.split('|').map(c=>c.trim());
      if(!inTable){
        html+='<div class="overflow-x-auto my-2 rounded-lg border border-gray-700"><table class="w-full text-sm border-collapse"><thead><tr>';
        cells.forEach(c=>html+=`<th class="px-3 py-2 text-left font-medium border-b border-gray-600" style="white-space:nowrap">${linkify(c,sources)}</th>`);
        html+='</tr></thead><tbody>';
        inTable=true;
        if(i+1<lines.length && lines[i+1].trim().match(/^\|?[\s\-:|]+\|?$/))i++;
      }else{
        html+='<tr class="hover:bg-gray-800 hover:bg-opacity-30">';
        cells.forEach(c=>html+=`<td class="px-3 py-1.5 border-b border-gray-800">${linkify(c,sources)}</td>`);
        html+='</tr>';
      }
      continue;
    }

    // Headers
    if(escaped.match(/^#{1,4}\s/)){
      const level=escaped.match(/^(#+)/)[1].length;
      const content=escaped.replace(/^#+\s*/,'');
      const sizes={1:'text-xl',2:'text-lg',3:'text-base',4:'text-sm'};
      html+=`<div class="${sizes[level]||'text-base'} font-bold mt-3 mb-1">${linkify(content,sources)}</div>`;
      continue;
    }

    // Unordered list
    if(escaped.match(/^\s*[-*]\s/)){
      if(!inList){html+='<ul class="list-disc ml-5 my-1 space-y-0.5">';inList=true;listType='ul';}
      html+=`<li>${linkify(escaped.replace(/^\s*[-*]\s/,''),sources)}</li>`;
      continue;
    }

    // Ordered list
    if(escaped.match(/^\s*\d+\.\s/)){
      if(!inList){html+='<ol class="list-decimal ml-5 my-1 space-y-0.5">';inList=true;listType='ol';}
      html+=`<li>${linkify(escaped.replace(/^\s*\d+\.\s/,''),sources)}</li>`;
      continue;
    }

    // Empty line
    if(escaped.trim()===''){html+='<div class="h-2"></div>';continue;}

    // Normal paragraph
    html+=`<div class="my-0.5">${linkify(escaped,sources)}</div>`;
  }

  // Close open blocks
  if(inCode)html+='</pre>';
  if(inTable)html+='</tbody></table></div>';
  if(inList)html+=listType==='ul'?'</ul>':'</ol>';

  return html;
}

function linkify(text, sources){
  // Bold
  text=text.replace(/\*\*(.+?)\*\*/g,'<strong>$1</strong>');
  // Inline code
  text=text.replace(/`([^`]+)`/g,'<code class="bg-gray-800 px-1 rounded text-xs">$1</code>');
  // [file references] → clickable
  text=text.replace(/\[([^\]]+)\]/g, function(match,ref){
    const src=(sources||[]).find(s=>{
      const n=s.name||'';
      return ref.includes(n)||n.includes(ref)||ref.replace(/\s*\(.*\)/,'').trim()===n;
    })||(sources||[]).find(s=>{
      const n=(s.name||'').replace(/\.[^.]+$/,'');
      return n.length>3&&ref.includes(n);
    });
    if(src&&src.path){
      const p=src.path.replace(/'/g,"\\'");
      return `<span class="src-link text-indigo-300 font-medium underline cursor-pointer" onclick="openFile('${p}')" onmouseenter="showPreview(this,'${p}')" onmouseleave="hidePreview()">[${ref}]</span>`;
    }
    return `<span class="text-indigo-300 font-medium">[${ref}]</span>`;
  });
  return text;
}

// File preview on hover
let previewCache={};
let previewTimer=null;

function showPreview(el, path){
  if(previewTimer)clearTimeout(previewTimer);
  previewTimer=setTimeout(async()=>{
    const tooltip=document.getElementById('preview-tooltip');
    if(!tooltip)return;
    // Position near element
    const rect=el.getBoundingClientRect();
    tooltip.style.left=Math.min(rect.left, window.innerWidth-520)+'px';
    tooltip.style.top=Math.min(rect.bottom+8, window.innerHeight-420)+'px';
    // Show loading
    tooltip.innerHTML='<div class="text-gray-400">Loading preview...</div>';
    tooltip.classList.add('visible');
    // Fetch preview (cached)
    if(!previewCache[path]){
      try{
        previewCache[path]=await fetchJSON('/api/preview?path='+encodeURIComponent(path)+'&max_chars=1500');
      }catch(e){
        previewCache[path]={preview:'Preview unavailable',title:path.split('/').pop()};
      }
    }
    const data=previewCache[path];
    if(data.type==='image'&&data.image){
      tooltip.innerHTML=`<div class="preview-title">${esc(data.title||'')}</div>
        <img src="data:image/png;base64,${data.image}" style="width:100%;border-radius:6px;">`;
    }else{
      tooltip.innerHTML=`<div class="preview-title">${esc(data.title||'')} <span class="text-xs text-gray-500">${data.type||''}</span></div>
        <div class="preview-text">${esc(data.preview||'').substring(0,800)}</div>`;
    }
  },400); // 400ms delay to avoid flicker
}

function hidePreview(){
  if(previewTimer)clearTimeout(previewTimer);
  previewTimer=null;
  const tooltip=document.getElementById('preview-tooltip');
  if(tooltip)tooltip.classList.remove('visible');
}

async function openFile(path){
  try{
    await fetchJSON('/api/open-file',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({path})});
  }catch(e){alert('Cannot open: '+e.message);}
}

// ---- Search ----
async function doSearch(){
  const q=document.getElementById('search-input').value.trim();
  if(!q)return;
  const type=document.getElementById('search-type').value;
  const el=document.getElementById('results');
  el.innerHTML=`<div class="text-gray-400">${curLang==='zh'?'正在搜索...':'Searching...'}</div>`;
  const data=await fetchJSON(`/api/search?q=${encodeURIComponent(q)}&type=${type}&top_k=15`);
  if(!data.results||data.results.length===0){el.innerHTML='<div class="text-gray-500 text-center py-8">No results</div>';return;}

  el.innerHTML=`<div class="text-sm text-gray-400 mb-3">${data.result_count} results (${data.methods_used.join('+')})</div>`+
    data.results.map((r,i)=>{
      const m=r.metadata||{};
      const score=(r.rrf_score||r.score||0).toFixed(4);
      const text=(r.text||'').substring(0,400);
      let highlighted=esc(text);
      q.split(/\s+/).filter(t=>t.length>1).forEach(t=>{
        highlighted=highlighted.replace(new RegExp('('+escRe(t)+')','gi'),'<span class="hl">$1</span>');
      });
      return `<div class="glass rounded-lg p-4 result-card transition fade-in" style="animation-delay:${i*40}ms">
        <div class="flex justify-between items-start mb-1">
          <span class="text-indigo-300 font-medium src-link" onclick="openFile('${(m.file_path||'').replace(/'/g,"\\'")}')">${m.file_name||'?'}
            <svg class="inline ml-1" width="12" height="12" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M3 7v10a2 2 0 002 2h14a2 2 0 002-2V9a2 2 0 00-2-2h-6l-2-2H5a2 2 0 00-2 2z"/></svg>
          </span>
          <span class="text-xs text-gray-500">${r.method}|${score}</span>
        </div>
        <div class="text-xs text-gray-600 mb-2">${m.file_path||''}</div>
        <div class="text-sm text-gray-300 leading-relaxed">${highlighted}${text.length>=400?'...':''}</div>
        ${m.slide?'<span class="text-xs text-purple-400">Slide '+m.slide+'</span>':''}
        ${m.page?'<span class="text-xs text-purple-400">Page '+m.page+'</span>':''}
      </div>`;
    }).join('');
}

// ---- SQL ----
async function doSQL(){
  const q=document.getElementById('sql-input').value.trim();
  if(!q)return;
  const el=document.getElementById('sql-results');
  el.innerHTML='<div class="text-gray-400">Running...</div>';
  const data=await fetchJSON(`/api/sql?q=${encodeURIComponent(q)}`);
  const r=data.results||data;
  if(r.error){el.innerHTML=`<div class="text-red-400 glass rounded-lg p-4">${r.error}</div>`;return;}
  if(!r.columns||!r.columns.length){el.innerHTML='<div class="text-gray-500">Empty</div>';return;}
  let h='<div class="overflow-x-auto glass rounded-lg"><table class="w-full text-sm"><thead><tr>'+r.columns.map(c=>`<th class="px-3 py-2 text-left text-indigo-300 border-b border-gray-700">${esc(c)}</th>`).join('')+'</tr></thead><tbody>';
  h+=r.rows.slice(0,100).map(row=>'<tr class="border-b border-gray-800">'+row.map(v=>`<td class="px-3 py-2 text-gray-300">${esc(String(v))}</td>`).join('')+'</tr>').join('');
  h+='</tbody></table></div>';
  if(r.rows.length>100)h+=`<div class="text-xs text-gray-500 mt-2">${r.rows.length} total</div>`;
  el.innerHTML=h;
}
async function loadTables(){
  const el=document.getElementById('table-list');
  el.innerHTML='<div class="text-gray-400">Loading...</div>';
  const data=await fetchJSON('/api/tables');
  if(!data.tables||!data.tables.length){el.innerHTML='<div class="text-gray-500 mt-4">No tables.</div>';return;}
  el.innerHTML='<div class="mt-4 space-y-2">'+data.tables.map(t=>
    `<div class="glass rounded-lg p-3 cursor-pointer hover:bg-gray-700 transition" onclick="document.getElementById('sql-input').value='SELECT * FROM \\"${t.table_name}\\" LIMIT 10'">
      <div class="text-indigo-300 font-medium text-sm">${t.table_name}</div>
      <div class="text-xs text-gray-500">${t.file_path}|${t.row_count} rows</div>
      <div class="text-xs text-gray-400 mt-1">${t.headers.join(', ')}</div>
    </div>`).join('')+'</div>';
}

// ---- Files ----
async function loadFiles(){
  const el=document.getElementById('file-list');
  el.innerHTML='<div class="text-gray-400">Loading...</div>';
  const data=await fetchJSON('/api/files');
  if(!data.files||!data.files.length){el.innerHTML='<div class="text-gray-500 text-center py-8">Empty</div>';return;}
  const g={};
  data.files.forEach(f=>{const d=f.source_dir||'?';if(!g[d])g[d]=[];g[d].push(f);});
  let h=`<div class="text-sm text-gray-400 mb-3">${data.count} files</div>`;
  Object.entries(g).forEach(([dir,files])=>{
    h+=`<div class="mb-4"><div class="text-xs text-gray-500 mb-1 font-mono">${esc(dir)}/</div>`;
    h+=files.map(f=>{
      const tc={'.pptx':'text-orange-300','.docx':'text-blue-300','.xlsx':'text-green-300','.pdf':'text-red-300','.md':'text-purple-300'}[f.file_type]||'text-gray-300';
      return `<div class="glass rounded px-3 py-2 mb-1 flex justify-between items-center text-sm result-card transition">
        <div class="flex items-center gap-2">
          <span class="${tc}">${f.file_type}</span>
          <span class="text-gray-300 src-link" onclick="openFile('${(f.file_path||'').replace(/'/g,"\\'")}')">${esc(f.file_name)}
            <svg class="inline ml-1" width="10" height="10" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M3 7v10a2 2 0 002 2h14a2 2 0 002-2V9a2 2 0 00-2-2h-6l-2-2H5a2 2 0 00-2 2z"/></svg>
          </span>
        </div>
        <div class="flex items-center gap-3">
          <span class="text-xs text-gray-500">${f.chunk_count}ch</span>
          <button onclick="removeFile('${f.file_id}')" class="text-red-400 hover:text-red-300 text-xs">x</button>
        </div>
      </div>`;
    }).join('');
    h+='</div>';
  });
  el.innerHTML=h;
}
async function removeFile(id){if(!confirm('Remove?'))return;await fetch(`/api/files/${id}`,{method:'DELETE'});loadFiles();loadStats();}

// ---- Ingest ----
async function doIngest(){
  const p=document.getElementById('ingest-path').value.trim();
  if(!p)return;
  const force=document.getElementById('ingest-force').checked;
  const el=document.getElementById('ingest-result');
  el.innerHTML=`<div class="glass rounded-lg p-4">
    <div class="flex justify-between mb-2"><span class="text-yellow-300" id="ingest-status">Starting...</span><span class="text-xs text-gray-500" id="ingest-pct">0%</span></div>
    <div class="w-full bg-gray-700 rounded-full h-2"><div id="ingest-bar" class="bg-indigo-500 h-2 rounded-full transition-all" style="width:0%"></div></div>
    <div class="text-xs text-gray-500 mt-2" id="ingest-file">-</div>
  </div>`;

  const url=`/api/ingest-stream?directory=${encodeURIComponent(p)}&force=${force}`;
  const evtSrc=new EventSource(url);
  evtSrc.onmessage=function(e){
    const d=JSON.parse(e.data);
    if(d.done){
      evtSrc.close();
      el.innerHTML=`<div class="glass rounded-lg p-4"><div class="text-green-300">Done ${d.elapsed_seconds||'?'}s</div>
        <div class="text-sm text-gray-300 mt-1">Processed:${d.processed||0} Skipped:${d.skipped||0} Failed:${d.failed||0} Total:${d.total||0}</div>
        ${d.errors&&d.errors.length?'<div class="text-xs text-red-400 mt-2 max-h-32 overflow-y-auto">'+d.errors.slice(0,10).map(e=>esc(e.file.split('/').pop())+': '+esc(e.error)).join('<br>')+'</div>':''}
      </div>`;
      loadStats();
    }else{
      const pct=d.total?Math.round(d.current/d.total*100):0;
      document.getElementById('ingest-bar').style.width=pct+'%';
      document.getElementById('ingest-pct').textContent=pct+'% ('+d.current+'/'+d.total+')';
      document.getElementById('ingest-status').textContent=d.status==='skipped'?'Skipping unchanged':'Processing';
      document.getElementById('ingest-file').textContent=d.name||'';
    }
  };
  evtSrc.onerror=function(){evtSrc.close();el.innerHTML+='<div class="text-red-400 text-xs mt-2">Stream ended</div>';loadStats();};
}
// File input labels
document.getElementById('upload-file').addEventListener('change',function(){
  const n=this.files.length;
  document.getElementById('upload-file-label').textContent=n?n+' file(s) selected':'No files selected';
});
document.getElementById('upload-folder').addEventListener('change',function(){
  const n=this.files.length;
  document.getElementById('upload-folder-label').textContent=n?n+' file(s) from folder':'No folder selected';
});

async function doUpload(){
  const inp=document.getElementById('upload-file');
  if(!inp.files.length)return;
  const el=document.getElementById('upload-result');
  el.innerHTML='<div class="text-yellow-300">Uploading '+inp.files.length+' file(s)...</div>';
  let ok=0,fail=0;
  for(const f of inp.files){
    const form=new FormData();form.append('file',f);
    try{const d=await fetchJSON('/api/add',{method:'POST',body:form});if(d.status==='ok')ok++;else fail++;}catch(e){fail++;}
  }
  el.innerHTML=`<div class="glass rounded-lg p-3 text-sm text-green-300">Done: ${ok} indexed, ${fail} failed</div>`;
  loadStats();
}

async function doUploadFolder(){
  const inp=document.getElementById('upload-folder');
  if(!inp.files.length)return;
  const el=document.getElementById('upload-result');
  el.innerHTML='<div class="text-yellow-300">Uploading '+inp.files.length+' file(s) from folder...</div>';
  let ok=0,fail=0;
  for(const f of inp.files){
    const form=new FormData();form.append('file',f);
    try{const d=await fetchJSON('/api/add',{method:'POST',body:form});if(d.status==='ok')ok++;else fail++;}catch(e){fail++;}
  }
  el.innerHTML=`<div class="glass rounded-lg p-3 text-sm text-green-300">Done: ${ok} indexed, ${fail} failed</div>`;
  loadStats();
}

// Directory browser
async function browseDir(){
  const current=document.getElementById('ingest-path').value||'~';
  const el=document.getElementById('dir-browser');
  el.classList.remove('hidden');
  await loadDir(current);
}

async function loadDir(path){
  const data=await fetchJSON(`/api/browse?path=${encodeURIComponent(path)}`);
  document.getElementById('dir-current').textContent=data.current||'/';
  document.getElementById('ingest-path').value=data.current||'';
  const el=document.getElementById('dir-list');
  if(!data.entries||!data.entries.length){el.innerHTML='<div class="text-xs text-gray-500">Empty</div>';return;}
  el.innerHTML=data.entries.map(e=>{
    const icon=e.type==='parent'?'&larr;':'&#128193;';
    return `<div class="flex items-center gap-2 px-2 py-1 rounded cursor-pointer hover:bg-gray-700 transition text-sm" onclick="loadDir('${e.path.replace(/'/g,"\\'")}')">
      <span>${icon}</span>
      <span class="text-gray-300">${esc(e.name)}</span>
    </div>`;
  }).join('') + `<div class="mt-2 pt-2 border-t border-gray-700">
    <button onclick="document.getElementById('dir-browser').classList.add('hidden')" class="text-xs px-3 py-1 bg-indigo-600 hover:bg-indigo-500 rounded transition">
      Use this directory
    </button>
  </div>`;
}

// ---- Settings ----
async function loadSettings(){
  const data=await fetchJSON('/api/settings');
  settings=data.settings||{};
  const embM=data.embedding_models||{}, whisM=data.whisper_models||{}, llmM=data.llm_providers||{}, buddyP=data.buddy_presets||{}, langP=data.language_profiles||{}, defs=data.defaults||{};

  _allLlmModels=llmM;
  renderLlmCards('llm-models',llmM,settings.llm_provider||'claude-sonnet');
  _allEmbModels=embM;
  _allWhisModels=whisM;
  renderGroupedCards('embedding-models',embM,settings.embedding_model||'bge-small-zh','embedding');
  renderGroupedCards('whisper-models',whisM,settings.whisper_model||'whisper-base','whisper');
  showModelConfig('embedding',settings.embedding_model||'bge-small-zh');
  showModelConfig('whisper',settings.whisper_model||'whisper-base');

  // Buddy presets
  const buddyEl=document.getElementById('buddy-presets');
  if(buddyEl){
    const selBuddy=settings.buddy_preset||'buddy';
    buddyEl.innerHTML=Object.entries(buddyP).map(([k,b])=>
      `<div class="model-card glass rounded-lg p-2 border ${k===selBuddy?'sel border-indigo-500':'border-transparent'}" onclick="pickModel(this,'buddy','${k}')">
        <span class="text-sm">${b.emoji||''} ${b.name}</span>
        <span class="text-xs text-gray-500 ml-2">${b.desc||''}</span>
      </div>`
    ).join('');
  }

  // Language profiles
  const langEl=document.getElementById('language-profiles');
  if(langEl){
    const selLang=settings.language||'zh-en';
    langEl.innerHTML=Object.entries(langP).map(([k,l])=>
      `<div class="model-card glass rounded-lg p-2 border ${k===selLang?'sel border-indigo-500':'border-transparent'}" onclick="pickModel(this,'language','${k}')">
        <div class="text-sm font-medium">${l.name}</div>
        <div class="text-xs text-gray-400">${l.desc||''}</div>
        <div class="text-xs text-gray-500 mt-1">${l.notes||''}</div>
      </div>`
    ).join('');
  }

  document.getElementById('memory-turns').value=settings.memory_turns||10;
  document.getElementById('chunk-max').value=settings.chunk_max_chars||defs.chunk_max_chars||1500;
  document.getElementById('chunk-overlap').value=settings.chunk_overlap_chars||defs.chunk_overlap_chars||200;
  // API keys and model-specific settings are now inline in LLM config panel

  // Chat dropdowns - fill with logo + name
  const sel=document.getElementById('chat-provider');
  if(sel){
    sel.innerHTML=Object.entries(llmM).map(([k,m])=>`<option value="${k}" ${k===(settings.llm_provider||'claude-sonnet')?'selected':''}>${m.name}</option>`).join('');
    sel.style.color='inherit';
  }
  const buddySel=document.getElementById('chat-buddy');
  if(buddySel){
    buddySel.innerHTML=Object.entries(buddyP).map(([k,b])=>`<option value="${k}" ${k===(settings.buddy_preset||'buddy')?'selected':''}>${b.name}</option>`).join('');
    buddySel.style.color='inherit';
  }
  console.log('Chat dropdowns filled:', sel?.options?.length, 'providers,', buddySel?.options?.length, 'buddies');

  // Update buddy greeting
  const curBuddy=buddyP[settings.buddy_preset||'buddy']||{};
  document.getElementById('buddy-avatar').textContent=curBuddy.emoji||'';
  const greetings=['Hey! Ask me anything about your files.','Ready when you are. What do you want to know?','Your files are loaded. Let\'s dig in!'];
  document.getElementById('buddy-greeting').textContent=greetings[Math.floor(Math.random()*greetings.length)];
  const statsEl=document.getElementById('buddy-status');
  const st=await fetchJSON('/api/status');
  statsEl.textContent=`${st.file_count} files | ${st.chunk_count} chunks | ${st.table_count} tables indexed`;
}

function renderLlmCards(containerId,models,selected){
  const el=document.getElementById(containerId);
  const groups={global:[],china:[],local:[]};
  Object.entries(models).forEach(([k,m])=>{
    const g=m.group||'global';
    if(!groups[g])groups[g]=[];
    groups[g].push([k,m]);
  });
  const groupLabels={global:'INTERNATIONAL','china':'CHINA (国内)',local:'LOCAL (本地)'};
  let html='';
  Object.entries(groups).forEach(([g,items])=>{
    if(!items.length)return;
    html+=`<div class="col-span-full text-xs font-medium text-gray-400 mt-3 mb-1 tracking-wider border-b border-gray-700 pb-1">${groupLabels[g]||g}</div>`;
    items.forEach(([key,m])=>{
      const isSel=selected===key;
      const logoHtml=m.logo
        ?`<img src="${m.logo}" alt="" style="width:20px;height:20px;object-fit:contain;border-radius:3px;">`
        :`<span style="font-size:16px;">${m.icon||'\u2699\ufe0f'}</span>`;
      html+=`<div class="model-card glass rounded-lg p-2 border ${isSel?'sel border-indigo-500':'border-transparent'}" onclick="pickModel(this,'llm','${key}')">
        <div class="flex items-center gap-2 mb-0.5">
          ${logoHtml}
          <span class="text-sm font-medium">${m.name}</span>
        </div>
        <div class="text-xs text-gray-500">${m.desc||''}</div>
      </div>`;
    });
  });
  el.innerHTML=html;
  showLlmConfig(selected);
}

let _allEmbModels={}, _allWhisModels={};

function showModelConfig(type, key){
  const models = type==='embedding' ? _allEmbModels : _allWhisModels;
  const m = models[key];
  const cfg = document.getElementById(type+'-config');
  if(!m || !cfg){cfg&&cfg.classList.add('hidden');return;}
  const tp = m.type||'';
  const isLocal = tp==='local' || tp==='faster-whisper';

  cfg.classList.remove('hidden');
  document.getElementById(type+'-config-title').textContent=key;

  let html='';

  if(isLocal){
    // Local model: show download status + button
    const modelName = m.name||'';
    html+=`<div id="${type}-dl-status" class="text-sm text-gray-400">Checking model status...</div>
      <div id="${type}-dl-progress" class="mt-2 hidden">
        <div class="w-full bg-gray-700 rounded-full h-2"><div id="${type}-dl-bar" class="bg-indigo-500 h-2 rounded-full transition-all" style="width:0%"></div></div>
        <div id="${type}-dl-msg" class="text-xs text-gray-400 mt-1"></div>
      </div>`;
    // Check if already downloaded
    setTimeout(async()=>{
      try{
        const st=await fetchJSON('/api/model-status/check?model_name='+encodeURIComponent(modelName));
        const el=document.getElementById(type+'-dl-status');
        if(!el)return;
        if(st.status==='downloaded'){
          el.innerHTML='<span class="text-green-400">Model ready</span> <span class="text-xs text-gray-500">('+modelName+')</span>';
        }else{
          el.innerHTML=`<span class="text-yellow-400">Model not downloaded</span>
            <button onclick="downloadModel('${type}','${modelName}')" class="ml-3 px-3 py-1 bg-indigo-600 hover:bg-indigo-500 rounded text-xs text-white transition">Download Now</button>
            <span class="text-xs text-gray-500 ml-2">${m.desc||''}</span>`;
        }
      }catch(e){
        const el=document.getElementById(type+'-dl-status');
        if(el)el.innerHTML='<span class="text-gray-500">Could not check status</span>';
      }
    },100);
  } else {
    // Cloud model: show API Key input
    if(m.key_env){
      const keyEnv=(m.key_env||'').toLowerCase();
      html+=`<label class="text-xs text-gray-400">API Key (${m.key_env})</label>
        <input id="${type}-cfg-key" type="password" placeholder="your-api-key" value="${settings[keyEnv]||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`;
    }
    if(m.base_url){
      html+=`<label class="text-xs text-gray-400 mt-2 block">API Base URL</label>
        <input id="${type}-cfg-url" type="text" value="${m.base_url}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1" readonly>`;
    }
    if(tp==='openai'){
      html+=`<label class="text-xs text-gray-400">OpenAI API Key</label>
        <input id="${type}-cfg-key" type="password" placeholder="sk-..." value="${settings.openai_api_key||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`;
    }
    if(tp==='voyageai'){
      html+=`<label class="text-xs text-gray-400">Voyage API Key</label>
        <input id="${type}-cfg-key" type="password" placeholder="pa-..." value="${settings.voyage_api_key||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`;
    }
    if(tp==='openai-api'){
      html+=`<label class="text-xs text-gray-400">OpenAI API Key</label>
        <input id="${type}-cfg-key" type="password" placeholder="sk-..." value="${settings.openai_api_key||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`;
    }
  }
  document.getElementById(type+'-config-fields').innerHTML=html;
}

function downloadModel(type, modelName){
  const progEl=document.getElementById(type+'-dl-progress');
  const barEl=document.getElementById(type+'-dl-bar');
  const msgEl=document.getElementById(type+'-dl-msg');
  const statusEl=document.getElementById(type+'-dl-status');
  if(progEl)progEl.classList.remove('hidden');
  if(statusEl)statusEl.innerHTML='<span class="text-yellow-300">Downloading...</span>';
  if(barEl)barEl.style.width='30%';

  const evtSrc=new EventSource('/api/model-download?model_name='+encodeURIComponent(modelName));
  evtSrc.onmessage=function(e){
    const d=JSON.parse(e.data);
    if(msgEl)msgEl.textContent=d.message||'';
    if(d.status==='downloading'){
      if(barEl)barEl.style.width='60%';
    }else if(d.status==='done'){
      evtSrc.close();
      if(barEl)barEl.style.width='100%';
      if(statusEl)statusEl.innerHTML='<span class="text-green-400">Model ready!</span>';
      setTimeout(()=>{if(progEl)progEl.classList.add('hidden');},2000);
    }else if(d.status==='error'){
      evtSrc.close();
      if(statusEl)statusEl.innerHTML=`<span class="text-red-400">Error: ${esc(d.message)}</span>`;
    }
  };
  evtSrc.onerror=function(){evtSrc.close();if(statusEl)statusEl.innerHTML='<span class="text-red-400">Download failed</span>';};
}

function renderGroupedCards(containerId,models,selected,type){
  const el=document.getElementById(containerId);
  const groups={china:[],global:[],local:[]};
  Object.entries(models).forEach(([k,m])=>{
    const g=m.group||'local';
    if(!groups[g])groups[g]=[];
    groups[g].push([k,m]);
  });
  const groupLabels={china:'CHINA (国内)',global:'INTERNATIONAL',local:'LOCAL (本地)'};
  let html='';
  Object.entries(groups).forEach(([g,items])=>{
    if(!items.length)return;
    html+=`<div class="col-span-full text-xs font-medium text-gray-400 mt-3 mb-1 tracking-wider border-b border-gray-700 pb-1">${groupLabels[g]||g}</div>`;
    items.forEach(([key,m])=>{
      const isSel=selected===key;
      const tp=m.type||'';
      const isLocal=tp==='local'||tp==='faster-whisper';
      const badge=isLocal?'<span class="badge-local text-xs px-2 py-0.5 rounded">Local</span>':'<span class="badge-cloud text-xs px-2 py-0.5 rounded">Cloud</span>';
      const logoHtml=m.logo?`<img src="${m.logo}" alt="" style="width:18px;height:18px;object-fit:contain;border-radius:3px;">`:'';
      html+=`<div class="model-card glass rounded-lg p-3 border ${isSel?'sel border-indigo-500':'border-transparent'}" onclick="pickModel(this,'${type}','${key}')">
        <div class="flex justify-between items-center mb-1"><div class="flex items-center gap-2">${logoHtml}<span class="text-sm font-medium">${key}</span></div>${badge}</div>
        <div class="text-xs text-gray-400">${m.desc||''}</div>
        ${m.dim?'<div class="text-xs text-gray-500 mt-1">dim:'+m.dim+'</div>':''}
      </div>`;
    });
  });
  el.innerHTML=html;
}

function renderModelCards(containerId,models,selected,type){
  const el=document.getElementById(containerId);
  el.innerHTML=Object.entries(models).map(([key,m])=>{
    const isSel=selected===key;
    const tp=m.type||'';
    const badge=tp==='local'||tp==='faster-whisper'?'<span class="badge-local text-xs px-2 py-0.5 rounded">Local</span>':'<span class="badge-cloud text-xs px-2 py-0.5 rounded">Cloud</span>';
    return `<div class="model-card glass rounded-lg p-3 border ${isSel?'sel border-indigo-500':'border-transparent'}" data-key="${key}" data-type="${type}"
      onclick="pickModel(this,'${type}','${key}')">
      <div class="flex justify-between items-center mb-1"><span class="text-sm font-medium text-gray-200">${key}</span>${badge}</div>
      <div class="text-xs text-gray-400">${m.desc||''}</div>
      ${m.dim?'<div class="text-xs text-gray-500 mt-1">dim:'+m.dim+'</div>':''}
    </div>`;
  }).join('');
}

let _allLlmModels={};
function pickModel(el,type,key){
  el.parentElement.querySelectorAll('.model-card').forEach(c=>{c.classList.remove('sel','border-indigo-500');c.classList.add('border-transparent');});
  el.classList.add('sel','border-indigo-500');el.classList.remove('border-transparent');
  if(type==='llm'){
    settings.llm_provider=key;
    showLlmConfig(key);
  }
  else if(type==='embedding'){settings.embedding_model=key;showModelConfig('embedding',key);}
  else if(type==='whisper'){settings.whisper_model=key;showModelConfig('whisper',key);}
  else if(type==='buddy')settings.buddy_preset=key;
  else if(type==='language')settings.language=key;
}

function showLlmConfig(key){
  const cfg=document.getElementById('llm-config');
  const m=_allLlmModels[key];
  if(!m){cfg.classList.add('hidden');return;}
  const tp=m.type||'';

  if(tp==='anthropic'){
    cfg.classList.remove('hidden');
    document.getElementById('llm-config-title').textContent=m.name;
    document.getElementById('llm-config-fields').innerHTML=`
      <label class="text-xs text-gray-400">Anthropic API Key</label>
      <input id="cfg-api-key" type="password" placeholder="sk-ant-..." value="${settings.anthropic_api_key||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      ${m.signup_url?`<a href="${m.signup_url}" target="_blank" rel="noopener" class="inline-block mt-2 text-xs text-indigo-400 hover:text-indigo-300 underline">Get API Key &rarr; console.anthropic.com</a>`:''}`;
  }else if(tp==='openai'){
    cfg.classList.remove('hidden');
    document.getElementById('llm-config-title').textContent=m.name;
    document.getElementById('llm-config-fields').innerHTML=`
      <label class="text-xs text-gray-400">OpenAI API Key</label>
      <input id="cfg-api-key" type="password" placeholder="sk-..." value="${settings.openai_api_key||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      ${m.signup_url?`<a href="${m.signup_url}" target="_blank" rel="noopener" class="inline-block mt-2 text-xs text-indigo-400 hover:text-indigo-300 underline">Get API Key &rarr; platform.openai.com</a>`:''}`;
  }else if(tp==='openai-compatible'){
    const keyEnv=(m.key_env||'').toLowerCase();
    const savedKey=settings[keyEnv]||'';
    const isCustom=(key==='custom');
    cfg.classList.remove('hidden');
    document.getElementById('llm-config-title').textContent=m.name;
    document.getElementById('llm-config-fields').innerHTML=`
      <label class="text-xs text-gray-400">API Key (${m.key_env||''})</label>
      <input id="cfg-api-key" type="password" placeholder="your-api-key" value="${savedKey}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      ${isCustom?`<label class="text-xs text-gray-400 mt-2 block">Model Name (模型名)</label>
      <input id="cfg-model-name" type="text" placeholder="gpt-4o / gemini-2.5-flash / qwen-max" value="${settings.custom_model||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`:''}
      <label class="text-xs text-gray-400 mt-2 block">API Base URL</label>
      <input id="cfg-base-url" type="text" value="${isCustom?(settings.custom_base_url||''):(m.base_url||'')}" placeholder="https://api.example.com/v1" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      ${m.signup_url?`<a href="${m.signup_url}" target="_blank" rel="noopener" class="inline-block mt-2 text-xs text-indigo-400 hover:text-indigo-300 underline">Get API Key &rarr; ${m.signup_url.replace(/https?:\/\//,'').split('/')[0]}</a>`:''}`;
  }else if(tp==='ollama'){
    cfg.classList.remove('hidden');
    document.getElementById('llm-config-title').textContent='Ollama Configuration';
    document.getElementById('llm-config-fields').innerHTML=`
      <label class="text-xs text-gray-400">Model Name (模型名)</label>
      <input id="cfg-ollama-model" type="text" placeholder="qwen2.5:7b" value="${settings.ollama_model||'qwen2.5:7b'}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      <label class="text-xs text-gray-400 mt-2 block">Ollama URL</label>
      <input id="cfg-ollama-url" type="text" value="${settings.ollama_url||'http://localhost:11434'}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">`;
  }else if(tp==='cli'){
    cfg.classList.remove('hidden');
    document.getElementById('llm-config-title').textContent=m.name+' Configuration';
    document.getElementById('llm-config-fields').innerHTML=`
      <label class="text-xs text-gray-400">CLI Command</label>
      <input id="cfg-cli-cmd" type="text" value="${settings.cli_command||m.cmd||'claude -p'}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1">
      <div class="text-xs text-gray-500 mt-1">No API Key needed. Uses local authentication.</div>`;
  }else{
    cfg.classList.add('hidden');
  }
}

async function saveSettings(){
  settings.chunk_max_chars=parseInt(document.getElementById('chunk-max').value)||1500;
  settings.chunk_overlap_chars=parseInt(document.getElementById('chunk-overlap').value)||200;
  settings.memory_turns=parseInt(document.getElementById('memory-turns').value)||10;
  // Read inline LLM config fields
  const m=_allLlmModels[settings.llm_provider]||{};
  const tp=m.type||'';
  const apiKeyEl=document.getElementById('cfg-api-key');
  if(apiKeyEl){
    if(tp==='anthropic') settings.anthropic_api_key=apiKeyEl.value;
    else if(tp==='openai') settings.openai_api_key=apiKeyEl.value;
    else if(tp==='openai-compatible'){
      const keyEnv=(m.key_env||'').toLowerCase();
      settings[keyEnv]=apiKeyEl.value;
    }
  }
  const baseUrlEl=document.getElementById('cfg-base-url');
  if(baseUrlEl) settings.custom_base_url=baseUrlEl.value;
  const modelNameEl=document.getElementById('cfg-model-name');
  if(modelNameEl) settings.custom_model=modelNameEl.value;
  const ollamaEl=document.getElementById('cfg-ollama-model');
  if(ollamaEl) settings.ollama_model=ollamaEl.value;
  const ollamaUrlEl=document.getElementById('cfg-ollama-url');
  if(ollamaUrlEl) settings.ollama_url=ollamaUrlEl.value;
  const cliEl=document.getElementById('cfg-cli-cmd');
  if(cliEl) settings.cli_command=cliEl.value;
  // Read embedding/whisper config keys
  ['embedding','whisper'].forEach(t=>{
    const keyEl=document.getElementById(t+'-cfg-key');
    if(keyEl&&keyEl.value){
      const models=t==='embedding'?_allEmbModels:_allWhisModels;
      const sel=t==='embedding'?settings.embedding_model:settings.whisper_model;
      const mm=models[sel]||{};
      if(mm.key_env)settings[mm.key_env.toLowerCase()]=keyEl.value;
      else if(mm.type==='openai'||mm.type==='openai-api')settings.openai_api_key=keyEl.value;
      else if(mm.type==='voyageai')settings.voyage_api_key=keyEl.value;
    }
  });
  await fetchJSON('/api/settings',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(settings)});
  const el=document.getElementById('settings-msg');
  el.innerHTML='<div class="text-green-300 py-2">Saved!</div>';
  setTimeout(()=>el.innerHTML='',3000);
  if(settings.llm_provider)document.getElementById('chat-provider').value=settings.llm_provider;
  if(settings.buddy_preset)document.getElementById('chat-buddy').value=settings.buddy_preset;
}

// ---- Tabs ----
function switchTab(name){
  ['chat','search','sql','files','ingest','connectors','settings'].forEach(t=>{
    document.getElementById('panel-'+t).classList.toggle('hidden',t!==name);
    document.getElementById('tab-'+t).classList.toggle('tab-active',t===name);
    document.getElementById('tab-'+t).classList.toggle('text-gray-400',t!==name);
  });
  if(name==='files')loadFiles();
  if(name==='sql')loadTables();
  if(name==='connectors')loadConnectors();
  if(name==='settings')loadSettings();
}

// ---- Shutdown ----
async function shutdown(){
  if(!confirm('Shutdown KBase server?'))return;
  try{await fetch('/api/shutdown',{method:'POST'});}catch(e){}
  document.body.innerHTML='<div class="flex items-center justify-center min-h-screen"><div class="text-center"><div class="text-4xl mb-4 text-gray-400">KBase Stopped</div><div class="text-gray-500">Server has been shut down. Close this tab.</div></div></div>';
}

// ---- Errors ----
async function showErrors(){
  document.getElementById('error-modal').classList.remove('hidden');
  const data=await fetchJSON('/api/errors');
  const sumEl=document.getElementById('error-summary');
  sumEl.innerHTML=`<div class="text-sm text-gray-400 mb-2">${data.total} errors total</div>
    <div class="space-y-1">${Object.entries(data.summary||{}).map(([k,v])=>
      `<div class="flex justify-between text-xs"><span class="text-gray-400 truncate mr-4">${esc(k)}</span><span class="text-red-400 font-medium">${v}</span></div>`
    ).join('')}</div>`;
  const listEl=document.getElementById('error-list');
  listEl.innerHTML='<div class="mt-3 space-y-1 max-h-60 overflow-y-auto">'+
    (data.errors||[]).slice(0,100).map(e=>
      `<div class="flex gap-2 text-xs py-1 border-b border-gray-800">
        <span class="text-indigo-300 shrink-0 cursor-pointer" onclick="openFile('${(e.file_path||'').replace(/'/g,"\\'")}')">${esc(e.file_name||'?')}</span>
        <span class="text-red-400 truncate">${esc(e.error||'')}</span>
      </div>`
    ).join('')+'</div>';
}

// ---- Theme ----
function toggleTheme(){
  document.body.classList.toggle('light');
  const isLight=document.body.classList.contains('light');
  document.getElementById('theme-btn').innerHTML=isLight?'&#9728;':'&#9790;';
  localStorage.setItem('kbase-theme',isLight?'light':'dark');
}
// Restore theme
if(localStorage.getItem('kbase-theme')==='light'){document.body.classList.add('light');document.getElementById('theme-btn').innerHTML='&#9728;';}

// ---- UI Language ----
const L={
  zh:{chat:'对话',search:'搜索',sql:'SQL查询',files:'文件',ingest:'导入',settings:'设置',
      send:'发送',clear:'清空',history:'历史',exit:'退出',search_ph:'搜索知识库...',
      chat_ph:'向知识库提问...',ingest_dir:'导入目录',browse:'浏览',force:'强制',
      upload:'上传',upload_files:'选择文件',upload_folder:'选择文件夹',
      save:'保存设置',no_files:'未选择文件',no_folder:'未选择文件夹'},
  en:{chat:'Chat',search:'Search',sql:'SQL Query',files:'Files',ingest:'Ingest',settings:'Settings',
      send:'Send',clear:'Clear',history:'History',exit:'Exit',search_ph:'Search knowledge base...',
      chat_ph:'Ask your knowledge base...',ingest_dir:'Ingest Directory',browse:'Browse',force:'Force',
      upload:'Upload',upload_files:'Choose Files',upload_folder:'Choose Folder',
      save:'Save Settings',no_files:'No files selected',no_folder:'No folder selected'}
};
let curLang=localStorage.getItem('kbase-ui-lang')||'zh';
function switchLang(lang){
  curLang=lang;
  localStorage.setItem('kbase-ui-lang',lang);
  const t=L[lang]||L.en;
  document.getElementById('tab-chat').textContent=t.chat;
  document.getElementById('tab-search').textContent=t.search;
  document.getElementById('tab-sql').textContent=t.sql;
  document.getElementById('tab-files').textContent=t.files;
  document.getElementById('tab-ingest').textContent=t.ingest;
  document.getElementById('tab-settings').textContent=t.settings;
  document.getElementById('chat-input').placeholder=t.chat_ph;
  document.getElementById('chat-btn').textContent=t.send;
  document.getElementById('search-input').placeholder=t.search_ph;
}
// Apply saved language on load
setTimeout(()=>{
  const savedLang=localStorage.getItem('kbase-ui-lang');
  if(savedLang){document.getElementById('ui-lang').value=savedLang;switchLang(savedLang);}
},100);

// ---- Chat History ----
async function showHistory(){
  const panel=document.getElementById('chat-history-panel');
  panel.classList.toggle('hidden');
  if(panel.classList.contains('hidden'))return;
  const data=await fetchJSON('/api/conversations');
  const el=document.getElementById('chat-history-list');
  if(!data.conversations||!data.conversations.length){el.innerHTML='<div class="text-xs text-gray-500">No conversations yet</div>';return;}
  const isLight=document.body.classList.contains('light');
  const hoverBg=isLight?'#e2e8f0':'#374151';
  const textClr=isLight?'#1e293b':'#d1d5db';
  const dimClr=isLight?'#64748b':'#6b7280';
  el.innerHTML=data.conversations.map(c=>
    `<div class="flex justify-between items-center px-2 py-1 rounded transition cursor-pointer text-xs" style="color:${textClr}" onmouseenter="this.style.background='${hoverBg}'" onmouseleave="this.style.background='transparent'">
      <div onclick="loadConversation('${c.id}')" class="flex-1 truncate">
        <span style="color:${textClr}">${esc(c.preview)}</span>
        <span style="color:${dimClr}" class="ml-2">${c.turns} turns</span>
      </div>
      <button onclick="deleteConversation('${c.id}')" class="text-red-400 hover:text-red-300 ml-2 shrink-0" title="Delete"><svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2"><path d="M3 6h18"/><path d="M8 6V4h8v2"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6"/><path d="M10 11v6"/><path d="M14 11v6"/></svg></button>
    </div>`
  ).join('');
}

async function loadConversation(cid){
  const data=await fetchJSON(`/api/conversations/${cid}`);
  const el=document.getElementById('chat-messages');
  el.innerHTML='';
  const welcome=document.getElementById('chat-welcome');
  if(welcome)welcome.style.display='none';
  (data.messages||[]).forEach(m=>{
    const cls=m.role==='user'?'chat-user':'chat-ai';
    el.innerHTML+=`<div class="chat-msg ${cls} rounded-lg px-4 py-3 text-sm">${m.role==='user'?esc(m.content):formatAnswer(m.content)}</div>`;
  });
  el.scrollTop=el.scrollHeight;
  document.getElementById('chat-history-panel').classList.add('hidden');
}

async function deleteConversation(cid){
  await fetch(`/api/conversations/${cid}`,{method:'DELETE'});
  showHistory();
}

// ---- Connectors ----
async function loadConnectors(){
  const data=await fetchJSON('/api/connectors');
  const el=document.getElementById('connector-list');
  el.innerHTML=Object.entries(data.connectors||{}).map(([k,c])=>{
    const connected=c.connected;
    const configured=c.configured;
    const comingSoon=c.coming_soon;
    const statusBadge=connected
      ?'<span class="badge-local text-xs px-2 py-0.5 rounded">Connected</span>'
      :configured
      ?'<span class="text-xs px-2 py-0.5 rounded" style="background:rgba(251,191,36,0.15);color:#d97706;">Configured</span>'
      :comingSoon
      ?'<span class="text-xs text-gray-500">Coming Soon</span>'
      :'<span class="text-xs text-gray-400">Not connected</span>';
    const clickHandler=comingSoon?'':`configConnector('${k}')`;
    const logoHtml=c.logo?`<img src="${c.logo}" alt="" style="width:24px;height:24px;object-fit:contain;border-radius:4px;">`:'';
    const bdr=connected?'border-green-600':configured?'border-yellow-500':'border-transparent';
    return `<div class="glass rounded-lg p-4 border ${bdr} ${comingSoon?'opacity-50':''} cursor-pointer model-card" onclick="${clickHandler}">
      <div class="flex justify-between items-center mb-2">
        <div class="flex items-center gap-2">${logoHtml}<span class="text-base font-medium">${c.name||k}</span></div>
        ${statusBadge}
      </div>
      <div class="text-xs text-gray-400">${c.desc||''}</div>
    </div>`;
  }).join('');
}

function configConnector(name){
  fetchJSON('/api/connectors').then(data=>{
    const c=data.connectors[name];
    if(!c)return;
    const cfg=document.getElementById('connector-config');
    cfg.classList.remove('hidden');
    document.getElementById('connector-config-title').textContent=c.name;
    let fieldsHtml=(c.fields||[]).map(f=>{
      if(f.type==='checkbox'){
        return `<div class="mb-2 flex items-center gap-2">
          <input id="conn-${f.key}" type="checkbox" class="w-4 h-4 rounded">
          <label class="text-sm">${f.label}</label></div>`;
      }
      return `<div class="mb-2"><label class="text-xs text-gray-400">${f.label}</label>
        <input id="conn-${f.key}" type="${f.type||'text'}" placeholder="${f.placeholder||''}" class="sinput w-full px-3 py-2 rounded-lg outline-none text-sm mt-1"></div>`;
    }).join('');
    if(c.signup_url){
      fieldsHtml+=`<div class="mt-2"><a href="${c.signup_url}" target="_blank" class="text-xs text-indigo-400 hover:text-indigo-300 underline">Create App / Get Credentials &rarr;</a></div>`;
    }
    if(c.setup_note){
      fieldsHtml+=`<div class="mt-3 p-3 rounded-lg border border-yellow-600 border-opacity-30 text-xs">
        <div class="font-medium text-yellow-400 mb-1">Setup Required:</div>
        <div class="text-gray-400">${c.setup_note}</div>
        <div class="mt-2 font-mono px-3 py-2 rounded select-all border border-indigo-400" style="background:#eef2ff; color:#4338ca; font-size:13px;">http://localhost:8765/api/connectors/feishu/callback</div>
        <div class="text-gray-500 mt-1">Copy this URL to your Feishu App &rarr; Security Settings &rarr; Redirect URLs</div>
      </div>`;
    }
    // Scope selector
    if(c.scopes&&c.scopes.length){
      fieldsHtml+=`<div class="mt-3"><div class="text-sm font-medium mb-2">Select Permissions / 选择权限 (OAuth)</div>
        <div class="text-xs text-gray-500 mb-2">Only check permissions your admin has approved. Unchecked = skip. / 只勾选管理员已开通的权限</div>
        <div class="space-y-1">${c.scopes.map(s=>
          `<label class="flex items-center gap-2 text-xs py-1">
            <input type="checkbox" class="scope-checkbox w-3.5 h-3.5" value="${s.key}" ${s.default?'checked':''}>
            <span>${s.label}</span>
            ${s.admin?'<span style="color:#d97706;font-size:10px;">(needs admin / 需管理员)</span>':'<span style="color:#059669;font-size:10px;">(free / 免审)</span>'}
          </label>`
        ).join('')}</div></div>`;
    }
    document.getElementById('connector-config-fields').innerHTML=fieldsHtml;
    document.getElementById('connector-config-actions').innerHTML=`
      <button onclick="connectSource('${name}')" class="px-4 py-2 rounded-lg text-sm transition" style="background:#4f46e5;color:#fff;">1. Save</button>
      ${name==='feishu'?`<button onclick="oauthFeishu()" class="px-4 py-2 rounded-lg text-sm transition" style="background:#7c3aed;color:#fff;">2. OAuth Login</button>`:''}
      ${name==='feishu'?`<button onclick="syncFeishu()" class="px-4 py-2 rounded-lg text-sm transition" style="background:#059669;color:#fff;">3. Sync Now</button>`:''}
      ${name==='feishu'?`<button onclick="window.open('/api/connectors/feishu/guide','_blank')" class="px-4 py-2 rounded-lg text-sm transition" style="background:#6b7280;color:#fff;">Setup Guide (教程)</button>`:''}`;
  });
}

async function connectSource(name){
  const fields=document.querySelectorAll('[id^="conn-"]');
  const body={};
  fields.forEach(f=>{const k=f.id.replace('conn-','');body[k]=f.type==='checkbox'?f.checked:f.value;});
  await fetchJSON(`/api/connectors/${name}/connect`,{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
  document.getElementById('connector-sync-result').innerHTML='<div class="text-green-300 text-sm">Saved!</div>';
  loadConnectors();
}

async function oauthFeishu(){
  try{
    await connectSource('feishu');
    // Collect selected scopes
    const checked=document.querySelectorAll('.scope-checkbox:checked');
    const scopes=Array.from(checked).map(c=>c.value).join(',');
    const data=await fetchJSON('/api/connectors/feishu/oauth-url?scopes='+encodeURIComponent(scopes));
    if(data.oauth_url){
      window.open(data.oauth_url,'_blank','width=600,height=700');
    }else{
      document.getElementById('connector-sync-result').innerHTML='<div style="color:#dc2626;" class="text-sm">Failed to get OAuth URL. Check App ID.</div>';
    }
  }catch(e){
    document.getElementById('connector-sync-result').innerHTML='<div style="color:#dc2626;" class="text-sm">Please fill in App ID and App Secret first, then Save.</div>';
  }
}

async function saveScopesToConfig(){
  const checked=document.querySelectorAll('.scope-checkbox:checked');
  const scopes=Array.from(checked).map(c=>c.value);
  // Save scopes to connector config
  const fields=document.querySelectorAll('[id^="conn-"]');
  const body={};
  fields.forEach(f=>{const k=f.id.replace('conn-','');body[k]=f.type==='checkbox'?f.checked:f.value;});
  body.selected_scopes=scopes;
  await fetchJSON('/api/connectors/feishu/connect',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
}

async function syncFeishu(){
  await saveScopesToConfig();
  const el=document.getElementById('connector-sync-result');
  el.innerHTML='<div class="text-sm" style="color:#d97706;">Syncing Feishu data...<br><span id="feishu-sync-progress" class="text-xs" style="color:#94a3b8;">Starting...</span></div>';
  const data=await fetchJSON('/api/connectors/feishu/sync',{method:'POST'});
  const perms=data.permissions||{};
  let permHtml='<div class="mt-2 space-y-1">';
  ['docs','chats','emails'].forEach(mod=>{
    const p=perms[mod];
    const labels={docs:'Cloud Docs (云文档)',chats:'Chat Messages (聊天记录)',emails:'Emails (邮件)'};
    if(!p||p.status==='skipped'){
      // Don't show skipped modules
      return;
    }else if(p.status==='ok'){
      let extra='';
      if(p.chat_errors&&p.chat_error_count){
        extra=` <span style="color:#d97706;">(${p.chat_error_count} chats failed to fetch messages)</span>`;
        if(p.chat_errors.length){extra+=`<div class="text-xs ml-4" style="color:#94a3b8;">${p.chat_errors.slice(0,3).map(e=>esc(e)).join('<br>')}</div>`;}
      }
      permHtml+=`<div class="text-xs"><span style="color:#059669;">OK</span> ${labels[mod]}: <strong>${p.count||0}</strong> items found${extra}</div>`;
    }else if(p.status==='no_permission'){
      permHtml+=`<div class="text-xs"><span style="color:#dc2626;">NO PERMISSION</span> ${labels[mod]}
        <div class="text-xs ml-4" style="color:#64748b;">Need: <code>${(p.needed||[]).join(' or ')}</code>
        ${p.admin?'<span style="color:#d97706;"> (needs admin / 需管理员)</span>':''}
        </div></div>`;
    }else{
      permHtml+=`<div class="text-xs"><span class="text-red-400">ERROR</span> ${labels[mod]}: ${esc((p.error||'').substring(0,100))}</div>`;
    }
  });
  permHtml+='</div>';
  el.innerHTML=`<div class="glass rounded-lg p-4 mt-2 text-sm">
    <div class="text-green-300 font-medium">Sync complete!</div>
    <div class="mt-1">Docs: <strong>${data.docs||0}</strong> | Chats: <strong>${data.chats||0}</strong> | Emails: <strong>${data.emails||0}</strong></div>
    ${permHtml}
    ${data.ingest?'<div class="text-xs mt-2 pt-2 border-t border-gray-700" style="color:#64748b;">Indexed: '+data.ingest.processed+' files into knowledge base</div>':''}
    ${Object.values(perms).some(p=>p.status==='no_permission')?
      '<div class="text-xs mt-2 pt-2 border-t border-gray-700" style="color:#d97706;"><strong>Tip:</strong> After adding permissions in Feishu Developer Console:<br>1. Re-publish the app version (重新发布版本)<br>2. Click <strong>2. OAuth Login</strong> again to re-authorize (重新授权)<br>3. Then <strong>3. Sync Now</strong> again</div>':''}
  </div>`;
  loadStats();
}

// ---- Init ----
loadStats();
loadSettings();
// Iframe mode: hide redundant UI elements
if(window!==window.top){
  // Hide header, stats bar, tabs — sidebar handles navigation
  document.querySelectorAll('.max-w-6xl > div:first-child, #stats-bar').forEach(el=>el.style.display='none');
  // Hide tab bar
  const tabBar=document.querySelector('.flex.gap-5.mb-5.border-b');
  if(tabBar)tabBar.style.display='none';
  // Remove top padding
  document.querySelector('.max-w-6xl').style.paddingTop='8px';
}
// Auto-switch to tab from hash
if(location.hash){
  const tab=location.hash.replace('#tab-','');
  if(tab)setTimeout(()=>switchTab(tab),100);
}
// Restore last conversation
(async()=>{
  try{
    const data=await fetchJSON(`/api/conversations/${convId}`);
    if(data.messages&&data.messages.length){
      const el=document.getElementById('chat-messages');
      const welcome=document.getElementById('chat-welcome');
      if(welcome)welcome.style.display='none';
      data.messages.forEach(m=>{
        const cls=m.role==='user'?'chat-user':'chat-ai';
        const content=m.role==='user'?esc(m.content):formatAnswer(m.content);
        el.innerHTML+=`<div class="chat-msg ${cls} rounded-lg px-4 py-3 text-sm">${content}</div>`;
      });
      chatTurns=data.turns||0;
      document.getElementById('chat-memory-badge').textContent=`Memory: ${chatTurns} turns`;
      el.scrollTop=el.scrollHeight;
    }
  }catch(e){}
})();
document.getElementById('chat-input').focus();
</script>
</body>
</html>
"""


def run_server(workspace: str = "default", host: str = "0.0.0.0", port: int = 8765):
    import uvicorn
    app = create_app(workspace)
    uvicorn.run(app, host=host, port=port)
